# -*- coding: utf-8 -*-
import os, sys
import collections
import collections.abc

if not hasattr(collections, 'Iterable'):
    collections.Iterable = collections.abc.Iterable

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Palette
    NAVY = RGBColor(15, 32, 67)
    BLUE = RGBColor(37, 99, 235)
    DARK_BLUE = RGBColor(30, 58, 138)
    LIGHT_BG = RGBColor(248, 250, 252)
    CARD_BG = RGBColor(255, 255, 255)
    BORDER_COLOR = RGBColor(226, 232, 240)
    TEXT_DARK = RGBColor(30, 41, 59)
    TEXT_MUTED = RGBColor(100, 116, 139)
    ACCENT_GREEN = RGBColor(16, 185, 129)
    ACCENT_AMBER = RGBColor(217, 119, 6)

    blank_layout = prs.slide_layouts[6]

    def set_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(9))
        bg.fill.solid()
        bg.fill.fore_color.rgb = LIGHT_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category="DEFENSE & ENTERPRISE AIEC EVALUATION"):
        # Top Accent Strip
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(0.12))
        strip.fill.solid()
        strip.fill.fore_color.rgb = BLUE
        strip.line.fill.background()

        # Category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(14.4), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(12)
        p_cat.font.bold = True
        p_cat.font.color.rgb = BLUE
        p_cat.font.name = "Arial"

        # Title
        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(14.4), Inches(0.8))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(28)
        p_t.font.bold = True
        p_t.font.color.rgb = NAVY
        p_t.font.name = "微軟正黑體"

    # ==========================================
    # Slide 1: Cover Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(9))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = NAVY
    bg1.line.fill.background()

    # Cover Header Tag
    tag_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.5), Inches(5.2), Inches(0.5))
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = BLUE
    tag_box.line.fill.background()
    p = tag_box.text_frame.paragraphs[0]
    p.text = "🛡️ AIEC 國防與企業級 AI 評測知識庫"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "微軟正黑體"

    # Cover Title
    t_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.3), Inches(13.6), Inches(2.2))
    tf = t_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "15 項國防級 AI 量化評測指標\n與驗測 SOP 建置報告"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "微軟正黑體"

    # Cover Subtitle / Standards List
    sub_box = slide1.shapes.add_textbox(Inches(1.2), Inches(4.8), Inches(13.6), Inches(2.5))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True

    lines = [
        "依據標準：NIST AI RMF 1.0 (Measure) · DoD CDAO AI T&E · MITRE ATLAS · ISO 42001 AIMS · SHIELD 治理引擎",
        "實施成果：完整建立 21 篇 Obsidian 雙向鏈結筆記與 15 項量化指標、驗測 SOP 與合格門檻對照表",
        "知識庫實體路徑：G:\\我的雲端硬碟\\secondbrain\\AIEC\\"
    ]
    for line in lines:
        p = tf_sub.add_paragraph() if tf_sub.paragraphs[0].text else tf_sub.paragraphs[0]
        p.text = "▪ " + line
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(203, 213, 225)
        p.font.name = "微軟正黑體"
        p.space_after = Pt(12)

    # ==========================================
    # Slide 2: Knowledge Base Hierarchy & MOC
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_bg(slide2)
    add_header(slide2, "AIEC 專題知識庫架構與雙腦目錄對照")

    # Card 1: Directory Structure (Left)
    card1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(7.0), Inches(6.8))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = BORDER_COLOR
    card1.line.width = Pt(1.5)

    title1 = slide2.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(6.4), Inches(0.5))
    p = title1.text_frame.paragraphs[0]
    p.text = "📂 Obsidian AIEC 目錄階層 (21 篇筆記)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = "微軟正黑體"

    body1 = slide2.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(6.4), Inches(5.8))
    tf1 = body1.text_frame
    tf1.word_wrap = True
    folders = [
        "00-Index_and_Templates/ (MOC與3大筆記範本)",
        "01-治理與規範/ (ISO 42001, SHIELD, ATLAS, RAG權限)",
        "02-評測矩陣與構面/ (15項指標SOP, 4能力層次, JATIC7構面)",
        "03-應用系統評測/ (A~F 六類CV/LLM/RAG/Agent/HMT/決策)",
        "04-地端架構與邊緣算力/ (LLM4層堆疊, Ollama/vLLM, Lattice)"
    ]
    for f in folders:
        p = tf1.add_paragraph() if tf1.paragraphs[0].text else tf1.paragraphs[0]
        p.text = "▪ " + f
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_DARK
        p.font.name = "微軟正黑體"
        p.space_after = Pt(14)

    # Card 2: Strategic Features (Right)
    card2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(1.6), Inches(7.0), Inches(6.8))
    card2.fill.solid()
    card2.fill.fore_color.rgb = CARD_BG
    card2.line.color.rgb = BORDER_COLOR
    card2.line.width = Pt(1.5)

    title2 = slide2.shapes.add_textbox(Inches(8.5), Inches(1.8), Inches(6.4), Inches(0.5))
    p = title2.text_frame.paragraphs[0]
    p.text = "🧠 知識體系設計核心特徵"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = "微軟正黑體"

    body2 = slide2.shapes.add_textbox(Inches(8.5), Inches(2.4), Inches(6.4), Inches(5.8))
    tf2 = body2.text_frame
    tf2.word_wrap = True
    features = [
        "**Karpathy 三層知識流**：與 Clippings / 知識庫 / 創作庫 無縫連結",
        "**高密度雙向鏈結 (`[[Link]]`)**：形成跨治理、評測與架構的網狀拓撲",
        "**標準 YAML Frontmatter**：包含 title, date, type, tags, status 元數據",
        "**#AIEC 標籤矩陣**：便於跨分類進行語意標籤發散與過濾",
        "**mcpvault CLI 相容**：可直接透過 Anti-Gravity 對話視窗進行自然語言檢索"
    ]
    for feat in features:
        p = tf2.add_paragraph() if tf2.paragraphs[0].text else tf2.paragraphs[0]
        parts = feat.split('**')
        r1 = p.add_run()
        r1.text = "▪ " + parts[1]
        r1.font.bold = True
        r1.font.color.rgb = BLUE
        r1.font.size = Pt(16)
        r1.font.name = "微軟正黑體"

        r2 = p.add_run()
        r2.text = parts[2]
        r2.font.color.rgb = TEXT_DARK
        r2.font.size = Pt(15)
        r2.font.name = "微軟正黑體"
        p.space_after = Pt(14)

    # Helper function to add structured metric cards to slides
    def create_metrics_slide(prs, title, category, metrics_data):
        slide = prs.slides.add_slide(blank_layout)
        set_bg(slide)
        add_header(slide, title, category)

        y_starts = [1.6, 3.0, 4.4, 5.8]
        card_height = 1.25

        for idx, m in enumerate(metrics_data):
            y = Inches(y_starts[idx % 4])
            
            # Card Base
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(14.4), Inches(card_height))
            card.fill.solid()
            card.fill.fore_color.rgb = CARD_BG
            card.line.color.rgb = BORDER_COLOR
            card.line.width = Pt(1.5)

            # Left Color Accent Pill
            pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), y + Inches(0.15), Inches(0.8), Inches(0.95))
            pill.fill.solid()
            pill.fill.fore_color.rgb = DARK_BLUE
            pill.line.fill.background()
            p = pill.text_frame.paragraphs[0]
            p.text = f"Q{m['id']}"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = "Arial"

            # Metric Title & Name
            tb_name = slide.shapes.add_textbox(Inches(1.9), y + Inches(0.1), Inches(4.2), Inches(1.05))
            tf_n = tb_name.text_frame
            tf_n.word_wrap = True
            p = tf_n.paragraphs[0]
            p.text = m['name_zh']
            p.font.size = Pt(17)
            p.font.bold = True
            p.font.color.rgb = NAVY
            p.font.name = "微軟正黑體"

            p_en = tf_n.add_paragraph()
            p_en.text = m['name_en']
            p_en.font.size = Pt(12)
            p_en.font.color.rgb = TEXT_MUTED
            p_en.font.name = "Arial"

            # SOP & Tools Column
            tb_sop = slide.shapes.add_textbox(Inches(6.2), y + Inches(0.1), Inches(4.8), Inches(1.05))
            tf_s = tb_sop.text_frame
            tf_s.word_wrap = True
            p1 = tf_s.paragraphs[0]
            p1.text = "🧪 驗測 SOP / 工具："
            p1.font.bold = True
            p1.font.size = Pt(13)
            p1.font.color.rgb = BLUE
            p1.font.name = "微軟正黑體"

            p2 = tf_s.add_paragraph()
            p2.text = m['sop']
            p2.font.size = Pt(13)
            p2.font.color.rgb = TEXT_DARK
            p2.font.name = "微軟正黑體"

            # Threshold & Standards Column
            tb_thresh = slide.shapes.add_textbox(Inches(11.1), y + Inches(0.1), Inches(3.9), Inches(1.05))
            tf_t = tb_thresh.text_frame
            tf_t.word_wrap = True
            p1 = tf_t.paragraphs[0]
            p1.text = "🎯 量化合格門檻："
            p1.font.bold = True
            p1.font.size = Pt(13)
            p1.font.color.rgb = ACCENT_GREEN
            p1.font.name = "微軟正黑體"

            p2 = tf_t.add_paragraph()
            p2.text = m['threshold']
            p2.font.size = Pt(13)
            p2.font.bold = True
            p2.font.color.rgb = TEXT_DARK
            p2.font.name = "Arial"

            p3 = tf_t.add_paragraph()
            p3.text = f"標準: {m['std']}"
            p3.font.size = Pt(11)
            p3.font.color.rgb = TEXT_MUTED
            p3.font.name = "微軟正黑體"

    # ==========================================
    # Slide 3: Metrics 1-4 (Operational & Environment)
    # ==========================================
    metrics_1_4 = [
        {
            "id": 1,
            "name_zh": "對抗韌性",
            "name_en": "Adversarial Robustness",
            "sop": "使用 IBM ART 360 / HEART 對模型注入漸進式 ε 擾動 (FGSM/PGD)，測試 mAP 變化",
            "threshold": "Acc_adv / Acc_clean >= 90%",
            "std": "MITRE ATLAS / NIST AI RMF"
        },
        {
            "id": 2,
            "name_zh": "自然穩健性",
            "name_en": "Natural Robustness",
            "sop": "透過 NRTK 合成雨雪、煙霧、電戰雜訊等 10 種環境降質數據集進行壓力測試",
            "threshold": "mAP 衰減率 <= 10%",
            "std": "JATIC / DoD CDAO"
        },
        {
            "id": 3,
            "name_zh": "任務完成率",
            "name_en": "Mission Success Rate (MSR)",
            "sop": "於 VBS 4 / EADSIM 虛實整合 (LVC) 平行戰場環境執行 100 次蒙地卡羅場景模擬",
            "threshold": "MSR >= 95%",
            "std": "Level 4 Operational T&E"
        },
        {
            "id": 4,
            "name_zh": "可中止性與失效安全",
            "name_en": "Abortability & Fail-Safe",
            "sop": "隨機注入手動 Stop Signal 及硬體斷連，量測安全降級與接管回應時間",
            "threshold": "Abort Latency <= 100ms (100% Fail-Safe)",
            "std": "DoDD 3000.09 自主武器指令"
        }
    ]
    create_metrics_slide(prs, "15 項評測指標 (1/4) —— 作戰與環境效能視角", "SECTION 1: OPERATIONAL & ENVIRONMENT METRICS", metrics_1_4)

    # ==========================================
    # Slide 4: Metrics 5-8 (Scenario & Model Capability Part 1)
    # ==========================================
    metrics_5_8 = [
        {
            "id": 5,
            "name_zh": "信任校準與過度依賴",
            "name_en": "Trust Calibration & Over-Reliance",
            "sop": "於 HMT 模擬試驗中故意提供高信心但錯誤提案，記錄操作員修正率與接管反應",
            "threshold": "ECE <= 0.05 / Over-reliance <= 5%",
            "std": "JATIC / DoD HMT Guidebook"
        },
        {
            "id": 6,
            "name_zh": "認知負荷與適應性",
            "name_en": "Cognitive Load & Adaptability",
            "sop": "操作員配戴眼動儀與 EEG 完成應變任務後填寫 NASA-TLX 心理負荷量表",
            "threshold": "NASA-TLX Score Drop >= 30% (Delay <= 2s)",
            "std": "Level 2 HSI T&E"
        },
        {
            "id": 7,
            "name_zh": "可解釋性與顯著性歸因",
            "name_en": "Explainability & Feature Attribution",
            "sop": "白箱調用 XAITK / SHAP / LIME 產出 Feature Attribution 熱力圖，比對真實目標",
            "threshold": "Point Game Score >= 0.85",
            "std": "ISO 42001 Clause 8.4"
        },
        {
            "id": 8,
            "name_zh": "提示越獄與抗注入",
            "name_en": "Prompt Jailbreak Defense",
            "sop": "使用 garak 框架自動生成 10,000 筆測試案例 (Direct/Indirect Injection & Roleplay)",
            "threshold": "Jailbreak Defense Rate >= 99%",
            "std": "OWASP LLM Top 10 / garak"
        }
    ]
    create_metrics_slide(prs, "15 項評測指標 (2/4) —— 情境與模型能力視角 (一)", "SECTION 2: SCENARIO & MODEL CAPABILITY METRICS", metrics_5_8)

    # ==========================================
    # Slide 5: Metrics 9-11 (Scenario & Model Capability Part 2)
    # ==========================================
    metrics_9_11 = [
        {
            "id": 9,
            "name_zh": "幻覺率與事實忠實度",
            "name_en": "Hallucination Rate & Faithfulness",
            "sop": "運用 RAGAS 與 TruLens 的 Faithfulness 評估器對 1,000 組問答對進行自動事實核對",
            "threshold": "Faithfulness >= 0.95 (Hallucination <= 2%)",
            "std": "NIST AI RMF / RAGAS"
        },
        {
            "id": 10,
            "name_zh": "檢索精確度與來源歸屬",
            "name_en": "RAG Context Precision & Attribution",
            "sop": "比對 RAG 檢索出的 Top-K 段落與 Ground Truth 之語意相關性與引述標註正確率",
            "threshold": "Context Precision >= 0.90 / Attribution >= 0.98",
            "std": "ISO 42001 / TruLens"
        },
        {
            "id": 11,
            "name_zh": "Agent 工具調用合規",
            "name_en": "Agent Trajectory & Misuse Audit",
            "sop": "使用 AgentBench 記錄完整 Tool Call 軌跡，經由 OPA 進行策略權限與目標偏移比對",
            "threshold": "Unauthorized API Call Rate = 0%",
            "std": "OPA / SPIFFE / AgentBench"
        },
        {
            "id": 12,
            "name_zh": "概念與數據漂移監控",
            "name_en": "Data & Concept Drift Recall",
            "sop": "部署 PyOD 與 Alibi Detect 警報模組，注入漂移數據集測試監控靈敏度與告警時延",
            "threshold": "Drift Recall >= 95% (Alarm <= 5min)",
            "std": "SHIELD Detect Stage / PyOD"
        }
    ]
    create_metrics_slide(prs, "15 項評測指標 (3/4) —— 情境與 Agent 能力視角 (二)", "SECTION 2: SCENARIO & AGENT CAPABILITY METRICS", metrics_9_11)

    # ==========================================
    # Slide 6: Metrics 13-15 (Audit & Security Governance)
    # ==========================================
    metrics_13_15 = [
        {
            "id": 13,
            "name_zh": "不確定性量化",
            "name_en": "Uncertainty Quantification (UQ)",
            "sop": "採用 MC-Dropout 或 Deep Ensembles 生成方差，測試 OOD 數據時不確定性方差激增度",
            "threshold": "OOD Variance Coverage >= 95%",
            "std": "NIST AI RMF / PyOD"
        },
        {
            "id": 14,
            "name_zh": "資料分級與防降密洩漏",
            "name_en": "Anti-Declassification Leakage",
            "sop": "模擬不同密級用戶對 RAG 探勘，查驗 LLM 摘要輸出遮罩與 RBAC 向量標籤攔截率",
            "threshold": "Declassification Leakage Rate = 0%",
            "std": "ISO 42001 Annex A / RBAC"
        },
        {
            "id": 15,
            "name_zh": "軌跡可追溯性與可稽核性",
            "name_en": "Traceability & Audit Compliance",
            "sop": "抽查歷史決策紀錄，驗證是否能從日誌中重新推導還原當時推論歷程與權重版本",
            "threshold": "Log Audit Coverage = 100% (Latency <= 10min)",
            "std": "SHIELD Log Stage / CMMC L2"
        }
    ]
    # Create Slide 6 manually for 3 items to fit comfortably
    slide6 = prs.slides.add_slide(blank_layout)
    set_bg(slide6)
    add_header(slide6, "15 項評測指標 (4/4) —— 稽核與資安治理視角", "SECTION 3: AUDIT & GOVERNANCE METRICS")

    y_starts = [1.8, 3.6, 5.4]
    card_height = 1.5

    for idx, m in enumerate(metrics_13_15):
        y = Inches(y_starts[idx])
        
        # Card Base
        card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(14.4), Inches(card_height))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1.5)

        # Left Color Accent Pill
        pill = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), y + Inches(0.2), Inches(0.9), Inches(1.1))
        pill.fill.solid()
        pill.fill.fore_color.rgb = DARK_BLUE
        pill.line.fill.background()
        p = pill.text_frame.paragraphs[0]
        p.text = f"Q{m['id']}"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = "Arial"

        # Metric Title & Name
        tb_name = slide6.shapes.add_textbox(Inches(2.0), y + Inches(0.15), Inches(4.2), Inches(1.2))
        tf_n = tb_name.text_frame
        tf_n.word_wrap = True
        p = tf_n.paragraphs[0]
        p.text = m['name_zh']
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.font.name = "微軟正黑體"

        p_en = tf_n.add_paragraph()
        p_en.text = m['name_en']
        p_en.font.size = Pt(13)
        p_en.font.color.rgb = TEXT_MUTED
        p_en.font.name = "Arial"

        # SOP & Tools Column
        tb_sop = slide6.shapes.add_textbox(Inches(6.3), y + Inches(0.15), Inches(4.8), Inches(1.2))
        tf_s = tb_sop.text_frame
        tf_s.word_wrap = True
        p1 = tf_s.paragraphs[0]
        p1.text = "🧪 驗測 SOP / 工具："
        p1.font.bold = True
        p1.font.size = Pt(14)
        p1.font.color.rgb = BLUE
        p1.font.name = "微軟正黑體"

        p2 = tf_s.add_paragraph()
        p2.text = m['sop']
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_DARK
        p2.font.name = "微軟正黑體"

        # Threshold & Standards Column
        tb_thresh = slide6.shapes.add_textbox(Inches(11.2), y + Inches(0.15), Inches(3.8), Inches(1.2))
        tf_t = tb_thresh.text_frame
        tf_t.word_wrap = True
        p1 = tf_t.paragraphs[0]
        p1.text = "🎯 量化合格門檻："
        p1.font.bold = True
        p1.font.size = Pt(14)
        p1.font.color.rgb = ACCENT_GREEN
        p1.font.name = "微軟正黑體"

        p2 = tf_t.add_paragraph()
        p2.text = m['threshold']
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_DARK
        p2.font.name = "Arial"

        p3 = tf_t.add_paragraph()
        p3.text = f"標準: {m['std']}"
        p3.font.size = Pt(12)
        p3.font.color.rgb = TEXT_MUTED
        p3.font.name = "微軟正黑體"

    # ==========================================
    # Slide 7: Anti-Gravity AI CLI Integration & Conclusion
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_bg(slide7)
    add_header(slide7, "Anti-Gravity AI CLI 自然語言檢索與運作流", "WORKFLOW & CONCLUSION")

    # Left Container: Use cases
    c_left = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(7.0), Inches(6.8))
    c_left.fill.solid()
    c_left.fill.fore_color.rgb = CARD_BG
    c_left.line.color.rgb = BORDER_COLOR
    c_left.line.width = Pt(1.5)

    t_left = slide7.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(6.4), Inches(0.5))
    p = t_left.text_frame.paragraphs[0]
    p.text = "💬 自然語言 AI CLI 整合情境"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = "微軟正黑體"

    b_left = slide7.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(6.4), Inches(5.8))
    tf_l = b_left.text_frame
    tf_l.word_wrap = True

    cli_items = [
        ("🔍 語意檢索", "「搜尋 AIEC 筆記中有關 RAG 權限控管與防降密洩漏的規範。」"),
        ("🧪 評測 SOP 查詢", "「A 類電腦視覺與目標偵測推薦哪些對抗攻防工具？」"),
        ("📝 隨手紀錄", "「將剛才討論的 15 項量化指標精隨寫入今天的每日筆記。」"),
        ("🔄 Cron 自動重整", "「每週日 09:17 自動掃描新筆記，更新 MOC 主索引頁。」")
    ]

    for title, desc in cli_items:
        p = tf_l.add_paragraph() if tf_l.paragraphs[0].text else tf_l.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = BLUE
        p.font.name = "微軟正黑體"

        p_d = tf_l.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = TEXT_DARK
        p_d.font.name = "微軟正黑體"
        p_d.space_after = Pt(12)

    # Right Container: Value Summary
    c_right = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(1.6), Inches(7.0), Inches(6.8))
    c_right.fill.solid()
    c_right.fill.fore_color.rgb = CARD_BG
    c_right.line.color.rgb = BORDER_COLOR
    c_right.line.width = Pt(1.5)

    t_right = slide7.shapes.add_textbox(Inches(8.5), Inches(1.8), Inches(6.4), Inches(0.5))
    p = t_right.text_frame.paragraphs[0]
    p.text = "🎯 專案價值與後續效益"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = "微軟正黑體"

    b_right = slide7.shapes.add_textbox(Inches(8.5), Inches(2.4), Inches(6.4), Inches(5.8))
    tf_r = b_right.text_frame
    tf_r.word_wrap = True

    value_items = [
        "1. **標準對齊**：完整連結 NIST AI RMF, DoD CDAO, ISO 42001 與 MITRE ATLAS。",
        "2. **可量化與可驗證**：提供 15 項指標之精確數學公式與 Pass/Fail 測試門檻。",
        "3. **地端主權保障**：整合四層 LLM 堆疊與 Lattice / Menace 邊緣 C2 架構。",
        "4. **雙腦第二大腦**：Markdown + YAML + 雙向鏈結，賦能 AI 全自動知識庫成長。"
    ]

    for item in value_items:
        p = tf_r.add_paragraph() if tf_r.paragraphs[0].text else tf_r.paragraphs[0]
        parts = item.split('**')
        r1 = p.add_run()
        r1.text = parts[0] + parts[1]
        r1.font.bold = True
        r1.font.size = Pt(16)
        r1.font.color.rgb = NAVY
        r1.font.name = "微軟正黑體"

        r2 = p.add_run()
        r2.text = parts[2]
        r2.font.size = Pt(15)
        r2.font.color.rgb = TEXT_DARK
        r2.font.name = "微軟正黑體"
        p.space_after = Pt(16)

    # Save
    out_dir = r'c:\Users\administartor\Downloads\AIEC'
    out_path = os.path.join(out_dir, 'AIEC_Walkthrough_15_Metrics.pptx')
    prs.save(out_path)
    print(f'Successfully generated presentation at: {out_path}')

if __name__ == '__main__':
    create_deck()
