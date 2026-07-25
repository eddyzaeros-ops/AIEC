# -*- coding: utf-8 -*-
import os, sys
import collections
import collections.abc

if not hasattr(collections, 'Iterable'):
    collections.Iterable = collections.abc.Iterable

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def is_ascii(s):
    return all(ord(c) < 128 for c in s)

def add_formatted_text(p, text, font_size=15, default_color=None, bold_color=None):
    NAVY = RGBColor(12, 35, 64)
    DARK_GRAY = RGBColor(50, 50, 50)
    if default_color is None: default_color = DARK_GRAY
    if bold_color is None: bold_color = NAVY
    
    parts = text.split('**')
    for idx, part in enumerate(parts):
        if not part:
            continue
        is_bold_part = (idx % 2 == 1)
        has_dot = ('.' in part) or ('。' in part) or (idx == 1 and '.' in text)
        
        # Dual font runs for ASCII vs Traditional Chinese
        # Split into sub-chunks of ASCII and non-ASCII to guarantee font assignment
        chunk = ""
        current_ascii = None
        for char in part:
            char_ascii = (ord(char) < 128)
            if current_ascii is None:
                current_ascii = char_ascii
                chunk += char
            elif current_ascii == char_ascii:
                chunk += char
            else:
                run = p.add_run()
                run.text = chunk
                run.font.bold = is_bold_part or has_dot
                run.font.color.rgb = bold_color if (is_bold_part or has_dot) else default_color
                run.font.size = Pt(font_size + 2) if (is_bold_part or has_dot) else Pt(font_size)
                run.font.name = "Arial" if current_ascii else "微軟正黑體"
                chunk = char
                current_ascii = char_ascii
        if chunk:
            run = p.add_run()
            run.text = chunk
            run.font.bold = is_bold_part or has_dot
            run.font.color.rgb = bold_color if (is_bold_part or has_dot) else default_color
            run.font.size = Pt(font_size + 2) if (is_bold_part or has_dot) else Pt(font_size)
            run.font.name = "Arial" if current_ascii else "微軟正黑體"

def add_formatted_bullets(tf, bullet_list, font_size=15, text_color=None, bold_color=None, bullet_char="▪"):
    NAVY = RGBColor(12, 35, 64)
    BLUE = RGBColor(37, 99, 235)
    DARK_GRAY = RGBColor(50, 50, 50)
    if text_color is None: text_color = DARK_GRAY
    if bold_color is None: bold_color = NAVY
    
    tf.word_wrap = True
    for idx, raw_item in enumerate(bullet_list):
        if idx == 0 and len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            
        p.space_after = Pt(6)
        p.space_before = Pt(2)
        
        is_sub = raw_item.startswith('  ') or raw_item.startswith('\t') or raw_item.strip().startswith('- ')
        clean_text = raw_item.strip().lstrip('*- ').strip()
        
        b_run = p.add_run()
        if is_sub:
            p.level = 1
            b_run.text = "  – "
            b_run.font.size = Pt(font_size - 1)
        else:
            p.level = 0
            b_run.text = bullet_char + " "
            b_run.font.size = Pt(font_size)
        b_run.font.bold = True
        b_run.font.color.rgb = BLUE
        b_run.font.name = "Arial"
        
        add_formatted_text(p, clean_text, font_size=font_size - 1 if is_sub else font_size, default_color=text_color, bold_color=bold_color)

def add_header(slide, title_text, category="DEFENSE & ENTERPRISE AIEC EVALUATION"):
    NAVY = RGBColor(12, 35, 64)
    BLUE = RGBColor(37, 99, 235)
    
    # Top Accent Strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(0.12))
    strip.fill.solid()
    strip.fill.fore_color.rgb = BLUE
    strip.line.fill.background()

    # Category
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(14.4), Inches(0.35))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.size = Pt(12)
    p_cat.font.bold = True
    p_cat.font.color.rgb = BLUE
    p_cat.font.name = "Arial"

    # Title
    t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(14.4), Inches(0.8))
    tf_t = t_box.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = title_text
    p_t.font.size = Pt(28)
    p_t.font.bold = True
    p_t.font.color.rgb = NAVY
    p_t.font.name = "微軟正黑體"

def add_icon_card(slide, x, y, width, height, icon_emoji, title_zh, title_en="", accent_color=RGBColor(37, 99, 235)):
    NAVY = RGBColor(12, 35, 64)
    CARD_BG = RGBColor(255, 255, 255)
    CARD_BORDER = RGBColor(218, 226, 236)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = CARD_BORDER
    card.line.width = Pt(1.5)
    
    # Dedicated Icon Box Container
    icon_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.2), Inches(y + 0.2), Inches(0.85), Inches(0.85))
    icon_box.fill.solid()
    icon_box.fill.fore_color.rgb = accent_color
    icon_box.line.fill.background()
    
    p_icon = icon_box.text_frame.paragraphs[0]
    p_icon.text = icon_emoji
    p_icon.alignment = PP_ALIGN.CENTER
    p_icon.font.size = Pt(22)
    p_icon.font.name = "微軟正黑體"
    
    # Title Text Box beside Icon Box
    t_box = slide.shapes.add_textbox(Inches(x + 1.2), Inches(y + 0.15), Inches(width - 1.35), Inches(0.85))
    tf_t = t_box.text_frame
    tf_t.word_wrap = True
    
    p1 = tf_t.paragraphs[0]
    p1.text = title_zh
    p1.font.size = Pt(17)
    p1.font.bold = True
    p1.font.color.rgb = NAVY
    p1.font.name = "微軟正黑體"
    
    if title_en:
        p2 = tf_t.add_paragraph()
        p2.text = title_en
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(100, 116, 139)
        p2.font.name = "Arial"
        
    return card

def set_pure_white_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(9))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255) # Pure White Background
    bg.line.fill.background()
    return bg

print("Helper functions compiled successfully.")
