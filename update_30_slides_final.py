# -*- coding: utf-8 -*-
import os, sys
import collections
import collections.abc

if not hasattr(collections, 'Iterable'):
    collections.Iterable = collections.abc.Iterable

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from pptx_helpers import set_pure_white_bg, add_header, add_icon_card, add_formatted_bullets, add_formatted_text

# Updated Master Acronym Database
ACRONYM_DB = {
    "AIEC": ("人工智慧評測與認證體系", "Artificial Intelligence Evaluation & Certification"),
    "AI": ("人工智慧", "Artificial Intelligence"),
    "NIST": ("美國國家標準暨技術研究院", "National Institute of Standards and Technology"),
    "RMF": ("風險管理框架", "Risk Management Framework"),
    "DoD": ("美國國防部", "Department of Defense"),
    "CDAO": ("數位與人工智慧長辦公室", "Chief Digital and Artificial Intelligence Office"),
    "T&E": ("測試與評估", "Test and Evaluation"),
    "DT&E": ("發展性測試與評估", "Developmental Test and Evaluation"),
    "OT&E": ("作戰性測試與評估", "Operational Test and Evaluation"),
    "ISO": ("國際標準化組織", "International Organization for Standardization"),
    "AIMS": ("人工智慧管理系統", "Artificial Intelligence Management System"),
    "MITRE": ("邁特公司", "MITRE Corporation"),
    "ATLAS": ("人工智慧系統對抗威脅圖譜", "Adversarial Threat Landscape for Artificial-Intelligence Systems"),
    "ATT&CK": ("對抗戰術、技術與通用知識", "Adversarial Tactics, Techniques, and Common Knowledge"),
    "DAGR": ("國防人工智慧治理與風險指南", "Defense Artificial Intelligence Governance & Risk Guidelines"),
    "SHIELD": ("設定、精煉、改進、評估、記錄、偵測六大治理階段", "Set, Hone, Improve, Evaluate, Log, Detect"),
    "SOC": ("關切事項聲明", "Statement of Concern"),
    "RAI": ("負責任的人工智慧", "Responsible Artificial Intelligence"),
    "CMMC": ("網路安全成熟度模型認證", "Cybersecurity Maturity Model Certification"),
    "GDPR": ("通用資料保護規則", "General Data Protection Regulation"),
    "ML": ("機器學習", "Machine Learning"),
    "LLM": ("大語言模型", "Large Language Model"),
    "CV": ("電腦視覺", "Computer Vision"),
    "RAG": ("檢索增強生成", "Retrieval-Augmented Generation"),
    "HMT": ("人機協同戰術團隊", "Human-Autonomy Teaming"),
    "HSI": ("人機系統整合", "Human-Systems Integration"),
    "TEVV": ("測試、評估、驗證與確認", "Test, Evaluation, Verification, and Validation"),
    "JATIC": ("聯合人工智慧測試中心", "Joint AI Test Center"),
    "mAP": ("平均精確度均值", "mean Average Precision"),
    "Acc_adv": ("對抗測試樣本準確率", "Adversarial Sample Accuracy"),
    "Acc_clean": ("乾淨測試樣本準確率", "Clean Sample Accuracy"),
    "ECE": ("期望校準誤差", "Expected Calibration Error"),
    "garak": ("生成式 AI 紅隊分析工具包", "Generative AI Redteam Analysis Kit"),
    "OOD": ("分布外數據", "Out-of-Distribution"),
    "NRTK": ("自然穩健性測試工具包", "Natural Robustness Toolkit"),
    "XAITK": ("可解釋人工智慧工具包", "Explainable AI Toolkit"),
    "ART": ("對抗韌性測試工具包", "Adversarial Robustness Toolbox"),
    "HEART": ("高爆對抗紅隊測試工具", "High-Explosive Adversarial Red Teaming"),
    "FGSM": ("快速梯度符號攻擊法", "Fast Gradient Sign Method"),
    "PGD": ("投影梯度下降法", "Projected Gradient Descent"),
    "MSR": ("任務完成率", "Mission Success Rate"),
    "LVC": ("實戰-虛擬-構造平行戰場整合", "Live-Virtual-Constructive"),
    "VBS": ("虛擬戰場模擬軟體", "Virtual Battlespace"),
    "EADSIM": ("延伸防空模擬系統", "Extended Air Defense Simulation"),
    "DoDD": ("國防部指令", "Department of Defense Directive"),
    "ToAST": ("自主系統測試工具", "Testing of Autonomous Systems Tool"),
    "IDA": ("國防分析研究所", "Institute for Defense Analyses"),
    "MIT-LL": ("麻省理工學院林肯實驗室", "Massachusetts Institute of Technology Lincoln Laboratory"),
    "EEG": ("腦電圖儀", "Electroencephalography"),
    "NASA-TLX": ("美國航太總署任務負荷指數", "NASA Task Load Index"),
    "SHAP": ("沙普利附加解釋值", "SHapley Additive exPlanations"),
    "LIME": ("局部可解釋模型無關說明", "Local Interpretable Model-agnostic Explanations"),
    "OWASP": ("開放 Web 應用程式安全計畫", "Open Web Application Security Project"),
    "RAGAS": ("檢索增強生成評估指標", "Retrieval Augmented Generation Assessment"),
    "API": ("應用程式介面", "Application Programming Interface"),
    "OPA": ("開放策略代理", "Open Policy Agent"),
    "SPIFFE": ("通用安全生產身份框架", "Secure Production Identity Framework for Everyone"),
    "SPIRE": ("SPIFFE 執行期環境", "SPIFFE Runtime Environment"),
    "PyOD": ("Python 異常值偵測庫", "Python Outlier Detection"),
    "UQ": ("不確定性量化", "Uncertainty Quantification"),
    "MC-Dropout": ("蒙地卡羅 Dropout", "Monte Carlo Dropout"),
    "RBAC": ("基於角色的存取控制", "Role-Based Access Control"),
    "MoE": ("專家混合架構", "Mixture of Experts"),
    "C2": ("指揮管制", "Command and Control"),
    "ISR": ("情報、監視與偵察", "Intelligence, Surveillance, and Reconnaissance"),
    "GUI": ("圖形使用者介面", "Graphical User Interface"),
    "GGUF": ("通用模型量化格式", "Georgi Gerganov Unified Format"),
    "MOC": ("內容主索引地圖", "Map of Content"),
    "Cron": ("定時排程服務", "Chronos Command Scheduler"),
    "MLOps": ("機器學習營運", "Machine Learning Operations"),
    "NeMo": ("輝達神經網路模組化工具包", "Neural Modules"),
    "HITL": ("人在迴路/人在紐中", "Human-in-the-Loop"),
    "AIF360": ("AI 公平性 360 工具包", "AI Fairness 360 Toolkit"),
    "MAITE": ("模型與 AI 測試評估基礎設施框架", "Model & AI Test and Evaluation Infrastructure Framework")
}

def add_acronym_footer(slide, acronym_keys, y_pos=7.95, height=0.9):
    NAVY = RGBColor(12, 35, 64)
    BLUE = RGBColor(37, 99, 235)
    MUTED = RGBColor(100, 116, 139)
    BG_FOOTER = RGBColor(241, 245, 249)
    BORDER_FOOTER = RGBColor(203, 213, 225)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y_pos), Inches(14.4), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = BG_FOOTER
    card.line.color.rgb = BORDER_FOOTER
    card.line.width = Pt(1.0)

    tf = card.text_frame
    tf.word_wrap = True

    p_head = tf.paragraphs[0]
    p_head.text = "📌 頁面英文縮寫與專有名詞對照 (Acronym & Term Footnotes):"
    p_head.font.size = Pt(10)
    p_head.font.bold = True
    p_head.font.color.rgb = BLUE
    p_head.font.name = "微軟正黑體"
    p_head.space_after = Pt(2)

    footer_parts = []
    for k in acronym_keys:
        if k in ACRONYM_DB:
            zh, en = ACRONYM_DB[k]
            footer_parts.append(f"**{k}**: {zh} ({en})")

    footer_str = "  |  ".join(footer_parts)

    p_body = tf.add_paragraph()
    add_formatted_text(p_body, footer_str, font_size=9, default_color=MUTED, bold_color=NAVY)

def add_infographic_card(slide, x, y, width, height, img_filename, label_zh):
    NAVY = RGBColor(12, 35, 64)
    CARD_BG = RGBColor(255, 255, 255)
    CARD_BORDER = RGBColor(218, 226, 236)
    img_dir = r"c:\Users\administartor\Downloads\AIEC\images"
    img_path = os.path.join(img_dir, img_filename)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = CARD_BORDER
    card.line.width = Pt(1.5)

    tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(width - 0.4), Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"🎨 {label_zh}"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = "微軟正黑體"

    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(x + 0.2), Inches(y + 0.7), width=Inches(width - 0.4), height=Inches(height - 0.85))

def build_30_ai_eval_deck():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank_layout = prs.slide_layouts[6]

    NAVY = RGBColor(12, 35, 64)
    BLUE = RGBColor(37, 99, 235)
    DARK_BLUE = RGBColor(30, 58, 138)
    CARD_BG = RGBColor(255, 255, 255)
    CARD_BORDER = RGBColor(218, 226, 236)

    def build_2card_slide_fn(title, category, icon1, t1_zh, t1_en, b1, icon2, t2_zh, t2_en, b2, acronym_keys, accent1=BLUE, accent2=DARK_BLUE):
        s = prs.slides.add_slide(blank_layout)
        set_pure_white_bg(s)
        add_header(s, title, category)

        add_icon_card(s, 0.8, 1.6, 7.0, 6.25, icon1, t1_zh, t1_en, accent_color=accent1)
        tb_b1 = s.shapes.add_textbox(Inches(1.1), Inches(2.7), Inches(6.4), Inches(5.0))
        add_formatted_bullets(tb_b1.text_frame, b1, font_size=13.5)

        add_icon_card(s, 8.2, 1.6, 7.0, 6.25, icon2, t2_zh, t2_en, accent_color=accent2)
        tb_b2 = s.shapes.add_textbox(Inches(8.5), Inches(2.7), Inches(6.4), Inches(5.0))
        add_formatted_bullets(tb_b2.text_frame, b2, font_size=13.5)

        add_acronym_footer(s, acronym_keys, y_pos=7.95, height=0.9)
        return s

    def build_infographic_slide(title, category, icon1, t1_zh, t1_en, b1, img_filename, img_title, acronym_keys):
        s = prs.slides.add_slide(blank_layout)
        set_pure_white_bg(s)
        add_header(s, title, category)

        add_icon_card(s, 0.8, 1.6, 7.0, 6.25, icon1, t1_zh, t1_en, accent_color=BLUE)
        tb_b1 = s.shapes.add_textbox(Inches(1.1), Inches(2.7), Inches(6.4), Inches(5.0))
        add_formatted_bullets(tb_b1.text_frame, b1, font_size=13.5)

        add_infographic_card(s, 8.2, 1.6, 7.0, 6.25, img_filename, img_title)

        add_acronym_footer(s, acronym_keys, y_pos=7.95, height=0.9)
        return s

    def build_metric_pair_slide_fn(title, category, m1, m2, acronym_keys):
        s = prs.slides.add_slide(blank_layout)
        set_pure_white_bg(s)
        add_header(s, title, category)

        metrics = [m1, m2]
        x_offsets = [0.8, 8.2]

        for idx, m in enumerate(metrics):
            x = x_offsets[idx]
            
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(7.0), Inches(6.25))
            card.fill.solid()
            card.fill.fore_color.rgb = CARD_BG
            card.line.color.rgb = CARD_BORDER
            card.line.width = Pt(1.5)

            ib = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.2), Inches(1.8), Inches(0.95), Inches(0.95))
            ib.fill.solid()
            ib.fill.fore_color.rgb = DARK_BLUE
            ib.line.fill.background()
            p = ib.text_frame.paragraphs[0]
            p.text = f"Q{m['id']}"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = "Arial"

            tb_t = s.shapes.add_textbox(Inches(x + 1.25), Inches(1.75), Inches(5.5), Inches(1.0))
            tf_t = tb_t.text_frame
            tf_t.word_wrap = True
            p1 = tf_t.paragraphs[0]
            p1.text = m['name_zh']
            p1.font.size = Pt(18)
            p1.font.bold = True
            p1.font.color.rgb = NAVY
            p1.font.name = "微軟正黑體"

            p2 = tf_t.add_paragraph()
            p2.text = m['name_en']
            p2.font.size = Pt(13)
            p2.font.color.rgb = RGBColor(100, 116, 139)
            p2.font.name = "Arial"

            tb_c = s.shapes.add_textbox(Inches(x + 0.3), Inches(2.85), Inches(6.4), Inches(4.8))
            tf_c = tb_c.text_frame
            tf_c.word_wrap = True

            calc_str = m.get('calc_formula', '')
            bullet_items = [
                f"**1. 指標定義與計算公式**：{m['def']}\n   **{calc_str}**",
                f"**2. 驗測 SOP 與工具**：{m['sop']}",
                f"**3. 量化合格門檻公式**：\n   **{m['thresh_formula']}**",
                f"**4. 對應國際標準**：{m['std']}"
            ]
            add_formatted_bullets(tf_c, bullet_items, font_size=11.5)

        add_acronym_footer(s, acronym_keys, y_pos=7.95, height=0.9)
        return s

    # ----------------------------------------------------
    # Slide 1: Cover Slide
    # ----------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s1)

    cover_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.6), Inches(14.4), Inches(7.8))
    cover_card.fill.solid()
    cover_card.fill.fore_color.rgb = NAVY
    cover_card.line.fill.background()

    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(1.1), Inches(6.2), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = BLUE
    badge.line.fill.background()
    p = badge.text_frame.paragraphs[0]
    p.text = "🛡️ AIEC 國防與企業級 AI 專屬評測與驗測 SOP 全集"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "微軟正黑體"

    tb1 = s1.shapes.add_textbox(Inches(1.3), Inches(1.9), Inches(13.4), Inches(2.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "AIEC 人工智慧專屬評測體系、\n15 項量化評測 SOP 與測試架構"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "微軟正黑體"

    tb_sub = s1.shapes.add_textbox(Inches(1.3), Inches(4.3), Inches(13.4), Inches(2.4))
    tf_sub = tb_sub.text_frame
    tf_sub.word_wrap = True

    cover_bullets = [
        "1. **評測標準整合**：精準融合 NIST AI RMF (Measure), DoD CDAO AI T&E, ISO 42001 與 MITRE ATLAS",
        "2. **15 項量化評測 SOP**：涵蓋對抗韌性、自然穩健性、MSR、可中止性、信任校準、garak/RAGAS 測試",
        "3. **6 大系統分類驗測**：包含 CV 對抗干擾、LLM 越獄防禦、RAG 三元組、Agent 軌跡與 HMT 人工介入評測"
    ]
    add_formatted_bullets(tf_sub, cover_bullets, font_size=16, text_color=RGBColor(226, 232, 240), bold_color=RGBColor(255, 255, 255))

    tb_foot1 = s1.shapes.add_textbox(Inches(1.3), Inches(6.8), Inches(13.4), Inches(1.2))
    tf_f1 = tb_foot1.text_frame
    tf_f1.word_wrap = True
    p_f1 = tf_f1.paragraphs[0]
    p_f1.text = "📌 頁面核心英文縮寫全稱 (Core Acronym Footnotes):"
    p_f1.font.size = Pt(11)
    p_f1.font.bold = True
    p_f1.font.color.rgb = BLUE
    p_f1.font.name = "微軟正黑體"

    s1_acronyms = ["AIEC", "NIST", "RMF", "DoD", "CDAO", "T&E", "ISO", "AIMS", "MITRE", "ATLAS", "LLM", "C2"]
    f1_str = "  |  ".join([f"**{k}**: {ACRONYM_DB[k][0]} ({ACRONYM_DB[k][1]})" for k in s1_acronyms if k in ACRONYM_DB])
    p_f1_b = tf_f1.add_paragraph()
    add_formatted_text(p_f1_b, f1_str, font_size=9, default_color=RGBColor(203, 213, 225), bold_color=RGBColor(255, 255, 255))

    # ----------------------------------------------------
    # Slide 2: Table of Contents
    # ----------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s2)
    add_header(s2, "AIEC 國防與企業級 AI 評測全集 — 簡報目錄", "TABLE OF CONTENTS")

    add_icon_card(s2, 0.8, 1.6, 7.0, 6.25, "📋", "簡報章節目錄 (一)", "Table of Contents - Part 1 to 3", accent_color=BLUE)
    tb_toc1 = s2.shapes.add_textbox(Inches(1.1), Inches(2.7), Inches(6.4), Inches(5.0))
    toc1_bullets = [
        "1. **第一區塊：AI 評測體系與治理框架 (P.1 - 6)**",
        "  - P.1 封面 | P.2 目錄 | P.3 雙支柱 | P.4 SHIELD 6階段 | P.5 ISO 42001 | P.6 MITRE ATLAS",
        "2. **第二區塊：評測矩陣、能力層次與方法論 (P.7 - 10)**",
        "  - P.7 T&E 四能力層次 | P.8 JATIC 7構面 | P.9 評測6大方法論 | P.10 RAG權限降密",
        "3. **第三區塊：六大應用系統 AI 專屬評測 SOP (P.11 - 16)**",
        "  - P.11 CV目標偵測 | P.12 LLM越獄 | P.13 RAG三元組 | P.14 Agent軌跡 | P.15 HMT人機 | P.16 預測分析"
    ]
    add_formatted_bullets(tb_toc1.text_frame, toc1_bullets, font_size=13.5)

    add_icon_card(s2, 8.2, 1.6, 7.0, 6.25, "📑", "簡報章節目錄 (二)", "Table of Contents - Part 4 to 5", accent_color=DARK_BLUE)
    tb_toc2 = s2.shapes.add_textbox(Inches(8.5), Inches(2.7), Inches(6.4), Inches(5.0))
    toc2_bullets = [
        "4. **第四區塊：15 項國防級 AI 量化評測指標與 SOP (P.17 - 25)**",
        "  - P.17 15項指標總覽 | P.18 Q1-Q2 對抗與自然穩健 | P.19 Q3-Q4 任務與失效 | P.20 Q5-Q6 信任與認知",
        "  - P.21 Q7-Q8 可解釋與越獄 | P.22 Q9-Q10 幻覺與精確 | P.23 Q11-Q12 Agent與漂移 | P.24 Q13-Q14 UQ與降密 | P.25 Q15 可追溯",
        "5. **第五區塊：評測平台與自動化 T&E 流水線 (P.26 - 30)**",
        "  - P.26 主權LLM測試架構 | P.27 地端推論評測 | P.28 C2與邊緣評測 | P.29 AI CLI與Cron | P.30 MAITE未來展望"
    ]
    add_formatted_bullets(tb_toc2.text_frame, toc2_bullets, font_size=13.5)

    add_acronym_footer(s2, ["AIEC", "T&E", "SOP", "MOC", "ISO", "AIMS", "MAITE"], y_pos=7.95, height=0.9)

    # ----------------------------------------------------
    # Slide 3: AIEC Dual Pillars
    # ----------------------------------------------------
    build_2card_slide_fn(
        "AIEC 評測雙支柱：DAGR 風險指南與 SHIELD 評測活動", "SECTION 1: EVALUATION FRAMEWORK & GOVERNANCE",
        "⚖️", "1. DAGR 風險評估指南", "DAGR Risk Evaluation Guidelines",
        [
            "1. **全生命週期風險矩陣**：提供研發、T&E 測試、部署至退役之危害識別",
            "2. **安全關鍵審查 (Safety-Critical)**：劃分高風險評測任務之道德與法律門檻",
            "3. **與 ISO 42001 勾稽**：對接國際人工智慧管理系統 AIMS 風險控制項",
            "4. **責任制落實 (Responsible AI)**：明確定義 AI 測試出錯時的評測責任"
        ],
        "🛡️", "2. SHIELD 評測循環活動", "SHIELD T&E Governance Activity",
        [
            "S. **Set Foundations**：設定倫理與政策基礎，產出 SOC 評測關切事項",
            "H. **Hone Operationalizations**：將政策轉化為量化 T&E 評測計畫",
            "I. **Improve & Innovate**：運用風險緩解工具進行測試與修復",
            "E. **Evaluate Status**：綜合評估危害解決程度與 T&E 量化滿意度",
            "L. **Log for Traceability**：全程文件化追溯評測數據與 Log 歷史",
            "D. **Detect via Monitoring**：上線後持續監控效能衰減與數據漂移"
        ],
        ["AIEC", "DAGR", "SHIELD", "SOC", "T&E", "ISO", "AIMS", "RAI"]
    )

    # ----------------------------------------------------
    # Slide 4: SHIELD 6 Stages Detailed
    # ----------------------------------------------------
    build_2card_slide_fn(
        "SHIELD 評測六大階段：從 T&E 計畫到持續測評", "SECTION 1: EVALUATION FRAMEWORK & GOVERNANCE",
        "🔄", "1. 前半週期：T&E 評測規劃與防範", "Set, Hone & Improve Stages",
        [
            "1. **Set Foundations (設定基礎)**：辨識 RAI 基礎，產出關切事項聲明 (SOC)",
            "2. **Hone Operationalizations (精煉操作化)**：對照 [[T&E 四大能力層次]] 制定 SOP",
            "3. **Improve & Innovate (改進與創新)**：導入 NeMo Guardrails 與對抗防禦緩解 SOC"
        ],
        "🔍", "2. 後半週期：量化評估、追溯與監控", "Evaluate, Log & Detect Stages",
        [
            "4. **Evaluate Status (評估狀態)**：比對 [[JATIC 七大共通評測構面]] 量化結果",
            "5. **Log for Traceability (記錄可追溯)**：合規 CMMC L2，完整保留演進與 Log 軌跡",
            "6. **Detect via Monitoring (持續監控)**：即時偵測模型過期與概念漂移 (Concept Drift)"
        ],
        ["SHIELD", "RAI", "SOC", "T&E", "SOP", "JATIC", "CMMC"]
    )

    # ----------------------------------------------------
    # Slide 5: ISO 42001 AIMS
    # ----------------------------------------------------
    build_2card_slide_fn(
        "ISO 42001 (AIMS) 人工智慧評測與控制項稽核要求", "SECTION 1: EVALUATION FRAMEWORK & GOVERNANCE",
        "📜", "1. AIMS 評測核心管理要求", "ISO/IEC 42001 Evaluation Clause",
        [
            "1. **Clause 6.1.2 AI 風險評估**：針對全生命週期進行威脅建模與處置驗證",
            "2. **Clause 8.4 透明度**：要求 AI 決策邏輯具備可追溯與可解釋性驗測",
            "3. **Clause 9 績效評估**：要求建立定期內部稽核與高階管理 T&E 審查"
        ],
        "🛡️", "2. Annex A 評測附錄控制項對映", "Annex A Controls & Verification",
        [
            "1. **偏見與公平性稽核**：審查訓練數據，防止對特定群體的隱性歧視測試",
            "2. **數據治理 (Data Governance)**：確保數據來源合法性與隱私保護 (GDPR)",
            "3. **ISO 17025 方法確效**：評測工具鏈與 SOP 必須經過標準檢定與確效"
        ],
        ["ISO", "AIMS", "AI", "ATLAS", "XAITK", "CV", "SHAP", "LIME", "GDPR", "SOP"]
    )

    # ----------------------------------------------------
    # Slide 6: MITRE ATLAS Threat Matrix
    # ----------------------------------------------------
    build_2card_slide_fn(
        "MITRE ATLAS 對抗威脅評測矩陣與攻防測評", "SECTION 1: EVALUATION FRAMEWORK & GOVERNANCE",
        "⚔️", "1. ATT&CK vs. ATLAS 對比", "Traditional IT vs. AI Security Testing",
        [
            "1. **評測關注對象**：ATT&CK 測試傳統網路；ATLAS 評測 AI/ML 模型與數據鏈",
            "2. **評測攻擊面**：ATLAS 聚焦模型權重、訓練資料集、Prompt 注入與向量庫",
            "3. **評測典型手法**：Data Poisoning, Jailbreak, Model Inversion, Adversarial Patch"
        ],
        "🎯", "2. 國防 AI ATLAS 防禦評測", "Defense ATLAS Tactics & Mitigations",
        [
            "1. **對抗干擾防禦測試**：針對邊緣 CV 的對抗貼片，實施 [[A類 - 電腦視覺評測]]",
            "2. **Prompt 越獄防禦測試**：針對 LLM 指管對答，導入 [[B類 - 生成式 AI 評測]]",
            "3. **經檢索注入防禦測試**：針對 RAG 向量庫，實施 [[RAG 權限控管與稽核]]"
        ],
        ["MITRE", "ATLAS", "ATT&CK", "ML", "LLM", "RAG", "CV", "SOP"]
    )

    # ----------------------------------------------------
    # Slide 7: T&E Four Capability Axes
    # ----------------------------------------------------
    build_2card_slide_fn(
        "T&E 四大能力層次評測 (Capability Axes T&E)", "SECTION 2: T&E MATRIX & METHODOLOGY",
        "📐", "1. Level 1 & Level 2 評測", "Model & HSI T&E Level",
        [
            "1. **Level 1 Model T&E (模型單體評測)**：演算法效能、對抗攻防、校準度與偏見測試",
            "2. **Level 2 HSI T&E (人機整合評測)**：Human-Systems Integration，測試認知負荷與過度依賴"
        ],
        "🌐", "2. Level 3 & Level 4 評測", "Systems & Operational T&E Level",
        [
            "3. **Level 3 Systems T&E (系統整合評測)**：端到端數據鏈、API 閘道與 Agent 軌跡測試",
            "4. **Level 4 Operational T&E (作戰實測)**：電戰干擾適應力、環境漂移與停用機制"
        ],
        ["T&E", "HSI", "TEVV", "API", "C2"]
    )

    # ----------------------------------------------------
    # Slide 8: JATIC 7 Common Dimensions
    # ----------------------------------------------------
    build_2card_slide_fn(
        "JATIC 七大跨系統共通評測構面", "SECTION 2: T&E MATRIX & METHODOLOGY",
        "🎯", "1. 共通技術基底評測 (構面 1-4)", "Robustness, Resiliency & Competence",
        [
            "1. **穩健性評測 (Robustness)**：在面對噪聲、OOD 輸入與對抗樣本時測試效能",
            "2. **韌性評測 (Resiliency)**：受到網路攻擊或電戰干擾時測試自動降級與復原",
            "3. **可解釋性評測 (Explainability)**：評測特徵歸因與熱力圖，消除黑盒子",
            "4. **勝任度評測 (Competence)**：在其指定作戰邊界內的任務完成率與精度"
        ],
        "⚖️", "2. 信任與治理基底評測 (構面 5-7)", "Fairness, Calibration & Drift",
        [
            "5. **公平性評測 (Fairness)**：審查與測試訓練數據中的隱性偏見與偏差",
            "6. **校準評測 (Calibration)**：測試信心度與實際正確率完全吻合 (ECE <= 0.05)",
            "7. **漂移監控評測 (Drift Monitoring)**：上線後持續追蹤概念漂移與數據衰減"
        ],
        ["JATIC", "AIEC", "OOD", "ECE"]
    )

    # ----------------------------------------------------
    # Slide 9: 6 T&E Methodologies
    # ----------------------------------------------------
    build_2card_slide_fn(
        "國防 AI 評測 6 大方法論矩陣", "SECTION 2: T&E MATRIX & METHODOLOGY",
        "🧪", "1. 開放與封閉評測 (黑箱/白箱/基準)", "Black-box, White-box & Benchmark",
        [
            "1. **黑箱測試 (Black-box)**：不存取權重，評測輸入輸出行為 (雲端 Tier-4 Gated Claude)",
            "2. **白箱測試 (White-box)**：存取梯度與特徵圖 (地端 Gemma 4 各層模型)",
            "3. **基準測試 (Benchmarking)**：使用標準化數據集量化比較基礎能力 (AgentBench)"
        ],
        "⚔️", "2. 安全與營運評測 (紅隊/人工/持續)", "Red Teaming, Human & Continuous",
        [
            "4. **對抗/紅隊測評 (Red Teaming)**：模擬敵方進行干擾與越獄攻擊 (IBM ART, garak)",
            "5. **人工評估 (Human Eval)**：主觀無真相時的評估 (HMT Guidebook)",
            "6. **持續監控測評 (Continuous)**：部署後即時偵測性能衰減 (Arize Phoenix, PyOD)"
        ],
        ["AIEC", "SOP", "API", "LLM", "garak", "XAITK", "SHAP", "ART", "HMT", "SHIELD", "PyOD"]
    )

    # ----------------------------------------------------
    # Slide 10: RAG Security & Data Classification
    # ----------------------------------------------------
    build_2card_slide_fn(
        "RAG 權限控管與資料分級降密評測機制", "SECTION 2: T&E MATRIX & METHODOLOGY",
        "🔐", "1. 向量庫 RBAC 評測與分級 Tag", "Vector Database Access Control T&E",
        [
            "1. **向量存儲層加標 (Metadata Tagging)**：於 Embedding 附加密級標籤與檢測",
            "2. **強制管道憑證查驗評測**：檢索前強制校驗用戶與 Agent 之數位證書",
            "3. **供應鏈安全審查**：審查開源 Embedder 模型與向量庫（Milvus）無後門"
        ],
        "🛡️", "2. 防範 LLM 摘要降密洩漏評測", "Anti-Declassification Leakage T&E",
        [
            "1. **跨文件摘要降密風險測試**：多篇限制級文件經 LLM 統整後降密推導測試",
            "2. **動態輸出遮罩評測 (Output Masking)**：驗證根據用戶權限自動遮蔽敏感實體",
            "3. **對照 ISO 42001 Annex A**：滿足機密性與權限隔離之合規評測要求"
        ],
        ["RAG", "RBAC", "AIEC", "LLM", "ISO", "AIMS"]
    )

    # ----------------------------------------------------
    # Slide 11: System A - Computer Vision
    # ----------------------------------------------------
    build_infographic_slide(
        "A類 - 電腦視覺與目標偵測對抗評測 (CV T&E)", "SECTION 3: SYSTEM-SPECIFIC AI EVALUATION SOP",
        "👁️", "A類 - 電腦視覺與目標偵測評測", "CV & Target Detection Evaluation",
        [
            "1. **對抗貼片測試 (Adversarial Patch)**：測試敵方貼附圖樣之欺騙誤判",
            "2. **自然穩健性測試**：雨雪、煙霧、電戰雜訊下 mAP 衰減測試 (NRTK)",
            "3. **分布外 (OOD) 目標測試**：新型偽裝目標誤標防範與 XAITK 特徵熱力圖驗證",
            "4. **代表工具鏈**：IBM ART 360, HEART, NRTK, XAITK"
        ],
        "adversarial_attack_defense.jpg", "對抗干擾與視覺防禦評測資訊圖表",
        ["CV", "mAP", "OOD", "SOP", "HEART", "NRTK", "XAITK", "ART"]
    )

    # ----------------------------------------------------
    # Slide 12: System B - GenAI & LLM
    # ----------------------------------------------------
    build_2card_slide_fn(
        "B類 - 生成式 AI 與大語言模型越獄評測 (GenAI & LLM)", "SECTION 3: SYSTEM-SPECIFIC AI EVALUATION SOP",
        "💬", "1. LLM 安全威脅與測試範疇", "LLM Vulnerability & Threat Scope",
        [
            "1. **幻覺測試 (Hallucination)**：評測生成虛構軍事情報與錯誤數據比例",
            "2. **越獄測試 (Jailbreak)**：發送對抗 Prompt 測試繞過系統安全護欄能力",
            "3. **敏感資料洩露測試**：評測機密上下文遭 Prompt Inversion 導出風險"
        ],
        "🛠️", "2. LLM 評測 SOP 與工具鏈", "LLM Testing Methodology & Tools",
        [
            "1. **紅隊模糊測試 (Red-teaming Fuzzing)**：自動發送萬筆越獄 Payload (garak)",
            "2. **護欄驗證測試**：測試 NeMo 護欄攔截成功率與反應時間",
            "3. **代表工具鏈**：garak (Vulnerability Scanner), NeMo Guardrails, PromptBench"
        ],
        ["AI", "LLM", "garak", "NeMo", "SOP"]
    )

    # ----------------------------------------------------
    # Slide 13: System C - RAG
    # ----------------------------------------------------
    build_infographic_slide(
        "C類 - 檢索增強生成 RAG 三元組評測 (RAG Triad T&E)", "SECTION 3: SYSTEM-SPECIFIC AI EVALUATION SOP",
        "📚", "C類 - 檢索增強生成 RAG 評測", "RAG Evaluation Framework",
        [
            "1. **知識衝突測試**：評測檢索不相關段落導致回答離題率",
            "2. **經檢索注入測試**：測試數據庫內植入對抗指令操控 LLM 風險",
            "3. **RAG Triad 評測三元組**：上下文精確度 (Precision)、忠實度 (Faithfulness) 與相關度",
            "4. **代表工具鏈**：RAGAS Assessment, TruLens Framework, Arize Phoenix"
        ],
        "rag_triad_eval.jpg", "RAG 檢索增強生成三元組評測架構圖",
        ["RAG", "LLM", "RAGAS", "SOP"]
    )

    # ----------------------------------------------------
    # Slide 14: System D - AI Agent
    # ----------------------------------------------------
    build_2card_slide_fn(
        "D類 - AI Agent 與多代理軌跡評測 (Multi-Agent T&E)", "SECTION 3: SYSTEM-SPECIFIC AI EVALUATION SOP",
        "🤖", "1. Agent 軌跡威脅與測試範疇", "Agent Threat Profile & Testing Scope",
        [
            "1. **工具誤用測試 (Tool Misuse)**：評測非法 API 調用或執行錯誤命令風險",
            "2. **目標偏移測試 (Goal Drift)**：評測多輪互動中偏離任務目標或死迴圈",
            "3. **Agent 爭奪測試**：多代理協同決策時發生邏輯死鎖或權限爭奪測試"
        ],
        "🛠️", "2. Agent 評測 SOP 與工具鏈", "Agent Testing Methodology & Tools",
        [
            "1. **軌跡稽核 (Trajectory Auditing)**：記錄完整 API 呼叫鏈與中介狀態",
            "2. **HITL 閘門測試**：驗證關鍵開火/變更授權點之強制人工介入能力",
            "3. **代表工具鏈**：AgentBench, SPIFFE/SPIRE 證書, Open Policy Agent (OPA)"
        ],
        ["AI", "T&E", "API", "HITL", "SPIFFE", "SPIRE", "OPA", "SOP"]
    )

    # ----------------------------------------------------
    # Slide 15: System E - Autonomous & HMT
    # ----------------------------------------------------
    build_infographic_slide(
        "E類 - 自主系統與人機協同 (HMT) 評測", "SECTION 3: SYSTEM-SPECIFIC AI EVALUATION SOP",
        "🚁", "E類 - 自主系統與 HMT 評測", "Autonomous Systems & HMT T&E",
        [
            "1. **非預期自主行為測試**：蜂群脫離邊界測試與緊急停用測試",
            "2. **過度依賴測試 (Over-reliance)**：評測操作員盲目信任與修正率",
            "3. **認知負荷測試 (NASA-TLX)**：評測介面警報過多引發之決策恐慌 (HITL)",
            "4. **評測規範與工具**：DoDD 3000.09 自主武器指令, ToAST, IDA HMT Guidebook"
        ],
        "hmt_trust_calibration.jpg", "HMT 人機協同信任校準與認知負荷評測儀表板",
        ["HMT", "AI", "HITL", "NASA-TLX", "DoDD", "ToAST", "IDA", "MIT-LL", "SOP"]
    )

    # ----------------------------------------------------
    # Slide 16: System F - Predictive Analytics
    # ----------------------------------------------------
    build_2card_slide_fn(
        "F類 - 預測分析與決策支援評測 (Predictive Analytics T&E)", "SECTION 3: SYSTEM-SPECIFIC AI EVALUATION SOP",
        "📈", "1. 預測模型威脅與測試範疇", "Predictive Analytics Risk Scope",
        [
            "1. **概念漂移測試 (Concept Drift)**：戰術態勢變化導致歷史模型失效測試",
            "2. **不確定性測試**：模型給出高信心度但實際為 OOD 預測風險測試",
            "3. **隱性偏見測試**：後勤或威脅排序模型受歷史數據偏差干擾測試"
        ],
        "🛠️", "2. 預測模型評測 SOP 與工具", "Predictive Testing Tools & SOP",
        [
            "1. **漂移與 OOD 偵測**：即時計算數據分布變化 (PyOD)",
            "2. **不確定性量化 (UQ)**：生成信心區間並進行特徵歸因分析",
            "3. **代表工具鏈**：AIF360, PyOD / Alibi Detect, SHAP / LIME (Feature Attribution)"
        ],
        ["OOD", "UQ", "PyOD", "AIF360", "SHAP", "LIME", "SOP"]
    )

    # ----------------------------------------------------
    # Slide 17: 15 Metrics Master Overview (3-Column Grid)
    # ----------------------------------------------------
    s17 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s17)
    add_header(s17, "15 項國防級 AI 量化評測指標與合格門檻總覽", "SECTION 4: 15 QUANTITATIVE EVALUATION METRICS & SOPS")

    add_icon_card(s17, 0.8, 1.6, 4.6, 6.25, "🎯", "一、作戰與環境效能", "Operational Performance (Q1-Q4)", accent_color=BLUE)
    tb_c1 = s17.shapes.add_textbox(Inches(0.95), Inches(2.7), Inches(4.3), Inches(5.0))
    c1_bullets = [
        "1. **Q1. 對抗韌性 (Robustness)**",
        "  - Acc_adv / Acc_clean >= 90%",
        "2. **Q2. 自然穩健性 (Natural)**",
        "  - mAP 衰減率 <= 10%",
        "3. **Q3. 任務完成率 (MSR)**",
        "  - MSR >= 95% (100次 LVC 模擬)",
        "4. **Q4. 可中止性 (Abortability)**",
        "  - Latency <= 100ms (100% Fail-Safe)"
    ]
    add_formatted_bullets(tb_c1.text_frame, c1_bullets, font_size=13.5)

    add_icon_card(s17, 5.7, 1.6, 4.6, 6.25, "🧠", "二、情境與模型能力", "Model Capabilities (Q5-Q10)", accent_color=DARK_BLUE)
    tb_c2 = s17.shapes.add_textbox(Inches(5.85), Inches(2.7), Inches(4.3), Inches(5.0))
    c2_bullets = [
        "1. **Q5. 信任校準**：ECE <= 0.05",
        "2. **Q6. 認知負荷**：NASA-TLX 下降 >= 30%",
        "3. **Q7. 可解釋性**：Point Game >= 0.85",
        "4. **Q8. 越獄防禦**：garak 防禦率 >= 99%",
        "5. **Q9. 幻覺控制**：忠實度 >= 0.95",
        "6. **Q10. RAG 精確度**：Precision >= 0.90"
    ]
    add_formatted_bullets(tb_c2.text_frame, c2_bullets, font_size=13)

    add_icon_card(s17, 10.6, 1.6, 4.6, 6.25, "🛡️", "三、稽核與資安治理", "Audit & Governance (Q11-Q15)", accent_color=BLUE)
    tb_c3 = s17.shapes.add_textbox(Inches(10.75), Inches(2.7), Inches(4.3), Inches(5.0))
    c3_bullets = [
        "1. **Q11. Agent 調用合規**：未授權 API = 0%",
        "2. **Q12. 概念數據漂移**：告警召回 >= 95%",
        "3. **Q13. 不確定性量化 (UQ)**：OOD覆蓋 >= 95%",
        "4. **Q14. 防降密洩漏**：防洩漏率 = 0% (RBAC)",
        "5. **Q15. 軌跡可追溯**：Log 稽核率 = 100%"
    ]
    add_formatted_bullets(tb_c3.text_frame, c3_bullets, font_size=13)

    add_acronym_footer(s17, ["Acc_adv", "Acc_clean", "mAP", "MSR", "ECE", "garak", "RAGAS", "UQ", "NIST", "RMF", "DoD", "CDAO", "T&E", "ISO", "AIMS", "MITRE", "ATLAS"], y_pos=7.95, height=0.9)

    # ----------------------------------------------------
    # Slide 18: Q1 & Q2
    # ----------------------------------------------------
    m1 = {
        "id": 1, "name_zh": "對抗韌性", "name_en": "Adversarial Robustness",
        "def": "模型遭受對抗貼片、FGSM/PGD 擾動攻擊時維護正確判讀能力",
        "calc_formula": "Robustness Ratio = Acc_adv(D_test, ε) / Acc_clean(D_test)",
        "sop": "使用 IBM ART 360 / HEART 對模型注入 ε 擾動，測試 mAP 變化",
        "thresh_formula": "PASS: Acc_adv / Acc_clean >= 90% (於 ε <= 0.05 條件下)",
        "std": "MITRE ATLAS / NIST AI RMF 1.0"
    }
    m2 = {
        "id": 2, "name_zh": "自然穩健性", "name_en": "Natural Robustness",
        "def": "模型面對自然環境干擾（雨雪、煙霧、電戰雜訊）時的效能維持度",
        "calc_formula": "ΔmAP = (mAP_clean - mAP_noise(η)) / mAP_clean",
        "sop": "透過 NRTK 合成 10 種等級的環境降質數據集進行壓力測試",
        "thresh_formula": "PASS: ΔmAP (mAP 衰減率) <= 10% (高噪聲測試條件下)",
        "std": "JATIC / DoD CDAO AI T&E"
    }
    build_metric_pair_slide_fn("15 項評測指標 (Q1 - Q2) —— 對抗與自然穩健性", "SECTION 4: 15 QUANTITATIVE EVALUATION METRICS & SOPS", m1, m2, ["Acc_adv", "Acc_clean", "FGSM", "PGD", "mAP", "ART", "HEART", "NRTK", "MITRE", "ATLAS", "NIST", "RMF", "JATIC", "DoD", "CDAO", "T&E"])

    # ----------------------------------------------------
    # Slide 19: Q3 & Q4
    # ----------------------------------------------------
    m3 = {
        "id": 3, "name_zh": "任務完成率", "name_en": "Mission Success Rate (MSR)",
        "def": "AI 系統在端到端戰術情境中成功執行完畢並閉合擊殺鏈的比例",
        "calc_formula": "MSR = (Σ_{i=1}^{N} S_i) / N,  S_i ∈ {0, 1}",
        "sop": "於 VBS 4 / EADSIM 虛實整合 (LVC) 平行戰場環境執行 100 次模擬",
        "thresh_formula": "PASS: MSR = 成功次數 / 總模擬數 N (100次) >= 95%",
        "std": "Level 4 Operational T&E"
    }
    m4 = {
        "id": 4, "name_zh": "可中止性與失效安全", "name_en": "Abortability & Fail-Safe Rate",
        "def": "當系統異常或接獲人工中斷指令時，即刻中斷並進入安全保護狀態",
        "calc_formula": "τ_abort = t_safe_state - t_signal_sent;  Fail-Safe Rate = N_safe / N_trigger",
        "sop": "隨機注入手動 Stop Signal 及硬體斷連，量測安全降級接管時間",
        "thresh_formula": "PASS: τ_abort <= 100ms 且 Fail-Safe Rate = 100%",
        "std": "DoDD 3000.09 自主武器指令"
    }
    build_metric_pair_slide_fn("15 項評測指標 (Q3 - Q4) —— 任務完成與失效安全", "SECTION 4: 15 QUANTITATIVE EVALUATION METRICS & SOPS", m3, m4, ["MSR", "LVC", "VBS", "EADSIM", "DoDD", "T&E"])

    # ----------------------------------------------------
    # Slide 20: Q5 & Q6
    # ----------------------------------------------------
    m5 = {
        "id": 5, "name_zh": "信任校準與過度依賴", "name_en": "Trust Calibration & Over-Reliance",
        "def": "操作員對 AI 信心度的理解符合實際能力，防止盲目信任或拒絕使用",
        "calc_formula": "ECE = Σ_{m=1}^{M} (|B_m|/N) * |acc(B_m) - conf(B_m)|;  R_overreliance = N_blind / N_false",
        "sop": "於 HMT 模擬試驗中故意提供高信心但錯誤提案，記錄操作員修正率",
        "thresh_formula": "PASS: ECE <= 0.05 且 R_overreliance <= 5%",
        "std": "JATIC / DoD HMT Guidebook"
    }
    m6 = {
        "id": 6, "name_zh": "認知負荷與適應性", "name_en": "Cognitive Load & Adaptability",
        "def": "AI 介面輸出對指揮官或操作員造成的心理負荷程度與決策時延",
        "calc_formula": "ΔTLX = (TLX_baseline - TLX_AI) / TLX_baseline;  Δt_decision = t_response",
        "sop": "操作員配戴眼動儀與 EEG 完成任務後填寫 NASA-TLX 量表",
        "thresh_formula": "PASS: ΔTLX (心理負荷下降) >= 30% 且 Δt_decision <= 2.0s",
        "std": "Level 2 HSI T&E"
    }
    build_metric_pair_slide_fn("15 項評測指標 (Q5 - Q6) —— 信任校準與認知負荷", "SECTION 4: 15 QUANTITATIVE EVALUATION METRICS & SOPS", m5, m6, ["HMT", "ECE", "EEG", "NASA-TLX", "JATIC", "DoD", "HSI", "T&E"])

    # ----------------------------------------------------
    # Slide 21: Q7 & Q8
    # ----------------------------------------------------
    m7 = {
        "id": 7, "name_zh": "模型可解釋性與顯著性歸因", "name_en": "Explainability & Point Game",
        "def": "AI 關鍵決策邏輯機能是否提供可被人類審計的特徵熱力圖 (Saliency)",
        "calc_formula": "Point Game Score = N_hit(argmax Saliency ∈ ROI) / N_total",
        "sop": "白箱調用 XAITK / SHAP / LIME 產出熱力圖，比對真實目標區域",
        "thresh_formula": "PASS: Point Game Score >= 0.85 (85%)",
        "std": "ISO 42001 Clause 8.4"
    }
    m8 = {
        "id": 8, "name_zh": "提示越獄與抗注入能力", "name_en": "Prompt Jailbreak Defense Rate",
        "def": "LLM 阻絕敵方對抗 Prompt 注入、越獄繞過與護欄突圍的能力",
        "calc_formula": "R_jailbreak_def = 1 - (N_successful_jailbreaks / N_total_attacks)",
        "sop": "使用 garak 框架執行 10,000 筆測試案例 (Direct/Indirect Injection)",
        "thresh_formula": "PASS: R_jailbreak_def (越獄防禦率) >= 99%",
        "std": "OWASP LLM Top 10 / garak"
    }
    build_metric_pair_slide_fn("15 項評測指標 (Q7 - Q8) —— 可解釋性與越獄防禦", "SECTION 4: 15 QUANTITATIVE EVALUATION METRICS & SOPS", m7, m8, ["XAITK", "SHAP", "LIME", "LLM", "garak", "OWASP", "ISO"])

    # ----------------------------------------------------
    # Slide 22: Q9 & Q10
    # ----------------------------------------------------
    m9 = {
        "id": 9, "name_zh": "幻覺率與事實忠實度", "name_en": "Hallucination Rate & Faithfulness",
        "def": "LLM 產出內容嚴格遵循檢索上下文與國防事實，無虛構捏造數據",
        "calc_formula": "Faithfulness = |Verified Statements| / |Total Statements|;  R_hallucination = 1 - Faithfulness",
        "sop": "運用 RAGAS 與 TruLens 的 Faithfulness 評估器進行自動事實核對",
        "thresh_formula": "PASS: Faithfulness Score >= 0.95 (且 R_hallucination <= 2%)",
        "std": "NIST AI RMF / RAGAS"
    }
    m10 = {
        "id": 10, "name_zh": "檢索精確度與來源歸屬", "name_en": "RAG Context Precision & Attribution",
        "def": "RAG 向量資料庫精確檢索權威規範段落並準確標註出處來源",
        "calc_formula": "Context Precision@K = (Σ_{k=1}^{K} Precision@k * v_k) / (Σ_{k=1}^{K} v_k)",
        "sop": "比對 RAG 檢索出的 Top-K 段落與 Ground Truth 之語意相關性",
        "thresh_formula": "PASS: Context Precision >= 0.90 且 Attribution Rate >= 0.98",
        "std": "ISO 42001 / TruLens"
    }
    build_metric_pair_slide_fn("15 項評測指標 (Q9 - Q10) —— 幻覺控制與 RAG 精確度", "SECTION 4: 15 QUANTITATIVE EVALUATION METRICS & SOPS", m9, m10, ["LLM", "RAG", "RAGAS", "NIST", "RMF", "ISO"])

    # ----------------------------------------------------
    # Slide 23: Q11 & Q12
    # ----------------------------------------------------
    m11 = {
        "id": 11, "name_zh": "Agent 工具調用與軌跡合規", "name_en": "Agent Trajectory Audit",
        "def": "自主 AI Agent 呼叫外部 API 時嚴格遵循權限邊界，無目標偏移",
        "calc_formula": "R_unauth_API = N_unauthorized_tool_calls / N_total_tool_calls",
        "sop": "使用 AgentBench 記錄 Tool Call 軌跡，經由 OPA 進行策略比對",
        "thresh_formula": "PASS: R_unauth_API = 0% (且 Task Success Rate >= 98%)",
        "std": "OPA / SPIFFE / AgentBench"
    }
    m12 = {
        "id": 12, "name_zh": "概念與數據漂移監控率", "name_en": "Data & Concept Drift Recall",
        "def": "系統在上線營運期間，即時捕捉數據分布變化與標籤漂移的靈敏度",
        "calc_formula": "Drift Recall = TP_drift / (TP_drift + FN_drift);  t_alarm_latency = t_alert - t_drift_occurred",
        "sop": "部署 PyOD 與 Alibi Detect 警報模組，注入漂移數據集測試反應時延",
        "thresh_formula": "PASS: Drift Recall >= 95% 且 t_alarm_latency <= 5min",
        "std": "SHIELD Detect Stage / PyOD"
    }
    build_metric_pair_slide_fn("15 項評測指標 (Q11 - Q12) —— Agent 合規與漂移監控", "SECTION 4: 15 QUANTITATIVE EVALUATION METRICS & SOPS", m11, m12, ["API", "OPA", "SPIFFE", "PyOD", "SHIELD"])

    # ----------------------------------------------------
    # Slide 24: Q13 & Q14
    # ----------------------------------------------------
    m13 = {
        "id": 13, "name_zh": "不確定性量化", "name_en": "Uncertainty Quantification (UQ)",
        "def": "模型對預測結果給出可靠信心區間，遇高不確定性時提示人類",
        "calc_formula": "Var_pred(x_OOD) > θ_variance;  OOD Coverage = N(Var_OOD > θ) / N_OOD_total",
        "sop": "採用 MC-Dropout 或 Deep Ensembles 生成方差，測試 OOD 方差激增度",
        "thresh_formula": "PASS: OOD Variance Coverage >= 95%",
        "std": "NIST AI RMF / PyOD"
    }
    m14 = {
        "id": 14, "name_zh": "資料分級與防降密洩漏", "name_en": "Anti-Declassification Leakage",
        "def": "多密級檢索時，防止低權限用戶或 LLM 摘要統整導出降密高密級資訊",
        "calc_formula": "R_declass_leak = N_unauthorized_high_classification_tokens / N_total_output_tokens",
        "sop": "模擬不同密級用戶對 RAG 探勘，查驗輸出遮罩與 RBAC 向量標籤",
        "thresh_formula": "PASS: R_declass_leak (防降密洩漏率) = 0% (RBAC Masking)",
        "std": "ISO 42001 Annex A / RBAC"
    }
    build_metric_pair_slide_fn("15 項評測指標 (Q13 - Q14) —— 不確定性與防降密洩漏", "SECTION 4: 15 QUANTITATIVE EVALUATION METRICS & SOPS", m13, m14, ["MC-Dropout", "OOD", "UQ", "RAG", "LLM", "RBAC", "ISO", "AIMS", "PyOD", "NIST", "RMF"])

    # ----------------------------------------------------
    # Slide 25: Q15
    # ----------------------------------------------------
    m15 = {
        "id": 15, "name_zh": "系統軌跡可追溯性與可稽核性", "name_en": "Traceability & Audit Compliance",
        "def": "AI 系統全生命週期的數據、權重、Prompt、API 軌跡與審核紀錄皆能完整追溯與合規重現",
        "calc_formula": "Log Coverage = N_logged_decision_traces / N_total_decisions",
        "sop": "抽查歷史決策紀錄，驗證是否能從日誌中重新推導並還原模型當時的推論歷程 (CMMC L2 / ISO 42001)",
        "thresh_formula": "PASS: Log Coverage = 100% 且 t_reproduction <= 10min",
        "std": "SHIELD Log Stage / CMMC L2"
    }
    s25 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s25)
    add_header(s25, "15 項評測指標 (Q15) —— 系統軌跡可追溯性與可稽核性", "SECTION 4: 15 QUANTITATIVE EVALUATION METRICS & SOPS")
    
    add_icon_card(s25, 0.8, 1.6, 14.4, 6.25, "📜", "Q15 - 系統軌跡可追溯性與可稽核性 (Traceability & Audit Compliance)", "CMMC Level 2 & ISO 42001 Full Auditability", accent_color=DARK_BLUE)
    tb_q15 = s25.shapes.add_textbox(Inches(1.1), Inches(2.7), Inches(13.8), Inches(5.0))
    add_formatted_bullets(tb_q15.text_frame, [
        f"**1. 指標定義與計算公式**：{m15['def']}\n   **{m15['calc_formula']}**",
        f"**2. 驗測 SOP 與方法**：{m15['sop']}",
        f"**3. 量化合格門檻公式**：\n   **{m15['thresh_formula']}**",
        f"**4. 對應國際標準**：{m15['std']}",
        "**5. 日誌鏈結規範**：包含用戶 Identity、Prompt 版本、RAG 檢索 Chunk Hash、LLM 權限 Hash 與決策輸出之全鏈結封裝"
    ], font_size=13.5)

    add_acronym_footer(s25, ["API", "RAG", "LLM", "CMMC", "ISO", "AIMS", "SHIELD"], y_pos=7.95, height=0.9)

    # ----------------------------------------------------
    # Slide 26: Sovereign AI Platform
    # ----------------------------------------------------
    build_infographic_slide(
        "評測專用四層主權 LLM 測試架構 (Tier 1~Tier 4)", "SECTION 5: EVALUATION PLATFORMS & AUTOMATED T&E",
        "🏰", "4 層 LLM 評測軟硬體架構", "4-Tier Sovereign LLM Test Stack",
        [
            "1. **Tier 1 核心評測**：Gemma 4 31B Dense (NVIDIA H200) — 地端白箱對抗測試",
            "2. **Tier 2 混合評測**：Gemma 4 26B MoE (AMD MI325X) — 地端白箱對抗測試",
            "3. **Tier 3 邊緣評測**：Gemma 4 E4B 本地推論 — 邊緣節點白箱基準測試",
            "4. **Tier 4 雲端評測**：Gated Cloud Claude API — 雲端黑箱基準測試 + 護欄閘門"
        ],
        "sovereign_ai_test_stack.jpg", "4-Tier 主權 LLM 白箱 vs 黑箱評測架構圖",
        ["AI", "LLM", "MoE", "API", "HITL", "OPA"]
    )

    # ----------------------------------------------------
    # Slide 27: Local LLM Inference Engines & Middleware
    # ----------------------------------------------------
    build_2card_slide_fn(
        "地端 LLM 推論評測引擎與 Middleware 工具鏈", "SECTION 5: EVALUATION PLATFORMS & AUTOMATED T&E",
        "🛠️", "1. 地端推論評測工具鏈對比", "Inference Evaluation Engine Comparison",
        [
            "1. **Ollama**：極簡指令集與 Docker 支援，適合 Modelfile 自訂與評測",
            "2. **LM Studio**：桌面 GUI，一鍵搜尋下載 Hugging Face GGUF 模型評測",
            "3. **llama.cpp**：地端 LLM 推論底層 C++ 核心引擎，具備最高效能評測控制",
            "4. **vLLM**：企業生產環境 High-Throughput 推論評測框架 (PagedAttention)"
        ],
        "📚", "2. Hugging Face 與 GGUF 評測量化", "Model Hub & Quantization T&E",
        [
            "1. **Hugging Face Model Hub**：全球 AI 社群模型圖書館與評測基準來源",
            "2. **GGUF 量化評測**：測試地端晶片 (Mac M系列 / NVIDIA) 推論速度與精度",
            "3. **Modelfile 評測封裝**：定義 Temperature, System Prompt 參數評測"
        ],
        ["LLM", "API", "GUI", "GGUF"]
    )

    # ----------------------------------------------------
    # Slide 28: Lattice C2 & Menace Edge Node
    # ----------------------------------------------------
    build_2card_slide_fn(
        "戰術 C2 (Lattice) 與邊緣算力節點 (Menace) 評測", "SECTION 5: EVALUATION PLATFORMS & AUTOMATED T&E",
        "⚡", "1. Lattice C2 實戰化 AI 評測", "Software-Defined C2 AI Evaluation",
        [
            "1. **任務工作流評測**：涵蓋 ISR, Strike Broker, High-Value Targeting 之 AI 測試",
            "2. **資料模型與 API 評測**：外部 Data Model API 與 Pub/Sub API 穩定度測試",
            "3. **知識圖譜評測**：測試 Entity Management 與 Sensor Fusion 之準確度"
        ],
        "💻", "2. Menace 邊緣算力節點評測", "Menace Tactical Edge Node T&E",
        [
            "1. **戰術加固評測**：測試極端氣候、強震與全地形軍規環境下之 AI 穩定度",
            "2. **分散式韌性評測**：測試 Mesh 無線網絡下 AI 無單點故障能力",
            "3. **雲-邊-端協同評測**：運行 Tier-3 輕量化模型 (Gemma 4 E4B) 邊緣目標測試"
        ],
        ["C2", "ISR", "API"]
    )

    # ----------------------------------------------------
    # Slide 29: Anti-Gravity AI CLI & Cron Schedule
    # ----------------------------------------------------
    build_2card_slide_fn(
        "Anti-Gravity AI CLI 自然語言評測檢索與 Cron 排程", "SECTION 5: EVALUATION PLATFORMS & AUTOMATED T&E",
        "💬", "1. 自然語言評測檢索工作流", "Natural Language AI T&E CLI Workflow",
        [
            "1. **對話即評測檢索**：「幫我搜尋 AIEC 筆記中有關 RAG 評測與越獄防禦的規定」",
            "2. **mcpvault API**：背景自動調用 `search_notes` 與 `write_note` 存取評測 Vault",
            "3. **隨手紀錄評測結果**：「將剛才討論的 15 項量化評測指標精隨寫入每日筆記」"
        ],
        "⏰", "2. 每週自動評測重整 Cron", "Weekly Evaluation Review Schedule",
        [
            "1. **排程時間**：每週日 09:17 (Cron：`17 9 * * 0`) 背景自動執行",
            "2. **4 大步驟**：盤點評測變動 -> 提煉 Clippings -> 評測健康檢查 -> 更新 MOC",
            "3. **評測雙腦成長**：自動偵測評測筆記缺口，產出 AI 評測週報"
        ],
        ["AI", "AIEC", "RAG", "API", "Cron", "MOC"]
    )

    # ----------------------------------------------------
    # Slide 30: Conclusion
    # ----------------------------------------------------
    build_2card_slide_fn(
        "結語與未來展望 — MAITE 自動化評測測試工廠", "SECTION 5: EVALUATION PLATFORMS & AUTOMATED T&E",
        "🎯", "1. AI 專屬評測核心貢獻", "AI Evaluation Impact & Value Summary",
        [
            "1. **權威評測標準整合**：確立 DoD CDAO, NIST AI RMF, ISO 42001 與 MITRE ATLAS 縱深評測",
            "2. **量化評測 SOP 落地**：建立 15 項量化指標、驗測 SOP 與明確 Pass/Fail 門檻",
            "3. **評測知識庫建置**：全套 21 篇雙向鏈結 Markdown 筆記完全存入 Obsidian Vault"
        ],
        "🚀", "2. 未來演進與推廣方向", "Future Roadmap & MAITE Test Factory",
        [
            "1. **MAITE 自動化評測測試工廠**：整合 garak, RAGAS, AgentBench 於單一 MLOps 流水線",
            "2. **LVC 平行戰場紅隊演練**：於 VBS 4 / EADSIM 中強化對抗貼片與電戰干擾紅隊演練",
            "3. **ISO 42001 / CMMC 評測認證**：備妥全套 SOP 與 Log 稽核軌跡，迎接正式合規檢定"
        ],
        ["DoD", "CDAO", "NIST", "RMF", "ISO", "AIMS", "MITRE", "ATLAS", "SOP", "MAITE", "MLOps", "LVC", "VBS", "EADSIM", "CMMC"]
    )

    out_dir = r'c:\Users\administartor\Downloads\AIEC'
    out_path = os.path.join(out_dir, 'AIEC_AI_Evaluation_30_Slides_NanoBanana.pptx')
    prs.save(out_path)
    print(f'Successfully updated 30-slide presentation with formulas at: {out_path}')

if __name__ == '__main__':
    build_30_ai_eval_deck()
