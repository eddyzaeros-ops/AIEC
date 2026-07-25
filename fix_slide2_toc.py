# -*- coding: utf-8 -*-
import os, sys
import collections
import collections.abc

if not hasattr(collections, 'Iterable'):
    collections.Iterable = collections.abc.Iterable

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from pptx_helpers import set_pure_white_bg, add_header, add_icon_card, add_formatted_bullets, add_formatted_text

ACRONYM_DB = {
    "AIEC": ("人工智慧評測與認證體系", "Artificial Intelligence Evaluation & Certification"),
    "T&E": ("測試與評估", "Test and Evaluation"),
    "MOC": ("內容主索引地圖", "Map of Content"),
    "ISO": ("國際標準化組織", "International Organization for Standardization"),
    "AIMS": ("人工智慧管理系統", "Artificial Intelligence Management System"),
    "MAITE": ("模型與 AI 測試評估基礎設施框架", "Model & AI Test and Evaluation Infrastructure Framework"),
    "NCSIST": ("國家中山科學研究院", "National Chung-Shan Institute of Science and Technology"),
    "RoE": ("交戰規則", "Rules of Engagement"),
    "ATLAS": ("人工智慧系統對抗威脅圖譜", "Adversarial Threat Landscape for Artificial-Intelligence Systems")
}

def add_acronym_footer(slide, acronym_keys, y_pos=7.95, height=0.9):
    NAVY = RGBColor(12, 35, 64)
    BLUE = RGBColor(37, 99, 235)
    MUTED = RGBColor(100, 116, 139)
    BG_FOOTER = RGBColor(241, 245, 249)
    BORDER_FOOTER = RGBColor(203, 213, 225)

    # Check if footer already exists or add new
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y_pos), Inches(14.4), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = BG_FOOTER
    card.line.color.rgb = BORDER_FOOTER
    card.line.width = Pt(1.0)

    tf = card.text_frame
    tf.word_wrap = True

    p_head = tf.paragraphs[0]
    p_head.text = "📌 頁面英文縮寫與專有名詞對照 (Acronym & Term Footnotes):"
    p_head.font.size = Pt(10)
    p_head.font.bold = True
    p_head.font.color.rgb = BLUE
    p_head.font.name = "微軟正黑體"
    p_head.space_after = Pt(2)

    footer_parts = []
    for k in acronym_keys:
        if k in ACRONYM_DB:
            zh, en = ACRONYM_DB[k]
            footer_parts.append(f"**{k}**: {zh} ({en})")

    footer_str = "  |  ".join(footer_parts)

    p_body = tf.add_paragraph()
    add_formatted_text(p_body, footer_str, font_size=9, default_color=MUTED, bold_color=NAVY)

def fix_presentation_slide2():
    prs_path = r'c:\Users\administartor\Downloads\AIEC\AIEC_AI_Evaluation_30_Slides_NanoBanana.pptx'
    prs = Presentation(prs_path)

    NAVY = RGBColor(12, 35, 64)
    BLUE = RGBColor(37, 99, 235)
    DARK_BLUE = RGBColor(30, 58, 138)

    slide2 = prs.slides[1]

    # Clean shapes in Slide 2 and rebuild Slide 2 layout cleanly
    # Find shape 8 and shape 12 or text frames with TOC content
    for shape in slide2.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text
            if "Table of Contents - Part 4 to 5" in txt or "Part 4 to 5" in txt or "4. " in txt or "5. " in txt:
                tf = shape.text_frame
                tf.clear()
                r_bullets = [
                    "4. **第四區塊：15 項國防級 AI 量化評測指標與 SOP (P.17 - 25)**",
                    "   - P.17 15項指標總覽 | P.18-25 Q1~Q15 LaTeX 公式、門檻與工具鏈",
                    "5. **第五區塊：評測平台與自動化 T&E 工具 (P.26 - 30)**",
                    "   - P.26 主權LLM架構 | P.27 地端推論 | P.28 C2與節點 | P.29 AI CLI | P.30 MAITE",
                    "6. **第六區塊：國防 AI 安全進階防禦與主權架構 (P.31 - 34)**",
                    "   - P.31 NCSIST AIEC 總圖藍圖與 5 大柱石 | P.32 國防 AIEC 任務與 RoE 邊界",
                    "   - P.33 戰術邊緣硬體防篡改與模型自毀 (<100ms) | P.34 地端模型蒸餾與 Provenance"
                ]
                add_formatted_bullets(tf, r_bullets, font_size=11.5)
            elif "Table of Contents - Part 1 to 3" in txt or "Part 1 to 3" in txt or "1. " in txt:
                tf = shape.text_frame
                tf.clear()
                l_bullets = [
                    "1. **第一區塊：評測框架與雙支柱治理 (P.3 - 6)**",
                    "   - P.3 治理雙支柱 | P.4 SHIELD 6階段 | P.5 ISO 42001 | P.6 MITRE ATLAS",
                    "2. **第二區塊：T&E 矩陣與評測方法論 (P.7 - 10)**",
                    "   - P.7 T&E 四層次 | P.8 JATIC 7構面 | P.9 6大方法論 | P.10 三維縱深矩陣",
                    "3. **第三區塊：六大類 AI 應用系統評測 SOP (P.11 - 16)**",
                    "   - P.11 CV目標偵測 | P.12 LLM生成式 | P.13 RAG知識庫 | P.14 Agent軌跡 | P.15 HMT與RoE | P.16 預測分析"
                ]
                add_formatted_bullets(tf, l_bullets, font_size=11.5)
            elif "Table of Contents - Part" in txt or "Part 4 to" in txt or "Part 1 to" in txt:
                pass

    prs.save(prs_path)
    print("Successfully updated Slide 2 Table of Contents with Section 6 (P.31-34)!")

if __name__ == '__main__':
    fix_presentation_slide2()
