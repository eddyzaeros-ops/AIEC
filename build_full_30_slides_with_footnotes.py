# -*- coding: utf-8 -*-
import os, sys
import collections
import collections.abc

if not hasattr(collections, 'Iterable'):
    collections.Iterable = collections.abc.Iterable

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from pptx_helpers import set_pure_white_bg, add_header, add_icon_card, add_formatted_bullets, add_formatted_text

# Master Acronym Database with Chinese and English Full Names
ACRONYM_DB = {
    "AIEC": ("人工智慧評測與認證體系", "Artificial Intelligence Evaluation & Certification"),
    "AI": ("人工智慧", "Artificial Intelligence"),
    "NIST": ("美國國家標準暨技術研究院", "National Institute of Standards and Technology"),
    "RMF": ("風險管理框架", "Risk Management Framework"),
    "DoD": ("美國國防部", "Department of Defense"),
    "CDAO": ("數位與人工智慧長辦公室", "Chief Digital and Artificial Intelligence Office"),
    "T&E": ("測試與評估", "Test and Evaluation"),
    "DT&E": ("發展性測試與評估", "Developmental Test and Evaluation"),
    "OT&E": ("作戰性測試與評估", "Operational Test and Evaluation"),
    "ISO": ("國際標準化組織", "International Organization for Standardization"),
    "AIMS": ("人工智慧管理系統", "Artificial Intelligence Management System"),
    "MITRE": ("邁特公司", "MITRE Corporation"),
    "ATLAS": ("人工智慧系統對抗威脅圖譜", "Adversarial Threat Landscape for Artificial-Intelligence Systems"),
    "ATT&CK": ("對抗戰術、技術與通用知識", "Adversarial Tactics, Techniques, and Common Knowledge"),
    "DAGR": ("國防人工智慧治理與風險指南", "Defense Artificial Intelligence Governance & Risk Guidelines"),
    "SHIELD": ("設定、精煉、改進、評估、記錄、偵測六大治理階段", "Set, Hone, Improve, Evaluate, Log, Detect"),
    "SOC": ("關切事項聲明", "Statement of Concern"),
    "RAI": ("負責任的人工智慧", "Responsible Artificial Intelligence"),
    "CMMC": ("網路安全成熟度模型認證", "Cybersecurity Maturity Model Certification"),
    "GDPR": ("通用資料保護規則", "General Data Protection Regulation"),
    "ML": ("機器學習", "Machine Learning"),
    "LLM": ("大語言模型", "Large Language Model"),
    "CV": ("電腦視覺", "Computer Vision"),
    "RAG": ("檢索增強生成", "Retrieval-Augmented Generation"),
    "HMT": ("人機協同戰術團隊", "Human-Autonomy Teaming"),
    "HSI": ("人機系統整合", "Human-Systems Integration"),
    "TEVV": ("測試、評估、驗證與確認", "Test, Evaluation, Verification, and Validation"),
    "JATIC": ("聯合人工智慧測試中心", "Joint AI Test Center"),
    "mAP": ("平均精確度均值", "mean Average Precision"),
    "OOD": ("分布外數據", "Out-of-Distribution"),
    "NRTK": ("自然穩健性測試工具包", "Natural Robustness Toolkit"),
    "XAITK": ("可解釋人工智慧工具包", "Explainable AI Toolkit"),
    "ART": ("對抗韌性測試工具包", "Adversarial Robustness Toolbox"),
    "HEART": ("高爆對抗紅隊測試工具", "High-Explosive Adversarial Red Teaming"),
    "FGSM": ("快速梯度符號攻擊法", "Fast Gradient Sign Method"),
    "PGD": ("投影梯度下降法", "Projected Gradient Descent"),
    "MSR": ("任務完成率", "Mission Success Rate"),
    "LVC": ("實戰-虛擬-構造平行戰場整合", "Live-Virtual-Constructive"),
    "VBS": ("虛擬戰場模擬軟體", "Virtual Battlespace"),
    "EADSIM": ("延伸防空模擬系統", "Extended Air Defense Simulation"),
    "DoDD": ("國防部指令", "Department of Defense Directive"),
    "ToAST": ("自主系統測試工具", "Testing of Autonomous Systems Tool"),
    "IDA": ("國防分析研究所", "Institute for Defense Analyses"),
    "MIT-LL": ("麻省理工學院林肯實驗室", "Massachusetts Institute of Technology Lincoln Laboratory"),
    "ECE": ("期望校準誤差", "Expected Calibration Error"),
    "EEG": ("腦電圖儀", "Electroencephalography"),
    "NASA-TLX": ("美國航太總署任務負荷指數", "NASA Task Load Index"),
    "SHAP": ("沙普利附加解釋值", "SHapley Additive exPlanations"),
    "LIME": ("局部可解釋模型無關說明", "Local Interpretable Model-agnostic Explanations"),
    "garak": ("生成式 AI 紅隊分析工具包", "Generative AI Redteam Analysis Kit"),
    "OWASP": ("開放 Web 應用程式安全計畫", "Open Web Application Security Project"),
    "RAGAS": ("檢索增強生成評估指標", "Retrieval Augmented Generation Assessment"),
    "API": ("應用程式介面", "Application Programming Interface"),
    "OPA": ("開放策略代理", "Open Policy Agent"),
    "SPIFFE": ("通用安全生產身份框架", "Secure Production Identity Framework for Everyone"),
    "SPIRE": ("SPIFFE 執行期環境", "SPIFFE Runtime Environment"),
    "PyOD": ("Python 異常值偵測庫", "Python Outlier Detection"),
    "UQ": ("不確定性量化", "Uncertainty Quantification"),
    "MC-Dropout": ("蒙地卡羅 Dropout", "Monte Carlo Dropout"),
    "RBAC": ("基於角色的存取控制", "Role-Based Access Control"),
    "MoE": ("專家混合架構", "Mixture of Experts"),
    "C2": ("指揮管制", "Command and Control"),
    "ISR": ("情報、監視與偵察", "Intelligence, Surveillance, and Reconnaissance"),
    "GUI": ("圖形使用者介面", "Graphical User Interface"),
    "GGUF": ("通用模型量化格式", "Georgi Gerganov Unified Format"),
    "MOC": ("內容主索引地圖", "Map of Content"),
    "Cron": ("定時排程服務", "Chronos Command Scheduler"),
    "MLOps": ("機器學習營運", "Machine Learning Operations")
}

def add_acronym_footer(slide, acronym_keys, y_pos=8.0, height=0.85):
    NAVY = RGBColor(12, 35, 64)
    BLUE = RGBColor(37, 99, 235)
    MUTED = RGBColor(100, 116, 139)
    BG_FOOTER = RGBColor(241, 245, 249)
    BORDER_FOOTER = RGBColor(203, 213, 225)

    # Footer Card Box
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y_pos), Inches(14.4), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = BG_FOOTER
    card.line.color.rgb = BORDER_FOOTER
    card.line.width = Pt(1.0)

    tf = card.text_frame
    tf.word_wrap = True

    p_head = tf.paragraphs[0]
    p_head.text = "📌 頁面英文縮寫對照 (Acronym Footnotes):"
    p_head.font.size = Pt(10)
    p_head.font.bold = True
    p_head.font.color.rgb = BLUE
    p_head.font.name = "微軟正黑體"
    p_head.space_after = Pt(2)

    footer_parts = []
    for k in acronym_keys:
        if k in ACRONYM_DB:
            zh, en = ACRONYM_DB[k]
            footer_parts.append(f"**{k}**: {zh} ({en})")

    footer_str = "  |  ".join(footer_parts)

    p_body = tf.add_paragraph()
    add_formatted_text(p_body, footer_str, font_size=9, default_color=MUTED, bold_color=NAVY)

def build_30_deck_with_footnotes():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank_layout = prs.slide_layouts[6]

    NAVY = RGBColor(12, 35, 64)
    BLUE = RGBColor(37, 99, 235)
    DARK_BLUE = RGBColor(30, 58, 138)
    CARD_BG = RGBColor(255, 255, 255)
    CARD_BORDER = RGBColor(218, 226, 236)
    TEXT_DARK = RGBColor(50, 50, 50)
    TEXT_MUTED = RGBColor(100, 116, 139)
    GREEN = RGBColor(16, 185, 129)

    # ----------------------------------------------------
    # Slide 1: Cover Slide
    # ----------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s1)

    cover_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.6), Inches(14.4), Inches(7.8))
    cover_card.fill.solid()
    cover_card.fill.fore_color.rgb = NAVY
    cover_card.line.fill.background()

    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(1.1), Inches(5.6), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = BLUE
    badge.line.fill.background()
    p = badge.text_frame.paragraphs[0]
    p.text = "🛡️ AIEC 國防與企業級 AI 評測與治理全集"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "微軟正黑體"

    tb1 = s1.shapes.add_textbox(Inches(1.3), Inches(1.9), Inches(13.4), Inches(2.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "AIEC 國防級 AI 評測體系、\n15 項量化指標與主權 AI 架構"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "微軟正黑體"

    tb_sub = s1.shapes.add_textbox(Inches(1.3), Inches(4.2), Inches(13.4), Inches(2.4))
    tf_sub = tb_sub.text_frame
    tf_sub.word_wrap = True
    cover_bullets = [
        "1. **權威標準引領**：結合 NIST AI RMF 1.0, DoD CDAO AI T&E, ISO 42001 AIMS 與 MITRE ATLAS",
        "2. **15 項量化評測 SOP**：涵蓋對抗韌性、自然穩健性、MSR、可中止性、信任校準、garak/RAGAS 測試",
        "3. **地端主權與戰術 C2**：整合 Gemma 4 四層 LLM 堆疊、Ollama/vLLM Middleware 與 Lattice/Menace 邊緣節點",
        "4. **Obsidian 第二大腦**：全套 21 篇雙向鏈結筆記實體建置於 G:\\我的雲端硬碟\\secondbrain\\AIEC"
    ]
    add_formatted_bullets(tf_sub, cover_bullets, font_size=15, text_color=RGBColor(226, 232, 240), bold_color=RGBColor(255, 255, 255))

    # Footnote Banner on Slide 1 Cover Card
    tb_foot1 = s1.shapes.add_textbox(Inches(1.3), Inches(6.8), Inches(13.4), Inches(1.2))
    tf_f1 = tb_foot1.text_frame
    tf_f1.word_wrap = True
    p_f1 = tf_f1.paragraphs[0]
    p_f1.text = "📌 頁面核心英文縮寫全稱 (Core Acronym Footnotes):"
    p_f1.font.size = Pt(11)
    p_f1.font.bold = True
    p_f1.font.color.rgb = BLUE
    p_f1.font.name = "微軟正黑體"

    s1_acronyms = ["AIEC", "NIST", "RMF", "DoD", "CDAO", "T&E", "ISO", "AIMS", "MITRE", "ATLAS", "LLM", "C2"]
    f1_str = "  |  ".join([f"**{k}**: {ACRONYM_DB[k][0]} ({ACRONYM_DB[k][1]})" for k in s1_acronyms if k in ACRONYM_DB])
    p_f1_b = tf_f1.add_paragraph()
    add_formatted_text(p_f1_b, f1_str, font_size=9, default_color=RGBColor(203, 213, 225), bold_color=RGBColor(255, 255, 255))

    # Helper function for 2-card slides with footnotes
    def build_2card_slide_fn(title, category, icon1, t1_zh, t1_en, b1, icon2, t2_zh, t2_en, b2, acronym_keys, accent1=BLUE, accent2=DARK_BLUE):
        s = prs.slides.add_slide(blank_layout)
        set_pure_white_bg(s)
        add_header(s, title, category)

        # Card 1 (Left) - Height reduced slightly to 6.1 inches to fit footer
        add_icon_card(s, 0.8, 1.6, 7.0, 6.25, icon1, t1_zh, t1_en, accent_color=accent1)
        tb_b1 = s.shapes.add_textbox(Inches(1.1), Inches(2.7), Inches(6.4), Inches(5.0))
        add_formatted_bullets(tb_b1.text_frame, b1, font_size=13.5)

        # Card 2 (Right)
        add_icon_card(s, 8.2, 1.6, 7.0, 6.25, icon2, t2_zh, t2_en, accent_color=accent2)
        tb_b2 = s.shapes.add_textbox(Inches(8.5), Inches(2.7), Inches(6.4), Inches(5.0))
        add_formatted_bullets(tb_b2.text_frame, b2, font_size=13.5)

        # Footer Box
        add_acronym_footer(s, acronym_keys, y_pos=7.95, height=0.9)
        return s

    # ----------------------------------------------------
    # Slide 2: Knowledge Vault Hierarchy & MOC
    # ----------------------------------------------------
    build_2card_slide_fn(
        "雙腦知識庫架構與 MOC 主索引頁", "SECTION 1: GOVERNANCE & ARCHITECTURE",
        "🧠", "1. Karpathy 雙腦流結構", "3-Tier Knowledge Architecture",
        [
            "1. **輸入層 (Clippings/)**：自動剪藏原始論文、軍事規範與講義",
            "2. **消化層 (知識庫/)**：AI 自動提煉概念、維護 index.md 百科目錄",
            "3. **輸出層 (創作庫/)**：生成專業評測報告、教案與資安 SOP"
        ],
        "🗂️", "2. AIEC 專題資料夾配置", "Vault Subfolder Hierarchy",
        [
            "00-Index_and_Templates/：**MOC 主索引與 3 大 Markdown 範本**",
            "01-治理與規範/：**ISO 42001, SHIELD, ATLAS, RAG 權限**",
            "02-評測矩陣與構面/：**15 項指標 SOP, 4 大能力層次, JATIC 7 構面**",
            "03-應用系統評測/：**A~F 六類 CV/LLM/RAG/Agent/HMT/決策評測**",
            "04-地端架構與邊緣算力/：**主權 LLM 堆疊, Ollama/vLLM, Lattice C2**"
        ],
        ["MOC", "AIEC", "ISO", "ATLAS", "RAG", "JATIC", "CV", "LLM", "HMT", "C2"]
    )

    # ----------------------------------------------------
    # Slide 3: AIEC Dual Pillars
    # ----------------------------------------------------
    build_2card_slide_fn(
        "AIEC 核心雙支柱：DAGR 風險指南與 SHIELD 治理引擎", "SECTION 1: GOVERNANCE & ARCHITECTURE",
        "⚖️", "1. DAGR 風險指南", "DAGR Risk Guidelines",
        [
            "1. **生命週期風險矩陣**：提供研發、測試、部署至退役的危害識別矩陣",
            "2. **安全關鍵審查 (Safety-Critical)**：劃分高風險任務的道德與法律審查門檻",
            "3. **與 ISO 42001 勾稽**：對接國際人工智慧管理系統 AIMS 風險控制項",
            "4. **責任制落實 (Responsible AI)**：明確定義 AI 出錯時的法律與管理責任"
        ],
        "🛡️", "2. SHIELD 治理引擎", "SHIELD Governance Engine",
        [
            "S. **Set Foundations**：設定倫理、法律、政策基礎與 SOC 聲明",
            "H. **Hone Operationalizations**：將政策轉化為量化 T&E 計畫",
            "I. **Improve & Innovate**：運用風險緩解工具處理 SOC 關切事項",
            "E. **Evaluate Status**：綜合評估危害解決程度與滿意度",
            "L. **Log for Traceability**：全程文件化追溯數據與模型演進歷史",
            "D. **Detect via Monitoring**：上線後持續監控效能衰減與數據漂移"
        ],
        ["AIEC", "DAGR", "SHIELD", "SOC", "T&E", "ISO", "AIMS", "RAI"]
    )

    # ----------------------------------------------------
    # Slide 4: SHIELD 6 Stages Detailed
    # ----------------------------------------------------
    build_2card_slide_fn(
        "SHIELD 六項治理循環活動 (S-H-I-E-L-D)", "SECTION 1: GOVERNANCE & ARCHITECTURE",
        "🔄", "1. 前半週期：基礎設定與評估", "Set, Hone & Improve Stages",
        [
            "1. **Set Foundations (設定基礎)**：辨識 RAI 基礎，產出關切事項聲明 (SOC)",
            "2. **Hone Operationalizations (精煉操作化)**：對照 [[T&E 四大能力層次]] 制定 SOP",
            "3. **Improve & Innovate (改進與創新)**：導入 NeMo Guardrails 與對抗防禦緩解 SOC"
        ],
        "🔍", "2. 後半週期：評估、追溯與監控", "Evaluate, Log & Detect Stages",
        [
            "4. **Evaluate Status (評估狀態)**：比對 [[JATIC 七大共通評測構面]] 量化結果",
            "5. **Log for Traceability (記錄可追溯)**：合規 CMMC L2，完整保留演進與 Log 軌跡",
            "6. **Detect via Monitoring (持續監控)**：即時偵測模型過期與概念漂移 (Concept Drift)"
        ],
        ["SHIELD", "RAI", "SOC", "T&E", "SOP", "JATIC", "CMMC"]
    )

    # ----------------------------------------------------
    # Slide 5: ISO 42001 AIMS
    # ----------------------------------------------------
    build_2card_slide_fn(
        "ISO 42001 (AIMS) 人工智慧管理系統與控制項對照", "SECTION 1: GOVERNANCE & ARCHITECTURE",
        "📜", "1. AIMS 核心管理要求", "ISO/IEC 42001 Core Requirements",
        [
            "1. **Clause 6.1.2 AI 風險評估**：針對全生命週期進行威脅建模與處置",
            "2. **Clause 8.4 透明度**：要求 AI 決策邏輯具備可追溯與可解釋性",
            "3. **Clause 9 績效評估**：要求建立定期內部稽核與高階管理審查流程"
        ],
        "🛡️", "2. Annex A 附錄控制項對映", "Annex A Controls Mapping",
        [
            "1. **偏見與公平性稽核**：審查訓練數據，防止對特定群體的隱性歧視",
            "2. **數據治理 (Data Governance)**：確保數據來源合法性與隱私保護 (GDPR)",
            "3. **ISO 17025 方法確效**：評測工具鏈與 SOP 必須經過標準檢定與確效"
        ],
        ["ISO", "AIMS", "AI", "ATLAS", "XAITK", "CV", "SHAP", "LIME", "GDPR", "SOP"]
    )

    # ----------------------------------------------------
    # Slide 6: MITRE ATLAS Threat Matrix
    # ----------------------------------------------------
    build_2card_slide_fn(
        "MITRE ATLAS 人工智慧資安威脅矩陣 (對照 ATT&CK)", "SECTION 1: GOVERNANCE & ARCHITECTURE",
        "⚔️", "1. ATT&CK vs. ATLAS 對比", "Traditional IT vs. AI/ML Security",
        [
            "1. **關注對象**：ATT&CK 關注傳統 IT 網路；ATLAS 關注 AI/ML 系統與數據鏈",
            "2. **攻擊面**：ATLAS 聚焦模型權重、訓練資料集、Prompt 注入與向量庫",
            "3. **典型手法**：Data Poisoning, Jailbreak, Model Inversion, Adversarial Patch"
        ],
        "🎯", "2. 國防 AI 主要 ATLAS 防禦", "Defense ATLAS Tactics & Mitigations",
        [
            "1. **對抗干擾防禦**：針對邊緣 CV 的對抗貼片，實施 [[A類 - 電腦視覺與目標偵測評測]]",
            "2. **Prompt 越獄防禦**：針對 LLM 指管對答，導入 [[B類 - 生成式 AI 與大語言模型評測]]",
            "3. **經檢索注入防禦**：針對 RAG 向量庫，實施 [[RAG 權限控管與資料分級稽核]]"
        ],
        ["MITRE", "ATLAS", "ATT&CK", "ML", "LLM", "RAG", "CV", "SOP"]
    )

    # ----------------------------------------------------
    # Slide 7: T&E Four Capability Axes
    # ----------------------------------------------------
    build_2card_slide_fn(
        "T&E 四大能力層次 (Capability Axes)", "SECTION 2: T&E MATRIX & METHODOLOGY",
        "📐", "1. Level 1 & Level 2 評測", "Model & HSI T&E Level",
        [
            "1. **Level 1 Model T&E (模型測評)**：單體模型演算法效能、對抗攻防、校準度與偏見",
            "2. **Level 2 HSI T&E (人機整合)**：Human-Systems Integration，測試認知負荷、信任度與過度依賴"
        ],
        "🌐", "2. Level 3 & Level 4 評測", "Systems & Operational T&E Level",
        [
            "3. **Level 3 Systems Integration (系統整合)**：端到端數據鏈串接、API 閘道穩定度與 Agent 軌跡",
            "4. **Level 4 Operational T&E (作戰測評)**：電戰干擾適應力、環境漂移韌性、緊急脱離與停用機制"
        ],
        ["T&E", "HSI", "TEVV", "API", "C2"]
    )

    # ----------------------------------------------------
    # Slide 8: JATIC 7 Common Dimensions
    # ----------------------------------------------------
    build_2card_slide_fn(
        "JATIC 七大跨系統共通評測構面", "SECTION 2: T&E MATRIX & METHODOLOGY",
        "🎯", "1. 共通技術基底 (構面 1-4)", "Robustness, Resiliency & Competence",
        [
            "1. **穩健性 (Robustness)**：在面對噪聲、OOD 輸入與對抗樣本時維持效能",
            "2. **韌性 (Resiliency)**：受到網路攻擊或電戰干擾時自動降級或安全復原",
            "3. **可解釋性 (Explainability)**：提供特徵歸因與熱力圖，消除黑盒子",
            "4. **勝任度 (Competence)**：在其指定作戰邊界內的任務完成率與精度"
        ],
        "⚖️", "2. 信任與治理基底 (構面 5-7)", "Fairness, Calibration & Drift",
        [
            "5. **公平性 (Fairness)**：防範訓練數據中的隱性偏見與群體偏差",
            "6. **校準 (Calibration)**：信心度與實際正確率完全吻合 (ECE <= 0.05)",
            "7. **漂移監控 (Drift Monitoring)**：上線後持續追蹤概念漂移與數據衰減"
        ],
        ["JATIC", "AIEC", "OOD", "ECE"]
    )

    # ----------------------------------------------------
    # Slide 9: 6 T&E Methodologies
    # ----------------------------------------------------
    build_2card_slide_fn(
        "國防 AI 評測 6 大方法論與匹配矩陣", "SECTION 2: T&E MATRIX & METHODOLOGY",
        "🧪", "1. 開放與封閉測試 (黑箱/白箱/基準)", "Black-box, White-box & Benchmark",
        [
            "1. **黑箱測試 (Black-box)**：不存取權重，評測輸入輸出行為 (雲端 Tier-4 Gated Claude)",
            "2. **白箱測試 (White-box)**：存取梯度與特徵圖 (地端 Gemma 4 各層模型)",
            "3. **基準測試 (Benchmarking)**：使用標準化數據集量化比較基礎能力 (AgentBench)"
        ],
        "⚔️", "2. 安全與營運測試 (紅隊/人工/持續)", "Red Teaming, Human & Continuous",
        [
            "4. **對抗/紅隊 (Red Teaming)**：模擬敵方進行干擾與越獄攻擊 (IBM ART 360, garak)",
            "5. **人工評估 (Human Eval)**：主觀無真相時的評估 (HMT Guidebook)",
            "6. **持續監控 (Continuous)**：部署後即時偵測性能衰減 (Arize Phoenix, PyOD)"
        ],
        ["AIEC", "SOP", "API", "LLM", "garak", "XAITK", "SHAP", "ART", "HMT", "SHIELD", "PyOD"]
    )

    # ----------------------------------------------------
    # Slide 10: RAG Security & Data Classification
    # ----------------------------------------------------
    build_2card_slide_fn(
        "RAG 權限控管與資料分級防降密稽核機制", "SECTION 2: T&E MATRIX & METHODOLOGY",
        "🔐", "1. 向量庫 RBAC 與分級 Tag", "Vector Database Access Control",
        [
            "1. **向量存儲層加標 (Metadata Tagging)**：於 Embedding 附加密級（機密、極機密、限制級）",
            "2. **強制管道憑證查驗**：檢索管道執行 Query 前，強制校驗用戶與 Agent 之數位證書",
            "3. **供應鏈安全審查**：審查開源 Embedder 模型與向量庫（Milvus, Qdrant）無後門"
        ],
        "🛡️", "2. 防範 LLM 摘要降密洩漏", "Anti-Declassification Leakage",
        [
            "1. **跨文件摘要降密風險**：多篇限制級文件經 LLM 統整後可能推導出機密結論",
            "2. **動態輸出遮罩 (Output Masking)**：根據用戶權限自動遮蔽高密級敏感實體",
            "3. **對照 ISO 42001 Annex A**：滿足機密性與權限隔離之合規稽核要求"
        ],
        ["RAG", "RBAC", "AIEC", "LLM", "ISO", "AIMS"]
    )

    # Helper function for System T&E slides (A~F) with footnotes
    def build_system_slide_fn(title, category, sys_code, sys_name, icon, risks, sop_steps, tools, acronym_keys):
        s = prs.slides.add_slide(blank_layout)
        set_pure_white_bg(s)
        add_header(s, title, category)

        add_icon_card(s, 0.8, 1.6, 7.0, 6.25, icon, f"{sys_code} - {sys_name}", "System Risk Profile & Scope", accent_color=BLUE)
        tb_r = s.shapes.add_textbox(Inches(1.1), Inches(2.7), Inches(6.4), Inches(5.0))
        add_formatted_bullets(tb_r.text_frame, risks, font_size=13.5)

        add_icon_card(s, 8.2, 1.6, 7.0, 6.25, "🛠️", "評測 SOP 與工具鏈", "Testing Methodology & Toolkits", accent_color=DARK_BLUE)
        tb_s = s.shapes.add_textbox(Inches(8.5), Inches(2.7), Inches(6.4), Inches(5.0))
        items = ["**測試實施 SOP**："] + ["  " + step for step in sop_steps] + ["**代表性工具鏈**："] + ["  " + tool for tool in tools]
        add_formatted_bullets(tb_s.text_frame, items, font_size=13.5)

        add_acronym_footer(s, acronym_keys, y_pos=7.95, height=0.9)
        return s

    # ----------------------------------------------------
    # Slide 11: System A - Computer Vision
    # ----------------------------------------------------
    build_system_slide_fn(
        "A類 - 電腦視覺與目標偵測評測 (CV & Target Detection)", "SECTION 3: SYSTEM-SPECIFIC T&E",
        "A類", "電腦視覺與目標偵測", "👁️",
        [
            "1. **對抗貼片 (Adversarial Patch)**：敵方貼附對抗圖樣引發欺騙誤判",
            "2. **自然穩健性下降**：雨雪、煙霧、電戰雜訊導致 mAP 暴降",
            "3. **分布外 (OOD) 偽裝目標**：新型偽裝目標被誤標為高威脅或背景"
        ],
        ["1. 白箱/黑箱對抗攻防：測試對抗干擾下目標框 (Bounding Box) 穩定度", "2. 顯著性可解釋性：產出 Heatmap 確保聚焦本體而非背景"],
        ["HEART (High-Explosive Adversarial Red Teaming)", "NRTK (Natural Robustness Toolkit)", "XAITK (Explainable AI Toolkit)", "IBM ART 360"],
        ["CV", "mAP", "OOD", "SOP", "HEART", "NRTK", "XAITK", "ART"]
    )

    # ----------------------------------------------------
    # Slide 12: System B - GenAI & LLM
    # ----------------------------------------------------
    build_system_slide_fn(
        "B類 - 生成式 AI 與大語言模型評測 (GenAI & LLM)", "SECTION 3: SYSTEM-SPECIFIC T&E",
        "B類", "生成式 AI 與 LLM", "💬",
        [
            "1. **幻覺 (Hallucination)**：生成虛構軍事情報與錯誤數據",
            "2. **越獄 (Jailbreak)**：對抗 Prompt 繞過系統安全護欄",
            "3. **敏感資料洩露**：訓練集或機密脈絡遭 Prompt Inversion 導出"
        ],
        ["1. 紅隊模糊測試 (Red-teaming Fuzzing)：自動發送萬筆越獄 Payload", "2. 護欄驗證：測試 NeMo 護欄攔截成功率與反應時間"],
        ["garak (LLM Vulnerability Scanner)", "NeMo Guardrails (NVIDIA 輸出入護欄)", "PromptBench & TrustLLM"],
        ["AI", "LLM", "garak", "SOP"]
    )

    # ----------------------------------------------------
    # Slide 13: System C - RAG
    # ----------------------------------------------------
    build_system_slide_fn(
        "C類 - 檢索增強生成 RAG 評測 (Retrieval-Augmented Gen)", "SECTION 3: SYSTEM-SPECIFIC T&E",
        "C類", "檢索增強生成 RAG", "📚",
        [
            "1. **知識衝突 (Knowledge Conflict)**：檢索不相關段落導致回答離題",
            "2. **經檢索注入 (Retrieved Injection)**：數據庫內植入對抗指令操控 LLM",
            "3. **來源歸屬錯誤**：引述錯誤的技術手冊或規範章節"
        ],
        ["1. RAG Triad 三元組評估：脈絡精確度、忠實度、答案相關度", "2. 軌跡追蹤：記錄向量搜尋 Top-K 比對分頁與時間戳"],
        ["RAGAS (RAG Assessment)", "TruLens Evaluation Framework", "Arize Phoenix (Tracing & Eval)"],
        ["RAG", "LLM", "RAGAS", "SOP"]
    )

    # ----------------------------------------------------
    # Slide 14: System D - AI Agent
    # ----------------------------------------------------
    build_system_slide_fn(
        "D類 - AI Agent 與多代理系統評測 (Multi-Agent Systems)", "SECTION 3: SYSTEM-SPECIFIC T&E",
        "D類", "AI Agent 與多代理系統", "🤖",
        [
            "1. **工具誤用 (Tool Misuse)**：非法 API 調用或執行錯誤命令",
            "2. **目標偏移 (Goal Drift)**：多輪互動中偏離原始任務目標或陷入死迴圈",
            "3. ** Agent 間衝突**：多代理協同決策時發生邏輯死鎖或權限爭奪"
        ],
        ["1. 軌跡稽核 (Trajectory Auditing)：記錄完整 API 呼叫鏈與中介狀態", "2. HITL 閘門測試：驗證關鍵開火/變更授權點之強制介入能力"],
        ["AgentBench (Agent Benchmark)", "SPIFFE/SPIRE (Agent 身份與證書)", "Open Policy Agent (OPA) 策略閘門"],
        ["AI", "API", "HITL", "SPIFFE", "SPIRE", "OPA", "SOP"]
    )

    # ----------------------------------------------------
    # Slide 15: System E - Autonomous Systems & HMT
    # ----------------------------------------------------
    build_system_slide_fn(
        "E類 - 自主系統與人機協同評測 (Human-Autonomy Teaming)", "SECTION 3: SYSTEM-SPECIFIC T&E",
        "E類", "自主系統與人機協同", "🚁",
        [
            "1. **非預期自主行為**：蜂群或無人車脫離指定邊界且無法緊急中止",
            "2. **過度依賴 (Over-reliance)**：操作員對 AI 盲目信任導致無警覺",
            "3. **認知負荷過載**：介面警報過多引發操作員決策恐慌"
        ],
        ["1. 遵循 DoDD 3000.09：美軍自主武器指令評測規範", "2. 脫離停用測試：測試手動 Stop Signal 及安全降級接管率"],
        ["ToAST (Testing of Autonomous Systems Tool)", "IDA HMT Guidebook", "MIT-LL HMT Testing Guide"],
        ["HMT", "AI", "DoDD", "ToAST", "IDA", "MIT-LL", "SOP"]
    )

    # ----------------------------------------------------
    # Slide 16: System F - Predictive Analytics
    # ----------------------------------------------------
    build_system_slide_fn(
        "F類 - 決策支援與預測分析評測 (Predictive Analytics)", "SECTION 3: SYSTEM-SPECIFIC T&E",
        "F類", "決策支援與預測分析", "📈",
        [
            "1. **概念漂移 (Concept Drift)**：戰術態勢變化導致歷史模型失效",
            "2. **不確定性未量化**：模型給出高信心度但實際為 OOD 預測",
            "3. **隱性群體偏見**：後勤或威脅排序模型受到歷史數據偏差干擾"
        ],
        ["1. 漂移與 OOD 偵測：即時計算數據分布變化 (PyOD)", "2. 不確定性量化 (UQ)：生成信心區間並進行特徵歸因分析"],
        ["AIF360 (AI Fairness 360)", "PyOD / Alibi Detect", "SHAP / LIME (Feature Attribution)"],
        ["OOD", "UQ", "PyOD", "SHAP", "LIME", "SOP"]
    )

    # Helper function for 15 Metrics pairs (Slides 18-24) with footnotes
    def build_metric_pair_slide_fn(title, category, m1, m2, acronym_keys):
        s = prs.slides.add_slide(blank_layout)
        set_pure_white_bg(s)
        add_header(s, title, category)

        metrics = [m1, m2]
        x_offsets = [0.8, 8.2]

        for idx, m in enumerate(metrics):
            x = x_offsets[idx]
            
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(7.0), Inches(6.25))
            card.fill.solid()
            card.fill.fore_color.rgb = CARD_BG
            card.line.color.rgb = CARD_BORDER
            card.line.width = Pt(1.5)

            ib = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.2), Inches(1.8), Inches(0.95), Inches(0.95))
            ib.fill.solid()
            ib.fill.fore_color.rgb = DARK_BLUE
            ib.line.fill.background()
            p = ib.text_frame.paragraphs[0]
            p.text = f"Q{m['id']}"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = "Arial"

            tb_t = s.shapes.add_textbox(Inches(x + 1.25), Inches(1.75), Inches(5.5), Inches(1.0))
            tf_t = tb_t.text_frame
            tf_t.word_wrap = True
            p1 = tf_t.paragraphs[0]
            p1.text = m['name_zh']
            p1.font.size = Pt(18)
            p1.font.bold = True
            p1.font.color.rgb = NAVY
            p1.font.name = "微軟正黑體"

            p2 = tf_t.add_paragraph()
            p2.text = m['name_en']
            p2.font.size = Pt(13)
            p2.font.color.rgb = TEXT_MUTED
            p2.font.name = "Arial"

            tb_c = s.shapes.add_textbox(Inches(x + 0.3), Inches(2.9), Inches(6.4), Inches(4.8))
            tf_c = tb_c.text_frame
            tf_c.word_wrap = True

            bullet_items = [
                f"**1. 指標定義**：{m['def']}",
                f"**2. 驗測 SOP 與工具**：{m['sop']}",
                f"**3. 量化合格門檻**：{m['thresh']}",
                f"**4. 對應國際標準**：{m['std']}"
            ]
            add_formatted_bullets(tf_c, bullet_items, font_size=12.5)

        add_acronym_footer(s, acronym_keys, y_pos=7.95, height=0.9)
        return s

    # ----------------------------------------------------
    # Slide 17: 15 Metrics Master Overview
    # ----------------------------------------------------
    build_2card_slide_fn(
        "15 項國防級 AI 量化評測指標架構總覽", "SECTION 4: 15 QUANTITATIVE METRICS & SOPS",
        "📊", "1. 三大驗測視角劃分", "Three Operational Perspectives",
        [
            "1. **作戰與環境效能 (Q1-Q4)**：對抗韌性、自然穩健性、任務完成率 (MSR)、可中止性",
            "2. **情境與模型能力 (Q5-Q10)**：信任校準、認知負荷、可解釋性、越獄防禦、幻覺率、RAG精確度",
            "3. **稽核、資安與治理 (Q11-Q15)**：Agent軌跡、漂移監控、不確定性量化、防降密洩漏、可追溯性"
        ],
        "⚖️", "2. 貫穿四大標準規範", "Aligned International Frameworks",
        [
            "1. **NIST AI RMF 1.0 (Measure)**：定義 Valid, Safe, Secure, Accountable 核心量化指標",
            "2. **DoD CDAO AI T&E Guidebook**：提供 DT&E / OT&E 實戰情境測試 SOP",
            "3. **ISO 42001 AIMS**：提供內部控制、偏見稽核與 Clause 9 持續監控要求",
            "4. **MITRE ATLAS**：提供對抗攻防與威脅向量對照防禦矩陣"
        ],
        ["MSR", "garak", "RAGAS", "UQ", "NIST", "RMF", "DoD", "CDAO", "T&E", "DT&E", "OT&E", "ISO", "AIMS", "MITRE", "ATLAS"]
    )

    # ----------------------------------------------------
    # Slide 18: Q1 & Q2
    # ----------------------------------------------------
    m1 = {
        "id": 1, "name_zh": "對抗韌性", "name_en": "Adversarial Robustness",
        "def": "模型遭受對抗貼片、FGSM/PGD 擾動攻擊時維持正確判讀能力",
        "sop": "使用 IBM ART 360 / HEART 對模型注入 ε 擾動，測試 mAP 變化",
        "thresh": "Acc_adv / Acc_clean >= 90% (於 ε <= 0.05 條件下)",
        "std": "MITRE ATLAS / NIST AI RMF 1.0"
    }
    m2 = {
        "id": 2, "name_zh": "自然穩健性", "name_en": "Natural Robustness",
        "def": "模型面對自然環境干擾（雨雪、煙霧、電戰雜訊）時的效能維持度",
        "sop": "透過 NRTK 合成 10 種等級的環境降質數據集進行壓力測試",
        "thresh": "mAP 衰減率 <= 10% (高噪聲測試條件下)",
        "std": "JATIC / DoD CDAO AI T&E"
    }
    build_metric_pair_slide_fn("15 項評測指標 (Q1 - Q2) —— 對抗與自然穩健性", "SECTION 4: 15 QUANTITATIVE METRICS & SOPS", m1, m2, ["FGSM", "PGD", "mAP", "ART", "HEART", "NRTK", "MITRE", "ATLAS", "NIST", "RMF", "JATIC", "DoD", "CDAO", "T&E"])

    # ----------------------------------------------------
    # Slide 19: Q3 & Q4
    # ----------------------------------------------------
    m3 = {
        "id": 3, "name_zh": "任務完成率", "name_en": "Mission Success Rate (MSR)",
        "def": "AI 系統在端到端戰術情境中成功執行完畢並閉合擊殺鏈的比例",
        "sop": "於 VBS 4 / EADSIM 虛實整合 (LVC) 平行戰場環境執行 100 次模擬",
        "thresh": "MSR = 成功次數 / 總模擬數 >= 95%",
        "std": "Level 4 Operational T&E"
    }
    m4 = {
        "id": 4, "name_zh": "可中止性與失效安全", "name_en": "Abortability & Fail-Safe Rate",
        "def": "當系統異常或接獲人工中斷指令時，即刻中斷並進入安全保護狀態",
        "sop": "隨機注入手動 Stop Signal 及硬體斷連，量測安全降級接管時間",
        "thresh": "Abort Latency <= 100ms (Fail-Safe Rate = 100%)",
        "std": "DoDD 3000.09 自主武器指令"
    }
    build_metric_pair_slide_fn("15 項評測指標 (Q3 - Q4) —— 任務完成與失效安全", "SECTION 4: 15 QUANTITATIVE METRICS & SOPS", m3, m4, ["MSR", "LVC", "VBS", "EADSIM", "DoDD", "T&E"])

    # ----------------------------------------------------
    # Slide 20: Q5 & Q6
    # ----------------------------------------------------
    m5 = {
        "id": 5, "name_zh": "信任校準與過度依賴", "name_en": "Trust Calibration & Over-Reliance",
        "def": "操作員對 AI 信心度的理解符合實際能力，防止盲目信任或拒絕使用",
        "sop": "於 HMT 模擬試驗中故意提供高信心但錯誤提案，記錄操作員修正率",
        "thresh": "ECE <= 0.05 / Over-reliance Rate <= 5%",
        "std": "JATIC / DoD HMT Guidebook"
    }
    m6 = {
        "id": 6, "name_zh": "認知負荷與適應性", "name_en": "Cognitive Load & Adaptability",
        "def": "AI 介面輸出對指揮官或操作員造成的心理負荷程度與決策時延",
        "sop": "操作員配戴眼動儀與 EEG 完成任務後填寫 NASA-TLX 量表",
        "thresh": "NASA-TLX Score Drop >= 30% (Decision Delay <= 2s)",
        "std": "Level 2 HSI T&E"
    }
    build_metric_pair_slide_fn("15 項評測指標 (Q5 - Q6) —— 信任校準與認知負荷", "SECTION 4: 15 QUANTITATIVE METRICS & SOPS", m5, m6, ["HMT", "ECE", "EEG", "NASA-TLX", "JATIC", "DoD", "HSI", "T&E"])

    # ----------------------------------------------------
    # Slide 21: Q7 & Q8
    # ----------------------------------------------------
    m7 = {
        "id": 7, "name_zh": "模型可解釋性與顯著性歸因", "name_en": "Explainability & Point Game",
        "def": "AI 關鍵決策邏輯機能是否提供可被人類審計的特徵熱力圖 (Saliency)",
        "sop": "白箱調用 XAITK / SHAP / LIME 產出熱力圖，比對真實目標區域",
        "thresh": "Point Game Score >= 0.85",
        "std": "ISO 42001 Clause 8.4"
    }
    m8 = {
        "id": 8, "name_zh": "提示越獄與抗注入能力", "name_en": "Prompt Jailbreak Defense Rate",
        "def": "LLM 阻絕敵方對抗 Prompt 注入、越獄繞過與護欄突圍的能力",
        "sop": "使用 garak 框架執行 10,000 筆測試案例 (Direct/Indirect Injection)",
        "thresh": "Jailbreak Defense Rate >= 99%",
        "std": "OWASP LLM Top 10 / garak"
    }
    build_metric_pair_slide_fn("15 項評測指標 (Q7 - Q8) —— 可解釋性與越獄防禦", "SECTION 4: 15 QUANTITATIVE METRICS & SOPS", m7, m8, ["XAITK", "SHAP", "LIME", "LLM", "garak", "OWASP", "ISO"])

    # ----------------------------------------------------
    # Slide 22: Q9 & Q10
    # ----------------------------------------------------
    m9 = {
        "id": 9, "name_zh": "幻覺率與事實忠實度", "name_en": "Hallucination Rate & Faithfulness",
        "def": "LLM 產出內容嚴格遵循檢索脈絡與國防事實，無虛構捏造數據",
        "sop": "運用 RAGAS 與 TruLens 的 Faithfulness 評估器進行自動事實核對",
        "thresh": "Faithfulness Score >= 0.95 (Hallucination <= 2%)",
        "std": "NIST AI RMF / RAGAS"
    }
    m10 = {
        "id": 10, "name_zh": "檢索精確度與來源歸屬", "name_en": "RAG Context Precision & Attribution",
        "def": "RAG 向量資料庫精確檢索權威規範段落並準確標註出處來源",
        "sop": "比對 RAG 檢索出的 Top-K 段落與 Ground Truth 之語意相關性",
        "thresh": "Context Precision >= 0.90 / Attribution >= 0.98",
        "std": "ISO 42001 / TruLens"
    }
    build_metric_pair_slide_fn("15 項評測指標 (Q9 - Q10) —— 幻覺控制與 RAG 精確度", "SECTION 4: 15 QUANTITATIVE METRICS & SOPS", m9, m10, ["LLM", "RAG", "RAGAS", "NIST", "RMF", "ISO"])

    # ----------------------------------------------------
    # Slide 23: Q11 & Q12
    # ----------------------------------------------------
    m11 = {
        "id": 11, "name_zh": "Agent 工具調用與軌跡合規", "name_en": "Agent Trajectory Audit",
        "def": "自主 AI Agent 呼叫外部 API 時嚴格遵循權限邊界，無目標偏移",
        "sop": "使用 AgentBench 記錄 Tool Call 軌跡，經由 OPA 進行策略比對",
        "thresh": "Unauthorized API Call Rate = 0% (Success >= 98%)",
        "std": "OPA / SPIFFE / AgentBench"
    }
    m12 = {
        "id": 12, "name_zh": "概念與數據漂移監控率", "name_en": "Data & Concept Drift Recall",
        "def": "系統在上線營運期間，即時捕捉數據分布變化與標籤漂移的靈敏度",
        "sop": "部署 PyOD 與 Alibi Detect 警報模組，注入漂移數據集測試反應時延",
        "thresh": "Drift Detection Recall >= 95% (Alarm <= 5min)",
        "std": "SHIELD Detect Stage / PyOD"
    }
    build_metric_pair_slide_fn("15 項評測指標 (Q11 - Q12) —— Agent 合規與漂移監控", "SECTION 4: 15 QUANTITATIVE METRICS & SOPS", m11, m12, ["API", "OPA", "SPIFFE", "PyOD", "SHIELD"])

    # ----------------------------------------------------
    # Slide 24: Q13 & Q14
    # ----------------------------------------------------
    m13 = {
        "id": 13, "name_zh": "不確定性量化", "name_en": "Uncertainty Quantification (UQ)",
        "def": "模型對預測結果給出可靠信心區間，遇高不確定性時提示人類",
        "sop": "採用 MC-Dropout 或 Deep Ensembles 生成方差，測試 OOD 方差激增度",
        "thresh": "OOD Variance Coverage >= 95%",
        "std": "NIST AI RMF / PyOD"
    }
    m14 = {
        "id": 14, "name_zh": "資料分級與防降密洩漏", "name_en": "Anti-Declassification Leakage",
        "def": "多密級檢索時，防止低權限用戶或 LLM 摘要統整導出降密高密級資訊",
        "sop": "模擬不同密級用戶對 RAG 探勘，查驗輸出遮罩與 RBAC 向量標籤",
        "thresh": "Declassification Leakage Rate = 0%",
        "std": "ISO 42001 Annex A / RBAC"
    }
    build_metric_pair_slide_fn("15 項評測指標 (Q13 - Q14) —— 不確定性與防降密洩漏", "SECTION 4: 15 QUANTITATIVE METRICS & SOPS", m13, m14, ["MC-Dropout", "OOD", "UQ", "RAG", "LLM", "RBAC", "ISO", "AIMS", "PyOD", "NIST", "RMF"])

    # ----------------------------------------------------
    # Slide 25: Q15
    # ----------------------------------------------------
    m15 = {
        "id": 15, "name_zh": "系統軌跡可追溯性與可稽核性", "name_en": "Traceability & Audit Compliance",
        "def": "AI 系統全生命週期的數據、權重、Prompt、API 軌跡與審核紀錄皆能完整追溯與合規重現",
        "sop": "抽查歷史決策紀錄，驗證是否能從日誌中重新推導並還原模型當時的推論歷程 (CMMC L2 / ISO 42001)",
        "thresh": "Log Audit Coverage = 100% (Reproduction Latency <= 10min)",
        "std": "SHIELD Log Stage / CMMC L2"
    }
    s25 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s25)
    add_header(s25, "15 項評測指標 (Q15) —— 系統軌跡可追溯性與可稽核性", "SECTION 4: 15 QUANTITATIVE METRICS & SOPS")
    
    add_icon_card(s25, 0.8, 1.6, 14.4, 6.25, "📜", "Q15 - 系統軌跡可追溯性與可稽核性 (Traceability & Audit Compliance)", "CMMC Level 2 & ISO 42001 Full Auditability", accent_color=DARK_BLUE)
    tb_q15 = s25.shapes.add_textbox(Inches(1.1), Inches(2.7), Inches(13.8), Inches(5.0))
    add_formatted_bullets(tb_q15.text_frame, [
        f"**1. 指標定義與背景**：{m15['def']}",
        f"**2. 驗測 SOP 與方法**：{m15['sop']}",
        f"**3. 量化合格門檻**：{m15['thresh']}",
        f"**4. 對應國際標準**：{m15['std']}",
        "**5. 日誌鏈結規範**：包含用戶 Identity、Prompt 版本、RAG 檢索 Chunk Hash、LLM 權限 Hash 與決策輸出之全鏈結封裝"
    ], font_size=14.5)

    add_acronym_footer(s25, ["API", "RAG", "LLM", "CMMC", "ISO", "AIMS", "SHIELD"], y_pos=7.95, height=0.9)

    # ----------------------------------------------------
    # Slide 26: Sovereign AI Platform & 4-Tier LLM Stack
    # ----------------------------------------------------
    build_2card_slide_fn(
        "主權 AI 平台與四層 LLM 堆疊 (Tier 1~Tier 4)", "SECTION 5: SOVEREIGN AI & EDGE C2",
        "🏰", "1. 4 層 LLM 軟硬體配置", "Sovereign LLM Architecture Stack",
        [
            "1. **Tier 1 核心**：Gemma 4 31B Dense (NVIDIA H200 / CUDA) — 白箱對抗測試",
            "2. **Tier 2 混合**：Gemma 4 26B MoE (AMD MI325X / ROCm) — 白箱對抗測試",
            "3. **Tier 3 邊緣**：Gemma 4 E4B 本地推論 (Menace 節點) — 白箱基準測試",
            "4. **Tier 4 雲端**：Gated Cloud Claude API (HITL 門禁閘門) — 黑箱基準測試"
        ],
        "⚙️", "2. 白箱 vs. 黑箱測試配置", "T&E Assignment by Tier",
        [
            "1. **地端模型 (Tier 1-3)**：採 **白箱測試 + 對抗紅隊**，存取梯度與中介特徵",
            "2. **雲端模型 (Tier 4)**：採 **黑箱測試 + 基準測試**，嚴格進行 API 閘門過濾",
            "3. **代理層 (Agent Layer)**：全堆疊納入 AgentBench 軌跡稽核與 OPA 策略管束"
        ],
        ["AI", "LLM", "MoE", "API", "HITL", "OPA"]
    )

    # ----------------------------------------------------
    # Slide 27: Local LLM Inference Engines & Middleware
    # ----------------------------------------------------
    build_2card_slide_fn(
        "地端 LLM 推論引擎與 Middleware 工具鏈", "SECTION 5: SOVEREIGN AI & EDGE C2",
        "🛠️", "1. 地端推論工具鏈對比", "Inference Engines & Frameworks",
        [
            "1. **Ollama**：極簡指令集與 Docker 支援，適合開發者與 Modelfile 自訂",
            "2. **LM Studio**：精美桌面視窗 GUI，一鍵搜尋下載 Hugging Face GGUF 模型",
            "3. **llama.cpp**：地端 LLM 推論底層 C++ 核心引擎，具備最高效能與硬體控制權",
            "4. **vLLM**：企業生產環境 High-Throughput 推論框架，支援 PagedAttention"
        ],
        "📚", "2. Hugging Face 與 GGUF 量化", "Model Hub & Quantization",
        [
            "1. **Hugging Face Model Hub**：全球 AI 社群模型圖書館，權重檔下載來源",
            "2. **GGUF 量化格式**：優化地端晶片 (Mac M系列 / NVIDIA 顯卡) 推論速度",
            "3. **Modelfile 封裝**：定義 Temperature, System Prompt, Context Window 參數"
        ],
        ["LLM", "API", "GUI", "GGUF"]
    )

    # ----------------------------------------------------
    # Slide 28: Lattice C2 & Menace Edge Node
    # ----------------------------------------------------
    build_2card_slide_fn(
        "Lattice 軟體定義戰術 C2 架構與 Menace 邊緣節點", "SECTION 5: SOVEREIGN AI & EDGE C2",
        "⚡", "1. Lattice C2 架構", "Software-Defined C2 Architecture",
        [
            "1. **任務工作流 (Mission Workflows)**：涵蓋 ISR, Strike Broker, High-Value Targeting",
            "2. **資料模型與 API**：外部 Data Model API 與內部 Pub/Sub API 整合",
            "3. **知識圖譜 (Knowledge Graph)**：整合 Entity Management 與 Sensor Fusion"
        ],
        "💻", "2. Menace 邊緣算力節點", "Menace Tactical Edge Node",
        [
            "1. **戰術加固**：適應極端氣候、強震與全地形軍規環境",
            "2. **分散式韌性**：支援 Mesh 無線網絡，防止單點故障 (No Single Point of Failure)",
            "3. **雲-邊-端協同**：運行 Tier-3 輕量化模型 (Gemma 4 E4B)，進行邊緣目標識別"
        ],
        ["C2", "ISR", "API"]
    )

    # ----------------------------------------------------
    # Slide 29: Anti-Gravity AI CLI & Cron Schedule
    # ----------------------------------------------------
    build_2card_slide_fn(
        "Anti-Gravity AI CLI 自然語言日常檢索與 Cron 排程", "SECTION 5: SOVEREIGN AI & EDGE C2",
        "💬", "1. 自然語言即 CLI 指令", "Natural Language AI CLI Workflow",
        [
            "1. **對話即存取**：不需要命令視窗，直接說「幫我搜尋 AIEC 筆記中有關 RAG 權限的規定」",
            "2. **mcpvault API**：背景自動調用 `search_notes` 與 `write_note` 精確讀寫 Vault",
            "3. **隨手紀錄與整理**：「將剛才討論的 15 項量化指標精隨寫入今天的每日筆記」"
        ],
        "⏰", "2. 每週自動知識重整 Cron", "Weekly Knowledge Review Schedule",
        [
            "1. **排程時間**：每週日 09:17 (Cron：`17 9 * * 0`) 背景自動執行",
            "2. **4 大步驟**：盤點過去 7 天變動 -> 提煉 Clippings 至知識庫 -> 健康檢查 -> 更新 index.md",
            "3. **雙腦自我成長**：自動偵測孤立筆記與知識缺口，產出知識週報"
        ],
        ["AI", "AIEC", "RAG", "API", "Cron", "MOC"]
    )

    # ----------------------------------------------------
    # Slide 30: Conclusion & Future Outlook
    # ----------------------------------------------------
    build_2card_slide_fn(
        "結語與未來展望 — 打造 Responsible & Trustworthy AI", "SECTION 5: SOVEREIGN AI & EDGE C2",
        "🎯", "1. 專案核心貢獻與價值", "Project Impact & Value Summary",
        [
            "1. **權威標準整合**：確立 DoD CDAO, NIST AI RMF, ISO 42001 與 MITRE ATLAS 縱深防禦",
            "2. **量化評測 SOP**：建立 15 項量化指標、測試 SOP 與明確合格門檻",
            "3. **雙腦知識庫落地**：全套 21 篇雙向鏈結 Markdown 筆記完全存入 Obsidian Vault"
        ],
        "🚀", "2. 未來演進與推廣方向", "Future Roadmap & Extension",
        [
            "1. **MAITE 測試工廠骨架**：持續整合 garak, RAGAS, AgentBench 於單一 MLOps 流水線",
            "2. **LVC 平行戰場紅隊**：於 VBS 4 / EADSIM 中強化對抗貼片與電戰干擾紅隊演練",
            "3. **ISO 42001 / CMMC 認證**：備妥全套 SOP 與 Log 稽核軌跡，迎接正式合規檢定"
        ],
        ["DoD", "CDAO", "NIST", "RMF", "ISO", "AIMS", "MITRE", "ATLAS", "SOP", "MLOps", "LVC", "VBS", "EADSIM", "CMMC"]
    )

    out_dir = r'c:\Users\administartor\Downloads\AIEC'
    out_path = os.path.join(out_dir, 'AIEC_Full_30_Slides_With_Acronym_Footnotes.pptx')
    prs.save(out_path)
    print(f'Successfully generated 30-slide presentation with acronym footnotes at: {out_path}')

if __name__ == '__main__':
    build_30_deck_with_footnotes()
