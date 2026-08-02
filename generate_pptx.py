import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette: Premium Tech Dark Theme
    COLOR_BG = RGBColor(15, 23, 42)       # Slate 900
    COLOR_CARD = RGBColor(30, 41, 59)     # Slate 800
    COLOR_PRIMARY = RGBColor(56, 189, 248) # Cyan 400
    COLOR_ACCENT = RGBColor(251, 146, 60) # Amber/Orange 400
    COLOR_TEXT = RGBColor(241, 245, 249)  # Slate 100
    COLOR_MUTED = RGBColor(148, 163, 184) # Slate 400
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_BORDER = RGBColor(51, 65, 85)   # Slate 700

    def add_background(slide):
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="Gartner & Technical Intelligence"):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_ACCENT
        p_cat.font.name = "Microsoft JhengHei"

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY
        p_title.font.name = "Microsoft JhengHei"

    def add_card(slide, left, top, width, height, title, points, bg_color=COLOR_CARD, border_color=COLOR_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1)

        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        tf.margin_top = Inches(0.3)
        tf.margin_bottom = Inches(0.3)

        p_head = tf.paragraphs[0]
        p_head.text = title
        p_head.font.size = Pt(18)
        p_head.font.bold = True
        p_head.font.color.rgb = COLOR_TEXT
        p_head.font.name = "Microsoft JhengHei"
        p_head.space_after = Pt(14)

        for pt in points:
            p = tf.add_paragraph()
            p.font.size = Pt(13)
            p.font.color.rgb = COLOR_MUTED
            p.font.name = "Microsoft JhengHei"
            p.space_after = Pt(8)
            
            if isinstance(pt, tuple):
                run_bold = p.add_run()
                run_bold.text = f"• {pt[0]}： "
                run_bold.font.bold = True
                run_bold.font.color.rgb = COLOR_WHITE
                
                run_desc = p.add_run()
                run_desc.text = pt[1]
            else:
                p.text = f"• {pt}"

    # ----------------------------------------------------
    # Slide 1: Title Slide
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide1)

    t_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.5))
    tf1 = t_box.text_frame
    tf1.word_wrap = True

    p0 = tf1.paragraphs[0]
    p0.text = "Claude Mythos 5 技術解析與 Gartner 專家觀點"
    p0.font.size = Pt(34)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_PRIMARY
    p0.font.name = "Microsoft JhengHei"
    p0.space_after = Pt(16)

    p1 = tf1.add_paragraph()
    p1.text = "從 Mythos 5 核心特色、Fable 5 深度區別到 Paul Furtado 國防防禦與資安研討會深度提問"
    p1.font.size = Pt(18)
    p1.font.color.rgb = COLOR_MUTED
    p1.font.name = "Microsoft JhengHei"
    p1.space_after = Pt(40)

    p2 = tf1.add_paragraph()
    p2.text = "Gartner 技術研討會特輯 | 2026 年 7 月"
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_ACCENT
    p2.font.name = "Microsoft JhengHei"

    # ----------------------------------------------------
    # Slide 2: Agenda (Updated for 11 slides)
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide2)
    add_header(slide2, "簡報議程大綱 Agenda")

    agenda_items = [
        ("01", "Claude Mythos 5 核心技術特色", "頂級 Mythos 級定位、超長 1M Context Window 與高階科學推理能力"),
        ("02", "雙重用途與 Project Glasswing", "解除安全過濾器、特許存取機制與 30 天強制資料保留規範"),
        ("03", "Claude Mythos 5 vs. Fable 5 區別", "底層同源與部署異構、安全分類器差異與自動降級機制"),
        ("04", "Mythos 5 vs. Fable 5 綜合對比矩陣", "6 大維度深度比較（過濾機制、觸發行為、開放對象與場景）"),
        ("05", "Paul Furtado 觀點 1 & 2", "資安臨界點（Tipping Point）與董事會溝通戰略（Drop the Geek Speak）"),
        ("06", "Paul Furtado 觀點 3 & 4", "AI 供應鏈與主權風險（Model-Agnostic）及 Agentic AI 權限治理"),
        ("07", "研討會現場 Q&A 聽會指南", "針對資安架構與自動化防禦的核心提問方向指南"),
        ("08", "國防領域進階技術提問 (Defense Q&A)", "C4ISR 系統自主漏洞修補與氣閘網合規衝突中英文深度提問"),
        ("09", "企業 AI 治理行動建議總結", "架構解耦、語言轉譯與動態邊界控管的 Exec Summary")
    ]

    card_w = Inches(3.6)
    card_h = Inches(1.5)
    for i, item in enumerate(agenda_items):
        col = i % 3
        row = i // 3
        left = Inches(0.8 + col * 4.0)
        top = Inches(1.7 + row * 1.7)

        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = COLOR_BORDER
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.15)
        
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"[{item[0]}] "
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = COLOR_ACCENT
        
        r2 = p.add_run()
        r2.text = item[1]
        r2.font.bold = True
        r2.font.size = Pt(14)
        r2.font.color.rgb = COLOR_TEXT
        
        p_sub = tf.add_paragraph()
        p_sub.text = item[2]
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = COLOR_MUTED

    # ----------------------------------------------------
    # Slide 3: Mythos 5 Key Features Part 1
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide3)
    add_header(slide3, "Claude Mythos 5 核心技術特色（一）", "Part 1: Anthropic Mythos-Class Capability")

    pts3_1 = [
        ("頂級模型定位", "屬於 Anthropic 最強的 'Mythos' 級別，階級超越傳統 Opus 系列，專為最高難度的任務設計。"),
        ("長時程代理能力", "具備極強的 Long-horizon Agentic 處理能力，能獨立完成跨多個步驟、跨領域的複雜任務。"),
        ("前沿科學推理", "具備前所未有的科學假說生成能力，能處理高階分子生物學、蛋白質設計與複雜的化學結構。")
    ]
    pts3_2 = [
        ("1M Context Window", "支援高達 1,000,000 (1M) Tokens 的超大輸入視窗，可直接載入整個大規模代碼庫或大型文獻。"),
        ("128k Output Limit", "單次請求支援高達 128,000 Tokens 的輸出長度，適合生成巨型軟體架構或完整的長篇研究。"),
        ("全模態支援", "原生支援文字與多圖文影像解析能力，滿足多維度科學與工程數據分析。")
    ]

    add_card(slide3, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), "智力等級與代理能力", pts3_1)
    add_card(slide3, Inches(6.8), Inches(1.7), Inches(5.6), Inches(5.0), "長脈絡與極致輸出規格", pts3_2)

    # ----------------------------------------------------
    # Slide 4: Mythos 5 Key Features Part 2
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide4)
    add_header(slide4, "Claude Mythos 5 核心技術特色（二）", "Part 1: Dual-Use & Security Access")

    pts4_1 = [
        ("解除安全分類器", "完整移除了針對高風險技術（如攻擊性網路安全、生物毒素研究）的限制過濾器，以徹底發揮科研能力。"),
        ("深度資安分析", "能夠深入作業系統底層進行進階漏洞挖掘、攻擊面分析與防禦代碼重構。"),
        ("極致軟體工程", "可處理數百萬行專案的重構、自動化微服務拆分與 Legacy 系統大規模遷移。")
    ]
    pts4_2 = [
        ("Project Glasswing", "僅開放給通過 Anthropic 與政府審核的特許夥伴，專為國家級與關鍵基礎設施防禦專案（Glasswing）提供服務。"),
        ("不對大眾開放", "因具備強大的潛在雙重用途（Dual-use）危險性，完全不開放給一般公眾或一般 API 訂閱者。"),
        ("30天強制資料保留", "所有 Mythos 級流量強制實施 30-day Data Retention，確保過程皆可進行安全稽核與記錄。")
    ]

    add_card(slide4, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), "雙重用途與極限工程", pts4_1)
    add_card(slide4, Inches(6.8), Inches(1.7), Inches(5.6), Inches(5.0), "特許存取與安全合規", pts4_2)

    # ----------------------------------------------------
    # Slide 5: Mythos 5 vs Fable 5 Architectural Overview
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide5)
    add_header(slide5, "Claude Mythos 5 vs. Fable 5 深度對比（一）", "Part 2: Model Architecture & Safety Alignment")

    pts5_1 = [
        ("底層架構完全同源", "Mythos 5 與 Fable 5 基於完全相同的 Mythos 級基礎模型（Base Model）與訓練權重。"),
        ("本質為部署配置差異", "兩者的核心智慧與邏輯推理能力相同，主要差異在於外圍的安全過濾（Guardrails）與發行存取政策。"),
        ("性能上限一致", "在未觸發安全限制的標準計算任務中，兩者表現出同等水準的頂級思考品質。")
    ]
    pts5_2 = [
        ("Mythos 5 (無過濾)", "移除內建安全分類器，允許處理高風險漏洞探索與雙重用途問題；存取需經過 Glasswing 特許審核。"),
        ("Fable 5 (內建過濾)", "內建嚴格的安全過濾機制，若提示涉及攻擊性資安或危險生化會直接阻斷，開放給大眾與企業訂閱者。"),
        ("Fable 5 自動降級", "當 Fable 5 觸發安全保護時，會自動拒絕回答或自動降級（Fallback）至 Opus 4.8 進行安全應答。")
    ]

    add_card(slide5, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), "底層同源 (Shared Base Model)", pts5_1)
    add_card(slide5, Inches(6.8), Inches(1.7), Inches(5.6), Inches(5.0), "安全機制與流轉行為", pts5_2)

    # ----------------------------------------------------
    # Slide 6: Comparison Matrix (Table Slide)
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide6)
    add_header(slide6, "Claude Mythos 5 vs. Fable 5 綜合對比矩陣", "Part 2: Comprehensive Comparison Matrix")

    table_shape = slide6.shapes.add_table(7, 3, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.0))
    table = table_shape.table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(4.766)
    table.columns[2].width = Inches(4.766)

    matrix_data = [
        ["比較維度", "Claude Mythos 5", "Claude Fable 5"],
        ["底層模型架構", "Mythos 級旗艦基礎模型 (完全同源)", "Mythos 級旗艦基礎模型 (完全同源)"],
        ["安全過濾器 (Classifiers)", "解禁 / 移除安全過濾分類器", "內建嚴格安全與防禦過濾機制"],
        ["敏感請求行為", "不阻斷，直接輸出深度雙重用途解析", "阻斷拒答或自動降級 (Fallback to Opus 4.8)"],
        ["開放與存取對象", "特許審核 (Project Glasswing / 政府 / 特許企業)", "大眾開放 (一般企業、開發者與 API 用戶)"],
        ["主要應用場景", "關鍵基礎設施防禦、漏洞挖掘、高階科研", "企業級自主 Agent、大規模工程、商業分析"],
        ["資料監控規範", "30 天強制資料保留 (Mandatory Retention)", "30 天強制資料保留 (Mandatory Retention)"]
    ]

    for r_idx, row in enumerate(matrix_data):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c_idx == 0 else PP_ALIGN.LEFT
            p.font.name = "Microsoft JhengHei"
            
            if r_idx == 0:
                cell.fill.fore_color.rgb = COLOR_PRIMARY
                p.font.bold = True
                p.font.size = Pt(14)
                p.font.color.rgb = COLOR_BG
            else:
                cell.fill.fore_color.rgb = COLOR_CARD if r_idx % 2 == 1 else RGBColor(23, 32, 47)
                p.font.size = Pt(12)
                p.font.color.rgb = COLOR_TEXT if c_idx > 0 else COLOR_ACCENT
                if c_idx == 0:
                    p.font.bold = True

    # ----------------------------------------------------
    # Slide 7: Paul Furtado Points 1 & 2
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide7)
    add_header(slide7, "Paul Furtado 論點（一）：威脅臨界點與董事會溝通", "Part 3: Gartner Keynote Insights (Part 1)")

    pts7_1 = [
        ("漏洞自動化探索", "Mythos 5 這類模型能以秒級速度分析代碼並找出系統隱藏漏洞。"),
        ("傳統 Patch Cycle 失效", "攻擊與漏洞武器化的速度已大幅超越企業傳統『發現 ➔ 測試 ➔ 修補』的作業週期。"),
        ("攻防動態重塑", "防守方必須全面導入自動化響應，否則將面臨維度上的代差劣勢。")
    ]
    pts7_2 = [
        ("Drop the Geek Speak", "CISO 向董事會報告時，必須摒棄 CVE 編號與技術術語，改用商業語言。"),
        ("專注業務影響", "董事會只在乎『營運中斷時間、財務衝擊與品牌信任損害』。"),
        ("重塑風險承受度", "AI Agent 的引進要求企業高層重新校準 Risk Appetite（風險承受度）。")
    ]

    add_card(slide7, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), "論點 1: 資安威脅的臨界點 (Tipping Point)", pts7_1)
    add_card(slide7, Inches(6.8), Inches(1.7), Inches(5.6), Inches(5.0), "論點 2: 董事會溝通與風險重構", pts7_2)

    # ----------------------------------------------------
    # Slide 8: Paul Furtado Points 3 & 4
    # ----------------------------------------------------
    slide8 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide8)
    add_header(slide8, "Paul Furtado 論點（二）：供應鏈韌性與 Agent 治理", "Part 3: Gartner Keynote Insights (Part 2)")

    pts8_1 = [
        ("地緣政治與封鎖風險", "從 Mythos 5 短暫遭遇出口管制封鎖事件可知，過度依賴單一模型極具危險。"),
        ("Model-Agnostic 架構", "企業應建立『模型中立』架構，使系統能在不同模型間無縫切換，保障營運連續性。"),
        ("AI 主權 (AI Sovereignty)", "確保關鍵核心業務具備地緣備援與多供應鏈備援方案。")
    ]
    pts8_2 = [
        ("從問答到自主執行", "AI 已演進為自主代理（Agentic AI），安全重點轉移至邊界與授權控管。"),
        ("最小權限原則", "必須為 AI Agent 設定嚴格的 Sandbox、權限上限與即時介入機制（Human-in-the-loop）。"),
        ("合規與審計追蹤", "滿足 mandatory 30-day data retention 規範，建立完善的 AI 行為日誌。")
    ]

    add_card(slide8, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), "論點 3: AI 供應鏈與模型中立架構", pts8_1)
    add_card(slide8, Inches(6.8), Inches(1.7), Inches(5.6), Inches(5.0), "論點 4: 自主代理 (Agentic AI) 安全治理", pts8_2)

    # ----------------------------------------------------
    # Slide 9: Q&A & Conference Guide
    # ----------------------------------------------------
    slide9 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide9)
    add_header(slide9, "研討會現場 Q&A 聽會指南", "Part 4: Technical Seminar Q&A Strategy")

    pts9_1 = [
        ("提問切入點 1", "『對於資源有限的中型企業，如何以最低成本構建 Model-Agnostic (模型中立) 的備援架構？』"),
        ("提問切入點 2", "『在既有系統與多雲環境中，如何平衡切換不同 LLM 產生的 API 成本與延遲代價？』"),
        ("關鍵學習點", "聽取 Gartner 針對多模型切換路由（Model Routing Gateway）的建議標準。")
    ]
    pts9_2 = [
        ("提問切入點 1", "『面對 Mythos 5 級別的自動化漏洞發現能力，企業最應優先投資哪些自動化防禦工具？』"),
        ("提問切入點 2", "『在推動董事會同意提高資安預算時，如何量化 AI 武器化帶來的潛在財務風險？』"),
        ("關鍵學習點", "獲取 Paul Furtado 提供的 CISO 董事會溝通模板與指標。")
    ]

    add_card(slide9, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), "方向一：AI 供應鏈與備援韌性", pts9_1)
    add_card(slide9, Inches(6.8), Inches(1.7), Inches(5.6), Inches(5.0), "方向二：自動化防禦與高層溝通", pts9_2)

    # ----------------------------------------------------
    # Slide 10: Defense Technical Question (Bilingual) NEW!
    # ----------------------------------------------------
    slide10 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide10)
    add_header(slide10, "國防領域進階技術提問 (中英文對照)", "Part 5: Defense & National Security Deep-Dive Q&A")

    # Left Box: Traditional Chinese
    card_zh = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    card_zh.fill.solid()
    card_zh.fill.fore_color.rgb = COLOR_CARD
    card_zh.line.color.rgb = COLOR_ACCENT
    card_zh.line.width = Pt(1.5)

    tf_zh = card_zh.text_frame
    tf_zh.word_wrap = True
    tf_zh.margin_left = Inches(0.25)
    tf_zh.margin_right = Inches(0.25)
    tf_zh.margin_top = Inches(0.2)

    p_zh_t = tf_zh.paragraphs[0]
    p_zh_t.text = "🇹🇼 繁體中文問題 (C4ISR 與氣閘網合規)"
    p_zh_t.font.size = Pt(15)
    p_zh_t.font.bold = True
    p_zh_t.font.color.rgb = COLOR_ACCENT
    p_zh_t.font.name = "Microsoft JhengHei"
    p_zh_t.space_after = Pt(8)

    zh_text_blocks = [
        ("背景", "Mythos 5 解除安全過濾器、具 1M Token 視窗與 128k 輸出，已被納入 Project Glasswing 防禦專案。"),
        ("核心問題", "在 SIPRNet/JWICS 等高機密獨立氣閘網絡進行 C4ISR 系統自主零日漏洞發現與動態修補時，如何處理解析下列衝突："),
        ("1. 代理幻覺防範", "在無防護過濾器下進行長脈絡推論時，如何防止因為 Agent 過度授權引發非預期系統毀損？"),
        ("2. 合規條款衝突", "Mandatory 30-Day Data Retention 與國防絕對零外洩氣閘網政策該如何達成合規平衡？"),
        ("3. 備援模型稀缺", "在 Mythos 5 超強反編譯能力下，市場上是否有同等級模型可實現 Model-Agnostic 架構？")
    ]

    for b_title, b_desc in zh_text_blocks:
        p = tf_zh.add_paragraph()
        p.space_after = Pt(6)
        r_b = p.add_run()
        r_b.text = f"• {b_title}： "
        r_b.font.bold = True
        r_b.font.size = Pt(11)
        r_b.font.color.rgb = COLOR_WHITE
        r_b.font.name = "Microsoft JhengHei"
        
        r_d = p.add_run()
        r_d.text = b_desc
        r_d.font.size = Pt(10.5)
        r_d.font.color.rgb = COLOR_MUTED
        r_d.font.name = "Microsoft JhengHei"

    # Right Box: English
    card_en = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.2))
    card_en.fill.solid()
    card_en.fill.fore_color.rgb = COLOR_CARD
    card_en.line.color.rgb = COLOR_PRIMARY
    card_en.line.width = Pt(1.5)

    tf_en = card_en.text_frame
    tf_en.word_wrap = True
    tf_en.margin_left = Inches(0.25)
    tf_en.margin_right = Inches(0.25)
    tf_en.margin_top = Inches(0.2)

    p_en_t = tf_en.paragraphs[0]
    p_en_t.text = "🇺🇸 English Question (Air-Gapped & C4ISR)"
    p_en_t.font.size = Pt(15)
    p_en_t.font.bold = True
    p_en_t.font.color.rgb = COLOR_PRIMARY
    p_en_t.font.name = "Arial"
    p_en_t.space_after = Pt(8)

    en_text_blocks = [
        ("Context", "Mythos 5 features lifted safety classifiers, 1M context, and 128k max output within Project Glasswing."),
        ("Key Dilemma", "When deployed in air-gapped networks (SIPRNet/JWICS) for C4ISR zero-day patching, how do we resolve:"),
        ("1. Agentic Hallucination", "How to prevent unintended code degradation from agentic over-privilege without classifier guardrails?"),
        ("2. Data Retention", "How to reconcile mandatory 30-day monitoring retention with zero-exfiltration air-gapped policies?"),
        ("3. Fallback Feasibility", "Are there viable secondary fallback models matching Mythos 5's depth for Model-Agnostic architecture?")
    ]

    for b_title, b_desc in en_text_blocks:
        p = tf_en.add_paragraph()
        p.space_after = Pt(6)
        r_b = p.add_run()
        r_b.text = f"• {b_title}: "
        r_b.font.bold = True
        r_b.font.size = Pt(11)
        r_b.font.color.rgb = COLOR_WHITE
        r_b.font.name = "Arial"
        
        r_d = p.add_run()
        r_d.text = b_desc
        r_d.font.size = Pt(10.5)
        r_d.font.color.rgb = COLOR_MUTED
        r_d.font.name = "Arial"

    # ----------------------------------------------------
    # Slide 11: Executive Summary & Action Plan
    # ----------------------------------------------------
    slide11 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide11)
    add_header(slide11, "企業 AI 治理行動建議總結", "Executive Summary & Next Steps")

    pts11_cards = [
        ("1. 架構解耦 (Decouple)", "擴展企業與國防 AI 應用，解耦對單一模型依賴，建立模型中立（Model-Agnostic）控制與路由層。"),
        ("2. 語言轉譯 (Translate)", "重塑 CISO 對高層報告結構，將資安臨界點轉譯為業務中斷時間與量化風險，放棄純技術語詞。"),
        ("3. 邊界控管 (Govern)", "針對 Agentic AI 實施動態邊界與最小權限沙盒，同時滿足 30 天數據審計與國防級合規標準。")
    ]

    for idx, (title, desc) in enumerate(pts11_cards):
        top = Inches(1.7 + idx * 1.7)
        card = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.733), Inches(1.4))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = COLOR_PRIMARY if idx == 0 else (COLOR_ACCENT if idx == 1 else COLOR_BORDER)
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.4)
        tf.margin_top = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY if idx == 0 else (COLOR_ACCENT if idx == 1 else COLOR_TEXT)
        p.font.name = "Microsoft JhengHei"
        p.space_after = Pt(6)

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = COLOR_MUTED
        p_desc.font.name = "Microsoft JhengHei"

    prs.save(output_path)
    print(f"Presentation successfully updated with 11 slides at: {output_path}")

if __name__ == "__main__":
    out_dir = r"c:\Users\administartor\Downloads\AIEC"
    out_file = os.path.join(out_dir, "Claude_Mythos5_Fable5_Gartner_Analysis.pptx")
    create_presentation(out_file)
