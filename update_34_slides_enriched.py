# -*- coding: utf-8 -*-
import os, sys
import collections
import collections.abc

if not hasattr(collections, 'Iterable'):
    collections.Iterable = collections.abc.Iterable

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from pptx_helpers import set_pure_white_bg, add_header, add_icon_card, add_formatted_bullets, add_formatted_text

# Master Acronym Database
ACRONYM_DB = {
    "AIEC": ("人工智慧評測與認證體系", "Artificial Intelligence Evaluation & Certification"),
    "AI": ("人工智慧", "Artificial Intelligence"),
    "NIST": ("美國國家標準暨技術研究院", "National Institute of Standards and Technology"),
    "RMF": ("風險管理框架", "Risk Management Framework"),
    "DoD": ("美國國防部", "Department of Defense"),
    "CDAO": ("數位與人工智慧長辦公室", "Chief Digital and Artificial Intelligence Office"),
    "T&E": ("測試與評估", "Test and Evaluation"),
    "DT&E": ("發展性測試與評估", "Developmental Test and Evaluation"),
    "OT&E": ("作戰性測試與評估", "Operational Test and Evaluation"),
    "ISO": ("國際標準化組織", "International Organization for Standardization"),
    "AIMS": ("人工智慧管理系統", "Artificial Intelligence Management System"),
    "MITRE": ("邁特公司", "MITRE Corporation"),
    "ATLAS": ("人工智慧系統對抗威脅圖譜", "Adversarial Threat Landscape for Artificial-Intelligence Systems"),
    "ATT&CK": ("對抗戰術、技術與通用知識", "Adversarial Tactics, Techniques, and Common Knowledge"),
    "DAGR": ("國防人工智慧治理與風險指南", "Defense Artificial Intelligence Governance & Risk Guidelines"),
    "SHIELD": ("設定、精煉、改進、評估、記錄、偵測六大治理階段", "Set, Hone, Improve, Evaluate, Log, Detect"),
    "SOC": ("關切事項聲明", "Statement of Concern"),
    "RAI": ("負責任的人工智慧", "Responsible Artificial Intelligence"),
    "CMMC": ("網路安全成熟度模型認證", "Cybersecurity Maturity Model Certification"),
    "GDPR": ("通用資料保護規則", "General Data Protection Regulation"),
    "ML": ("機器學習", "Machine Learning"),
    "LLM": ("大語言模型", "Large Language Model"),
    "CV": ("電腦視覺", "Computer Vision"),
    "RAG": ("檢索增強生成", "Retrieval-Augmented Generation"),
    "HMT": ("人機協同戰術團隊", "Human-Autonomy Teaming"),
    "HSI": ("人機系統整合", "Human-Systems Integration"),
    "TEVV": ("測試、評估、驗證與確認", "Test, Evaluation, Verification, and Validation"),
    "JATIC": ("聯合人工智慧測試中心", "Joint AI Test Center"),
    "mAP": ("平均精確度均值", "mean Average Precision"),
    "Acc_adv": ("對抗測試樣本準確率", "Adversarial Sample Accuracy"),
    "Acc_clean": ("乾淨測試樣本準確率", "Clean Sample Accuracy"),
    "ECE": ("期望校準誤差", "Expected Calibration Error"),
    "garak": ("生成式 AI 紅隊分析工具包", "Generative AI Redteam Analysis Kit"),
    "OOD": ("分布外數據", "Out-of-Distribution"),
    "NRTK": ("自然穩健性測試工具包", "Natural Robustness Toolkit"),
    "XAITK": ("可解釋人工智慧工具包", "Explainable AI Toolkit"),
    "ART": ("對抗韌性測試工具包", "Adversarial Robustness Toolbox"),
    "HEART": ("高爆對抗紅隊測試工具", "High-Explosive Adversarial Red Teaming"),
    "FGSM": ("快速梯度符號攻擊法", "Fast Gradient Sign Method"),
    "PGD": ("投影梯度下降法", "Projected Gradient Descent"),
    "MSR": ("任務完成率", "Mission Success Rate"),
    "LVC": ("實戰-虛擬-構造平行戰場整合", "Live-Virtual-Constructive"),
    "VBS": ("虛擬戰場模擬軟體", "Virtual Battlespace"),
    "EADSIM": ("延伸防空模擬系統", "Extended Air Defense Simulation"),
    "DoDD": ("國防部指令", "Department of Defense Directive"),
    "ToAST": ("自主系統測試工具", "Testing of Autonomous Systems Tool"),
    "IDA": ("國防分析研究所", "Institute for Defense Analyses"),
    "MIT-LL": ("麻省理工學院林肯實驗室", "Massachusetts Institute of Technology Lincoln Laboratory"),
    "EEG": ("腦電圖儀", "Electroencephalography"),
    "NASA-TLX": ("美國航太總署任務負荷指數", "NASA Task Load Index"),
    "SHAP": ("沙普利附加解釋值", "SHapley Additive exPlanations"),
    "LIME": ("局部可解釋模型無關說明", "Local Interpretable Model-agnostic Explanations"),
    "OWASP": ("開放 Web 應用程式安全計畫", "Open Web Application Security Project"),
    "RAGAS": ("檢索增強生成評估指標", "Retrieval Augmented Generation Assessment"),
    "API": ("應用程式介面", "Application Programming Interface"),
    "OPA": ("開放策略代理", "Open Policy Agent"),
    "SPIFFE": ("通用安全生產身份框架", "Secure Production Identity Framework for Everyone"),
    "SPIRE": ("SPIFFE 執行期環境", "SPIFFE Runtime Environment"),
    "PyOD": ("Python 異常值偵測庫", "Python Outlier Detection"),
    "UQ": ("不確定性量化", "Uncertainty Quantification"),
    "MC-Dropout": ("蒙地卡羅 Dropout", "Monte Carlo Dropout"),
    "RBAC": ("基於角色的存取控制", "Role-Based Access Control"),
    "MoE": ("專家混合架構", "Mixture of Experts"),
    "C2": ("指揮管制", "Command and Control"),
    "ISR": ("情報、監視與偵察", "Intelligence, Surveillance, and Reconnaissance"),
    "GUI": ("圖形使用者介面", "Graphical User Interface"),
    "GGUF": ("通用模型量化格式", "Georgi Gerganov Unified Format"),
    "MOC": ("內容主索引地圖", "Map of Content"),
    "Cron": ("定時排程服務", "Chronos Command Scheduler"),
    "MLOps": ("機器學習營運", "Machine Learning Operations"),
    "NeMo": ("輝達神經網路模組化工具包", "Neural Modules"),
    "HITL": ("人在紐中/人在迴路", "Human-in-the-Loop"),
    "HOTL": ("人在紐上/實時監督", "Human-on-the-Loop"),
    "HOOTL": ("完全自主/人外迴路", "Human-out-of-the-Loop"),
    "AIF360": ("AI 公平性 360 工具包", "AI Fairness 360 Toolkit"),
    "MAITE": ("模型與 AI 測試評估基礎設施框架", "Model & AI Test and Evaluation Infrastructure Framework"),
    "NCSIST": ("國家中山科學研究院", "National Chung-Shan Institute of Science and Technology"),
    "GPS": ("全球定位系統", "Global Positioning System"),
    "Mesh": ("網狀拓撲通訊網絡", "Mesh Topology Network"),
    "TPM": ("可信賴平台模組", "Trusted Platform Module"),
    "RAM": ("隨機存取記憶體", "Random-Access Memory"),
    "Flash": ("快閃記憶體", "Flash Memory"),
    "LoRA": ("低秩適應微調", "Low-Rank Adaptation")
}

def render_latex_to_png(latex_str, output_filename, fontsize=18, color_hex='#1e3a8a'):
    s = '$' + latex_str.strip('$') + '$'
    os.makedirs('math_images', exist_ok=True)
    out_path = os.path.join('math_images', output_filename)
    fig, ax = plt.subplots(figsize=(6.5, 0.8))
    ax.axis('off')
    ax.text(0.5, 0.5, s, fontsize=fontsize, color=color_hex, ha='center', va='center')
    plt.savefig(out_path, bbox_inches='tight', dpi=300, transparent=True)
    plt.close(fig)
    return out_path

def add_acronym_footer(slide, acronym_keys, y_pos=7.95, height=0.9):
    NAVY = RGBColor(12, 35, 64)
    BLUE = RGBColor(37, 99, 235)
    MUTED = RGBColor(100, 116, 139)
    BG_FOOTER = RGBColor(241, 245, 249)
    BORDER_FOOTER = RGBColor(203, 213, 225)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y_pos), Inches(14.4), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = BG_FOOTER
    card.line.color.rgb = BORDER_FOOTER
    card.line.width = Pt(1.0)

    tf = card.text_frame
    tf.word_wrap = True

    p_head = tf.paragraphs[0]
    p_head.text = "📌 頁面英文縮寫與專有名詞對照 (Acronym & Term Footnotes):"
    p_head.font.size = Pt(10)
    p_head.font.bold = True
    p_head.font.color.rgb = BLUE
    p_head.font.name = "微軟正黑體"
    p_head.space_after = Pt(2)

    footer_parts = []
    for k in acronym_keys:
        if k in ACRONYM_DB:
            zh, en = ACRONYM_DB[k]
            footer_parts.append(f"**{k}**: {zh} ({en})")

    footer_str = "  |  ".join(footer_parts)

    p_body = tf.add_paragraph()
    add_formatted_text(p_body, footer_str, font_size=9, default_color=MUTED, bold_color=NAVY)

def generate_34_slides_enriched_deck():
    # Load original 30-slide presentation
    orig_path = r'c:\Users\administartor\Downloads\AIEC\AIEC_AI_Evaluation_30_Slides_NanoBanana.pptx'
    prs = Presentation(orig_path)
    blank_layout = prs.slide_layouts[6]

    NAVY = RGBColor(12, 35, 64)
    BLUE = RGBColor(37, 99, 235)
    DARK_BLUE = RGBColor(30, 58, 138)

    # ----------------------------------------------------
    # Update Slide 2: TABLE OF CONTENTS
    # ----------------------------------------------------
    slide2 = prs.slides[1]
    # Update Right Card in Slide 2 to include Section 6
    for shape in slide2.shapes:
        if shape.has_text_frame and "5. **第五區塊" in shape.text_frame.text:
            tf = shape.text_frame
            tf.clear()
            c2_bullets = [
                "4. **第四區塊：15 項量化評測指標與 SOP (P.17 - P.25)**",
                "   - Q1~Q15 完整量化計算公式 (300 DPI LaTeX 渲染)",
                "   - Pass/Fail 合格門檻與代表性測試工具鏈",
                "5. **第五區塊：評測平台與地端工具 (P.26 - P.30)**",
                "   - 主權 AI 平台、地端 LLM 引擎與戰術 C2 對接",
                "6. **第六區塊：國防 AI 安全進階防禦與主權架構 (P.31 - P.34)**",
                "   - NCSIST AIEC 總圖藍圖與 5 大底層戰術柱石",
                "   - 國防 AIEC 核心任務與交戰規則 (RoE) 人機授權邊界",
                "   - 戰術邊緣硬體防篡改與緊急模型自毀 (<100ms)",
                "   - 地端模型蒸餾 (Distillation) 與 Data/Model Provenance"
            ]
            add_formatted_bullets(tf, c2_bullets, font_size=12.5)

    # ----------------------------------------------------
    # SECTION 6: Slide 31 - NCSIST AIEC 國防 AI 評測總體藍圖與 5 大戰術柱石
    # ----------------------------------------------------
    s31 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s31)
    add_header(s31, "NCSIST AIEC 國防 AI 評測總體藍圖與 5 大戰術柱石", "SECTION 6: ADVANCED DEFENSE AI SECURITY & INFRASTRUCTURE")

    add_icon_card(s31, 0.8, 1.6, 7.0, 6.25, "🏛️", "NCSIST AIEC 架構圖與主動防禦", "NCSIST Defense Architecture Blueprint", accent_color=BLUE)
    tb_s31_l = s31.shapes.add_textbox(Inches(1.1), Inches(2.7), Inches(6.4), Inches(4.8))
    s31_l_bullets = [
        "1. **代碼注入與數據污染主動防禦**：針對視覺辨識 (YOLO, Mobile SAM) 與 LLM 模型，建立代碼注入 (Code Injection) 與數據污染 (Poisoning) 過濾閘門。",
        "2. **Cyber Range 對抗演練**：結合 [[MITRE ATLAS 人工智慧對抗威脅矩陣]] 與 CALDERA Arsenal 外掛，建立動態資安脆弱性模擬。",
        "3. **國家級機構中樞**：由國家中山科學研究院 (NCSIST) 主導主權算力與國防 AIEC 評測中心規範放行。"
    ]
    add_formatted_bullets(tb_s31_l.text_frame, s31_l_bullets, font_size=13.5)

    add_icon_card(s31, 8.2, 1.6, 7.0, 6.25, "🛡️", "底層 5 大戰術 SOP 柱石", "Five Tactical Operational Pillar SOPs", accent_color=DARK_BLUE)
    tb_s31_r = s31.shapes.add_textbox(Inches(8.5), Inches(2.7), Inches(6.4), Inches(4.8))
    s31_r_bullets = [
        "1. **擬定 ROE 交戰規則**：明確定義人機協同與可控性 (HITL/HOTL/HOOTL)。",
        "2. **驗證分層式 AI 架構**：確效系統能在極端作戰條件下完成任務。",
        "3. **執行紅隊對抗/演訓**：測試系統承受攻擊與受損後之彈性恢復力。",
        "4. **研發/軍事資料集分級**：落實清洗規則、版本控管與 Data/Model Provenance。",
        "5. **供應鏈安全檢查**：審查開源模型與第三方 SDK 後門及邏輯合規。"
    ]
    add_formatted_bullets(tb_s31_r.text_frame, s31_r_bullets, font_size=13.5)
    add_acronym_footer(s31, ["NCSIST", "AIEC", "ATLAS", "RoE", "HITL", "HOTL", "HOOTL", "CMMC", "TRL", "SOP"], y_pos=7.95, height=0.9)

    # ----------------------------------------------------
    # SECTION 6: Slide 32 - 國防 AIEC 核心任務與交戰規則 (RoE) 授權邊界
    # ----------------------------------------------------
    s32 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s32)
    add_header(s32, "國防 AIEC 核心任務與交戰規則 (RoE) 授權邊界", "SECTION 6: ADVANCED DEFENSE AI SECURITY & INFRASTRUCTURE")

    add_icon_card(s32, 0.8, 1.6, 7.0, 6.25, "🎯", "國防 AIEC 五大核心管制任務", "Five Core Control Tasks of Defense AIEC", accent_color=BLUE)
    tb_s32_l = s32.shapes.add_textbox(Inches(1.1), Inches(2.7), Inches(6.4), Inches(4.8))
    s32_l_bullets = [
        "1. **制定交戰規則 (RoE)**：劃分全系統人機自動化權能授權邊界。",
        "2. **審核主權基礎設施**：審核民雄院區三層算力與電力/物理韌性。",
        "3. **策劃 LVC 紅軍測試**：於 VBS 4 / EADSIM 平行戰場測試對抗防禦。",
        "4. **RAG 降密稽核**：審查研發資料集探勘，防止資訊降密洩漏 (Q14)。",
        "5. **供應鏈後門審查**：嚴格審查開源模型與 SDK 演算法後門。"
    ]
    add_formatted_bullets(tb_s32_l.text_frame, s32_l_bullets, font_size=13.5)

    add_icon_card(s32, 8.2, 1.6, 7.0, 6.25, "⚖️", "RoE 人機授權三階權能邊界", "Three-Tier Autonomy Boundary (RoE)", accent_color=DARK_BLUE)
    tb_s32_r = s32.shapes.add_textbox(Inches(8.5), Inches(2.7), Inches(6.4), Inches(4.8))
    s32_r_bullets = [
        "1. **Human-in-the-loop (HITL / 人在紐中)**：最終火力打擊與目標授權必須由人類指揮官手動執行。",
        "2. **Human-on-the-loop (HOTL / 人在紐上)**：AI 執行戰術建議，人類具備實時監控與強制中斷權 (Abort Button)。",
        "3. **Human-out-of-the-loop (HOOTL / 完全自主)**：僅限於蜂群無人機航路規劃與偵察等非致傷性任務。"
    ]
    add_formatted_bullets(tb_s32_r.text_frame, s32_r_bullets, font_size=13.5)
    add_acronym_footer(s32, ["RoE", "HITL", "HOTL", "HOOTL", "LVC", "VBS", "EADSIM", "RAG", "RBAC", "TRL"], y_pos=7.95, height=0.9)

    # ----------------------------------------------------
    # SECTION 6: Slide 33 - 戰術邊緣硬體安全與模型緊急自毀機制
    # ----------------------------------------------------
    s33 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s33)
    add_header(s33, "戰術邊緣硬體安全與模型緊急自毀機制", "SECTION 6: ADVANCED DEFENSE AI SECURITY & INFRASTRUCTURE")

    add_icon_card(s33, 0.8, 1.6, 7.0, 6.25, "🛡️", "無 GPS 網狀通訊與硬體級防篡改", "Hardware Tamper-Resistance in Mesh Networks", accent_color=BLUE)
    tb_s33_l = s33.shapes.add_textbox(Inches(1.1), Inches(2.7), Inches(6.4), Inches(4.8))
    s33_l_bullets = [
        "1. **網狀通訊被俘獲風險 (Mesh Network)**：在前線電戰干擾與無 GPS 網狀環境下，無人機與邊緣運算設備極易遭敵方俘獲。",
        "2. **硬體級防篡改 (Tamper-Resistance)**：採用軍規 Secure Enclave / TPM 晶片與物理開箱感測電路 (Tamper Mesh)。",
        "3. **物理防護障壁**：設備遭受異常解體或外殼穿透時，立即切斷主板電源並解發邏輯保護。"
    ]
    add_formatted_bullets(tb_s33_l.text_frame, s33_l_bullets, font_size=13.5)

    add_icon_card(s33, 8.2, 1.6, 7.0, 6.25, "💣", "模型緊急自毀與物理零化 (<100ms)", "Emergency Self-Destruct & Zeroization (<100ms)", accent_color=DARK_BLUE)
    tb_s33_r = s33.shapes.add_textbox(Inches(8.5), Inches(2.7), Inches(6.4), Inches(4.8))
    s33_r_bullets = [
        "1. **緊急權重複寫指令**：當系統偵測到物理入侵、通訊斷連超過安全時限或接獲中斷信號時，在 $<100\mathrm{ms}$ 內觸發自毀。",
        "2. **Flash/RAM 雜訊衝刷**：以隨機雜訊迅速覆寫模型記憶體與 Flash 快閃記憶體，徹底撕毀演算法與權重。",
        "3. **物理零化 (Zeroization)**：隨機化金鑰熔絲 (Key Fusing)，防止敵方逆向工程與模型竊取。"
    ]
    add_formatted_bullets(tb_s33_r.text_frame, s33_r_bullets, font_size=13.5)
    add_acronym_footer(s33, ["GPS", "Mesh", "TPM", "RAM", "Flash"], y_pos=7.95, height=0.9)

    # ----------------------------------------------------
    # SECTION 6: Slide 34 - 地端模型蒸餾、資料與模型溯源 SOP
    # ----------------------------------------------------
    s34 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s34)
    add_header(s34, "地端模型蒸餾、資料與模型溯源 SOP", "SECTION 6: ADVANCED DEFENSE AI SECURITY & INFRASTRUCTURE")

    add_icon_card(s34, 0.8, 1.6, 7.0, 6.25, "⚙️", "地端模型蒸餾與輕量化微調", "On-Premise Distillation & LoRA Fine-Tuning", accent_color=BLUE)
    tb_s34_l = s34.shapes.add_textbox(Inches(1.1), Inches(2.7), Inches(6.4), Inches(4.8))
    s34_l_bullets = [
        "1. **知識蒸餾 (Knowledge Distillation)**：將地端 70B 主權教師模型之推理能力蒸餾至 7B / 3B 輕量化邊緣學生模型。",
        "2. **Gemma 4 LoRA 本地微調**：針對戰術任務進行輕量化適應微調 (LoRA)，僅更新 $<1\%$ 參數，兼顧時延與精度。",
        "3. **GGUF 4-bit 量化**：透過 llama.cpp 實施 GGUF 量化部署，適應邊緣低功耗晶片。"
    ]
    add_formatted_bullets(tb_s34_l.text_frame, s34_l_bullets, font_size=13.5)

    add_icon_card(s34, 8.2, 1.6, 7.0, 6.25, "📊", "Data & Model Provenance 溯源與 Confidence", "Data & Model Provenance & Confidence Calibration", accent_color=DARK_BLUE)
    tb_s34_r = s34.shapes.add_textbox(Inches(8.5), Inches(2.7), Inches(6.4), Inches(4.8))
    s34_r_bullets = [
        "1. **Data Provenance (資料溯源)**：詳細記錄訓練資料集來源、邏輯檢查規則、清洗演算法版本號 (`v1.4.2`) 與標籤品質。",
        "2. **Model Provenance (模型溯源)**：追溯模型訓練超參數與 Checkpoint 雜湊值；若發現戰術誤判，可於 10 分鐘內回溯隔離受污染數據。",
        "3. **信心分數 (Confidence Score) 校準**：所有模型輸出伴隨 Confidence Score (score)，確保期望校準誤差 $\mathrm{ECE} \le 0.05$。"
    ]
    add_formatted_bullets(tb_s34_r.text_frame, s34_r_bullets, font_size=13.5)
    add_acronym_footer(s34, ["LLM", "LoRA", "GGUF", "ECE", "SOP"], y_pos=7.95, height=0.9)

    out_dir = r'c:\Users\administartor\Downloads\AIEC'
    out_path = os.path.join(out_dir, 'AIEC_AI_Evaluation_30_Slides_NanoBanana.pptx')
    prs.save(out_path)
    print(f'Successfully generated 34-slide presentation at: {out_path}')

if __name__ == '__main__':
    generate_34_slides_enriched_deck()
