# -*- coding: utf-8 -*-
import os, sys
import collections
import collections.abc

if not hasattr(collections, 'Iterable'):
    collections.Iterable = collections.abc.Iterable

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from pptx_helpers import set_pure_white_bg, add_header, add_icon_card, add_formatted_bullets, add_formatted_text

# Master Acronym Database
ACRONYM_DB = {
    "AIEC": ("人工智慧評測與認證體系", "Artificial Intelligence Evaluation & Certification"),
    "AI": ("人工智慧", "Artificial Intelligence"),
    "NIST": ("美國國家標準暨技術研究院", "National Institute of Standards and Technology"),
    "RMF": ("風險管理框架", "Risk Management Framework"),
    "DoD": ("美國國防部", "Department of Defense"),
    "CDAO": ("數位與人工智慧長辦公室", "Chief Digital and Artificial Intelligence Office"),
    "T&E": ("測試與評估", "Test and Evaluation"),
    "DT&E": ("發展性測試與評估", "Developmental Test and Evaluation"),
    "OT&E": ("作戰性測試與評估", "Operational Test and Evaluation"),
    "ISO": ("國際標準化組織", "International Organization for Standardization"),
    "AIMS": ("人工智慧管理系統", "Artificial Intelligence Management System"),
    "MITRE": ("邁特公司", "MITRE Corporation"),
    "ATLAS": ("人工智慧系統對抗威脅圖譜", "Adversarial Threat Landscape for Artificial-Intelligence Systems"),
    "ATT&CK": ("對抗戰術、技術與通用知識", "Adversarial Tactics, Techniques, and Common Knowledge"),
    "DAGR": ("國防人工智慧治理與風險指南", "Defense Artificial Intelligence Governance & Risk Guidelines"),
    "SHIELD": ("設定、精煉、改進、評估、記錄、偵測六大治理階段", "Set, Hone, Improve, Evaluate, Log, Detect"),
    "SOC": ("關切事項聲明", "Statement of Concern"),
    "RAI": ("負責任的人工智慧", "Responsible Artificial Intelligence"),
    "CMMC": ("網路安全成熟度模型認證", "Cybersecurity Maturity Model Certification"),
    "GDPR": ("通用資料保護規則", "General Data Protection Regulation"),
    "ML": ("機器學習", "Machine Learning"),
    "LLM": ("大語言模型", "Large Language Model"),
    "CV": ("電腦視覺", "Computer Vision"),
    "RAG": ("檢索增強生成", "Retrieval-Augmented Generation"),
    "HMT": ("人機協同戰術團隊", "Human-Autonomy Teaming"),
    "HSI": ("人機系統整合", "Human-Systems Integration"),
    "TEVV": ("測試、評估、驗證與確認", "Test, Evaluation, Verification, and Validation"),
    "JATIC": ("聯合人工智慧測試中心", "Joint AI Test Center"),
    "mAP": ("平均精確度均值", "mean Average Precision"),
    "Acc_adv": ("對抗測試樣本準確率", "Adversarial Sample Accuracy"),
    "Acc_clean": ("乾淨測試樣本準確率", "Clean Sample Accuracy"),
    "ECE": ("期望校準誤差", "Expected Calibration Error"),
    "garak": ("生成式 AI 紅隊分析工具包", "Generative AI Redteam Analysis Kit"),
    "OOD": ("分布外數據", "Out-of-Distribution"),
    "NRTK": ("自然穩健性測試工具包", "Natural Robustness Toolkit"),
    "XAITK": ("可解釋人工智慧工具包", "Explainable AI Toolkit"),
    "ART": ("對抗韌性測試工具包", "Adversarial Robustness Toolbox"),
    "HEART": ("高爆對抗紅隊測試工具", "High-Explosive Adversarial Red Teaming"),
    "FGSM": ("快速梯度符號攻擊法", "Fast Gradient Sign Method"),
    "PGD": ("投影梯度下降法", "Projected Gradient Descent"),
    "MSR": ("任務完成率", "Mission Success Rate"),
    "RAGAS": ("檢索增強生成評估指標", "Retrieval Augmented Generation Assessment"),
    "LVC": ("實戰-虛擬-構造平行戰場整合", "Live-Virtual-Constructive"),
    "VBS": ("虛擬戰場模擬軟體", "Virtual Battlespace"),
    "EADSIM": ("延伸防空模擬系統", "Extended Air Defense Simulation"),
    "DoDD": ("國防部指令", "Department of Defense Directive"),
    "ToAST": ("自主系統測試工具", "Testing of Autonomous Systems Tool"),
    "IDA": ("國防分析研究所", "Institute for Defense Analyses"),
    "MIT-LL": ("麻省理工學院林肯實驗室", "Massachusetts Institute of Technology Lincoln Laboratory"),
    "EEG": ("腦電圖儀", "Electroencephalography"),
    "NASA-TLX": ("美國航太總署任務負荷指數", "NASA Task Load Index"),
    "SHAP": ("沙普利附加解釋值", "SHapley Additive exPlanations"),
    "LIME": ("局部可解釋模型無關說明", "Local Interpretable Model-agnostic Explanations"),
    "OWASP": ("開放 Web 應用程式安全計畫", "Open Web Application Security Project"),
    "API": ("應用程式介面", "Application Programming Interface"),
    "OPA": ("開放策略代理", "Open Policy Agent"),
    "SPIFFE": ("通用安全生產身份框架", "Secure Production Identity Framework for Everyone"),
    "SPIRE": ("SPIFFE 執行期環境", "SPIFFE Runtime Environment"),
    "PyOD": ("Python 異常值偵測庫", "Python Outlier Detection"),
    "UQ": ("不確定性量化", "Uncertainty Quantification"),
    "MC-Dropout": ("蒙地卡羅 Dropout", "Monte Carlo Dropout"),
    "RBAC": ("基於角色的存取控制", "Role-Based Access Control"),
    "MoE": ("專家混合架構", "Mixture of Experts"),
    "C2": ("指揮管制", "Command and Control"),
    "ISR": ("情報、監視與偵察", "Intelligence, Surveillance, and Reconnaissance"),
    "GUI": ("圖形使用者介面", "Graphical User Interface"),
    "GGUF": ("通用模型量化格式", "Georgi Gerganov Unified Format"),
    "MOC": ("內容主索引地圖", "Map of Content"),
    "Cron": ("定時排程服務", "Chronos Command Scheduler"),
    "MLOps": ("機器學習營運", "Machine Learning Operations"),
    "NeMo": ("輝達神經網路模組化工具包", "Neural Modules"),
    "HITL": ("人在紐中/人在迴路", "Human-in-the-Loop"),
    "HOTL": ("人在紐上/實時監督", "Human-on-the-Loop"),
    "HOOTL": ("完全自主/人外迴路", "Human-out-of-the-Loop"),
    "AIF360": ("AI 公平性 360 工具包", "AI Fairness 360 Toolkit"),
    "MAITE": ("模型與 AI 測試評估基礎設施框架", "Model & AI Test and Evaluation Infrastructure Framework"),
    "NCSIST": ("國家中山科學研究院", "National Chung-Shan Institute of Science and Technology"),
    "GPS": ("全球定位系統", "Global Positioning System"),
    "Mesh": ("網狀拓撲通訊網絡", "Mesh Topology Network"),
    "TPM": ("可信賴平台模組", "Trusted Platform Module"),
    "RAM": ("隨機存取記憶體", "Random-Access Memory"),
    "Flash": ("快閃記憶體", "Flash Memory"),
    "LoRA": ("低秩適應微調", "Low-Rank Adaptation"),
    "YOLO": ("單次目標偵測演算法", "You Only Look Once"),
    "SAM": ("通用分割模型", "Segment Anything Model"),
    "ROE": ("交戰規則", "Rules of Engagement"),
    "RoE": ("交戰規則", "Rules of Engagement")
}

def render_latex_to_png(latex_str, output_filename, fontsize=18, color_hex='#1e3a8a'):
    s = '$' + latex_str.strip('$') + '$'
    os.makedirs('math_images', exist_ok=True)
    out_path = os.path.join('math_images', output_filename)
    fig, ax = plt.subplots(figsize=(6.5, 0.8))
    ax.axis('off')
    ax.text(0.5, 0.5, s, fontsize=fontsize, color=color_hex, ha='center', va='center')
    plt.savefig(out_path, bbox_inches='tight', dpi=300, transparent=True)
    plt.close(fig)
    return out_path

def add_acronym_footer(slide, acronym_keys, y_pos=7.75, height=0.95):
    NAVY = RGBColor(12, 35, 64)
    BLUE = RGBColor(37, 99, 235)
    MUTED = RGBColor(100, 116, 139)
    BG_FOOTER = RGBColor(241, 245, 249)
    BORDER_FOOTER = RGBColor(203, 213, 225)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y_pos), Inches(14.4), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = BG_FOOTER
    card.line.color.rgb = BORDER_FOOTER
    card.line.width = Pt(1.0)

    tf = card.text_frame
    tf.word_wrap = True

    p_head = tf.paragraphs[0]
    p_head.text = "📌 頁面英文縮寫與專有名詞對照 (Acronym & Term Footnotes):"
    p_head.font.size = Pt(10)
    p_head.font.bold = True
    p_head.font.color.rgb = BLUE
    p_head.font.name = "微軟正黑體"
    p_head.space_after = Pt(2)

    footer_parts = []
    for k in acronym_keys:
        if k in ACRONYM_DB:
            zh, en = ACRONYM_DB[k]
            footer_parts.append(f"**{k}**: {zh} ({en})")

    footer_str = "  |  ".join(footer_parts)

    p_body = tf.add_paragraph()
    add_formatted_text(p_body, footer_str, font_size=9, default_color=MUTED, bold_color=NAVY)

def generate_33_slides_master_deck():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank_layout = prs.slide_layouts[6]

    NAVY = RGBColor(12, 35, 64)
    BLUE = RGBColor(37, 99, 235)
    DARK_BLUE = RGBColor(30, 58, 138)

    # ----------------------------------------------------
    # Slide 1: TITLE SLIDE (Requirement 3 & 1)
    # ----------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s1)

    cover_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.6), Inches(14.4), Inches(7.8))
    cover_card.fill.solid()
    cover_card.fill.fore_color.rgb = NAVY
    cover_card.line.fill.background()

    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(1.1), Inches(7.5), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = BLUE
    badge.line.fill.background()
    p = badge.text_frame.paragraphs[0]
    p.text = "🛡️ AIEC 國防與企業級 AI 專屬評測與驗證 SOP 體系"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "微軟正黑體"

    tb1 = s1.shapes.add_textbox(Inches(1.3), Inches(1.9), Inches(13.4), Inches(2.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    # Requirement 3: Change title to "AIEC 國防級 AI 專屬評測與驗證 SOP 全集"
    p.text = "AIEC 國防級 AI 專屬評測與驗證 SOP 全集"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "微軟正黑體"

    tb_sub = s1.shapes.add_textbox(Inches(1.3), Inches(4.2), Inches(13.4), Inches(2.5))
    tf_sub = tb_sub.text_frame
    tf_sub.word_wrap = True

    cover_bullets = [
        "1. **國家級總體架構藍圖對齊 (NCSIST AIEC Blueprint)**：整合代碼注入/數據污染防禦、底層 5 大戰術柱石 (RoE, 分層架構, 紅隊演練, 資料分級, 供應鏈合規)",
        "2. **安全、保密與審計三維縱深矩陣 (3D Defense Matrix)**：對抗防禦 (ATLAS)、邊緣防篡改與模型自毀 (<100ms)、聯邦學習「模型移動，資料不動」、Data/Model Provenance 溯源",
        "3. **權威標準與量化指標全集**：精準融合 ISO 42001 (AIMS), NIST AI RMF 1.0 (Measure), DoD CDAO T&E, MITRE ATLAS 與 15 項 LaTeX 量化評測公式"
    ]
    add_formatted_bullets(tf_sub, cover_bullets, font_size=15.0, text_color=RGBColor(226, 232, 240), bold_color=RGBColor(255, 255, 255))

    # Requirement 1: Add MSR, garak, RAGAS to Slide 1 Footnote
    s1_acronyms = ["AIEC", "NCSIST", "NIST", "RMF", "DoD", "CDAO", "T&E", "ISO", "AIMS", "MITRE", "ATLAS", "MSR", "garak", "RAGAS", "SOP"]
    add_acronym_footer(s1, s1_acronyms, y_pos=7.65, height=1.0)

    # ----------------------------------------------------
    # Slide 2: TABLE OF CONTENTS (Requirement 4 & 2)
    # ----------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s2)
    # Requirement 4: Header title "AIEC 國防級 AI 評測全集 - 簡報目錄"
    add_header(s2, "AIEC 國防級 AI 評測全集 - 簡報目錄", "AIEC GOVERNANCE & EVALUATION MASTER DECK")

    # Requirement 2: Clean, non-overlapping 2-card layout
    add_icon_card(s2, 0.8, 1.6, 7.0, 5.9, "📋", "簡報導覽與核心區塊 (1-3)", "Governance, Methodology & System SOPs", accent_color=BLUE)
    tb_c1 = s2.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(6.4), Inches(4.7))
    c1_bullets = [
        "1. **第一區塊：評測框架與雙支柱治理 (P.3 - 6)**",
        "   - P.3 治理雙支柱 | P.4 SHIELD 6階段 | P.5 ISO 42001 | P.6 MITRE ATLAS",
        "2. **第二區塊：T&E 評測矩陣與方法論 (P.7 - 10)**",
        "   - P.7 T&E 四層次 | P.8 JATIC 7構面 | P.9 6大評測方法 | P.10 三維縱深矩陣",
        "3. **第三區塊：六大類 AI 應用系統評測 SOP (P.11 - 16)**",
        "   - P.11 CV目標偵測 | P.12 LLM生成式 | P.13 RAG知識庫 | P.14 Agent軌跡 | P.15 HMT與RoE | P.16 預測分析"
    ]
    add_formatted_bullets(tb_c1.text_frame, c1_bullets, font_size=12.0)

    add_icon_card(s2, 8.2, 1.6, 7.0, 5.9, "📐", "簡報導覽與核心區塊 (4-6)", "Quantitative Metrics & Platform Architecture", accent_color=DARK_BLUE)
    tb_c2 = s2.shapes.add_textbox(Inches(8.5), Inches(2.6), Inches(6.4), Inches(4.7))
    c2_bullets = [
        "4. **第四區塊：15 項量化評測指標與 SOP (P.17 - 25)**",
        "   - P.17 15項指標總覽 | P.18-25 Q1~Q15 數學公式、門檻與工具鏈",
        "5. **第五區塊：評測平台與地端工具 (P.26 - 29)**",
        "   - P.26 主權LLM架構 | P.27 地端推論 | P.28 C2與節點 | P.29 聯邦學習",
        "6. **第六區塊：國防 AI 安全進階防禦與主權架構 (P.30 - 33)**",
        "   - P.30 NCSIST AIEC 總藍圖與5大柱石 | P.31 國防 AIEC 任務與 RoE 邊界",
        "   - P.32 戰術邊緣硬體防篡改與模型自毀 (<100ms) | P.33 地端模型蒸餾與 Provenance"
    ]
    add_formatted_bullets(tb_c2.text_frame, c2_bullets, font_size=12.0)
    add_acronym_footer(s2, ["AIEC", "NCSIST", "SHIELD", "ISO", "AIMS", "ATLAS", "JATIC", "RAGAS", "garak", "PyOD"], y_pos=7.65, height=1.0)

    # ----------------------------------------------------
    # SECTION 1: Slide 3 - AIEC 規範與治理雙支柱
    # ----------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s3)
    add_header(s3, "AIEC 規範與治理雙支柱 (Governance Dual Pillars)", "SECTION 1: EVALUATION FRAMEWORK & GOVERNANCE")

    add_icon_card(s3, 0.8, 1.6, 7.0, 5.9, "🏛️", "支柱一：NCSIST AIEC 藍圖與 DAGR 風險指南", "NCSIST AIEC Blueprint & DAGR Guidelines", accent_color=BLUE)
    tb_s3_l = s3.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(6.4), Inches(4.7))
    s3_l_bullets = [
        "1. **代碼注入與數據污染主動防禦**：針對視覺辨識 (YOLO, Mobile SAM) 與 LLM 注入對抗過濾與 Cyber Range 演練。",
        "2. **底層 5 大戰術 SOP 柱石**：",
        "   - **擬定 RoE 交戰規則**：明訂人機協同與可控性 (HITL/HOTL/HOOTL)。",
        "   - **驗證分層式 AI 架構**：確效作戰條件下完成任務之能力。",
        "   - **執行紅隊對抗/演訓**：測試承受攻擊與受損彈性恢復。",
        "   - **研發/軍事資料集分級**：落實清洗規則與 Data/Model Provenance。",
        "   - **供應鏈安全檢查**：審查開源模型與第三方 SDK 後門。"
    ]
    add_formatted_bullets(tb_s3_l.text_frame, s3_l_bullets, font_size=12.5)

    add_icon_card(s3, 8.2, 1.6, 7.0, 5.9, "🔄", "支柱二：SHIELD 六大治理循環", "SHIELD 6-Stage Governance Lifecycle", accent_color=DARK_BLUE)
    tb_s3_r = s3.shapes.add_textbox(Inches(8.5), Inches(2.6), Inches(6.4), Inches(4.7))
    s3_r_bullets = [
        "1. **Set (目標界定)**：界定 A~F 類應用邊界、RoE 授權與 CMMC L2 等級。",
        "2. **Hone (精煉調校)**：資料清洗、Gemma 4 LoRA 輕量化微調與 GGUF 量化。",
        "3. **Improve (連續改進)**：根據紅軍演練漏洞持續進行模型重訓練。",
        "4. **Evaluate (量化評測)**：執行 [[AIEC 15 項量化指標]] 並核發 TRL 證書。",
        "5. **Log (高保真日誌)**：強制保存傳感器輸入、Confidence Score 與 COA 行動方案。",
        "6. **Detect (漂移偵測)**：即時監控數據分布與概念漂移 (Concept Drift)。"
    ]
    add_formatted_bullets(tb_s3_r.text_frame, s3_r_bullets, font_size=12.5)
    add_acronym_footer(s3, ["NCSIST", "AIEC", "DAGR", "SHIELD", "CMMC", "RoE", "HITL", "HOTL", "HOOTL", "TRL", "COA"], y_pos=7.65, height=1.0)

    # ----------------------------------------------------
    # SECTION 1: Slide 4 - SHIELD 治理循環活動
    # ----------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s4)
    add_header(s4, "SHIELD 治理循環活動 (SHIELD Governance Lifecycle)", "SECTION 1: EVALUATION FRAMEWORK & GOVERNANCE")

    add_icon_card(s4, 0.8, 1.6, 14.4, 5.9, "🔄", "SHIELD 六大治理階段與審計控制點", "Six Governance Stages & Audit Control Points", accent_color=BLUE)
    tb_s4 = s4.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(13.8), Inches(4.7))
    s4_bullets = [
        "1. **Set (目標界定)**：定義戰術邊界、CMMC L2 資安等級與 RoE 人機權能授權 (HITL/HOTL/HOOTL)。",
        "2. **Hone (精煉調校)**：執行資料清洗規則、安全對齊、Gemma 4 本地 LoRA 微調與 GGUF 4-bit 量化。",
        "3. **Improve (連續改進)**：根據 CALDERA Arsenal 與 garak 漏洞演練回饋，實施動態補強與護欄重訓練。",
        "4. **Evaluate (量化評測)**：經由獨立 AIEC 測試中心執行 Q1~Q15 量化指標，產出具備 TRL 放行資格之確效報告。",
        "5. **Log (高保真日誌)**：記錄原生的傳感器輸入數據、XAI 歸因熱力圖、Confidence Score 與最終 COA 決策路徑。",
        "6. **Detect (漂移監控)**：營運期間部署 PyOD 與 Alibi 模組，實現漂移告警極低時延。"
    ]
    add_formatted_bullets(tb_s4.text_frame, s4_bullets, font_size=13.0)
    add_acronym_footer(s4, ["SHIELD", "CMMC", "RoE", "HITL", "HOTL", "HOOTL", "TRL", "COA", "PyOD", "AIEC"], y_pos=7.65, height=1.0)

    # ----------------------------------------------------
    # SECTION 1: Slide 5 - ISO 42001 人工智慧管理系統
    # ----------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s5)
    add_header(s5, "ISO 42001 (AIMS) 與 AIEC 國防角色劃分", "SECTION 1: EVALUATION FRAMEWORK & GOVERNANCE")

    add_icon_card(s5, 0.8, 1.6, 7.0, 5.9, "🌐", "ISO 42001：通用框架與方法論", "Framework & Methodology (AIMS)", accent_color=BLUE)
    tb_s5_l = s5.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(6.4), Inches(4.7))
    s5_l_bullets = [
        "1. **國際通用標準**：提供組織建立 AI 風險評估與生命週期管理 SOP 的底層邏輯。",
        "2. **AI Impact Assessment (AIIA)**：要求實施 AI 影響評估，分析戰術部署對安全與人命之影響。",
        "3. **Clause 8.4 透明度要求**：要求關鍵決策模型必須具備 XAI 特徵熱力圖歸因能力 (Point Game)。",
        "4. **Annex A 控制項**：規範資料治理 (A.6)、對抗防禦 (A.7) 與存取控制 (A.8)。"
    ]
    add_formatted_bullets(tb_s5_l.text_frame, s5_l_bullets, font_size=13.0)

    add_icon_card(s5, 8.2, 1.6, 7.0, 5.9, "🎯", "AIEC：領域化審查閘門與規則制定者", "Domain-Specific Gatekeeper & TRL Enforcement", accent_color=DARK_BLUE)
    tb_s5_r = s5.shapes.add_textbox(Inches(8.5), Inches(2.6), Inches(6.4), Inches(4.7))
    s5_r_bullets = [
        "1. **標準戰術化**：將 ISO 通用標準詮釋為武器系統、C2 指管與情報分析中的具體 SOP。",
        "2. **TRL 放行閘門**：決定 AI 專案是否具備足夠的技術成熟度 (TRL) 可部署於作戰環境。",
        "3. **國防 AIEC 核心任務**：",
        "   - 制定 RoE 人機授權邊界 (HITL/HOTL/HOOTL)",
        "   - 審核民雄院區三層算力與主權基礎設施",
        "   - 執行 VBS 4 / EADSIM LVC 紅軍演練測試",
        "   - 研發資料集分級與 RAG 抗降密洩漏稽核"
    ]
    add_formatted_bullets(tb_s5_r.text_frame, s5_r_bullets, font_size=13.0)
    add_acronym_footer(s5, ["ISO", "AIMS", "AIIA", "AIEC", "TRL", "RoE", "HITL", "HOTL", "HOOTL", "LVC", "VBS", "EADSIM", "RAG"], y_pos=7.65, height=1.0)

    # ----------------------------------------------------
    # SECTION 1: Slide 6 - MITRE ATLAS 人工智慧對抗威脅矩陣
    # ----------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s6)
    add_header(s6, "MITRE ATLAS 人工智慧對抗威脅矩陣", "SECTION 1: EVALUATION FRAMEWORK & GOVERNANCE")

    add_icon_card(s6, 0.8, 1.6, 7.0, 5.9, "⚔️", "ATLAS 威脅矩陣與 16 大戰術鏈", "ATLAS Matrix & 16 Tactics Chain", accent_color=BLUE)
    tb_s6_l = s6.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(6.4), Inches(4.7))
    s6_l_bullets = [
        "1. **專屬對抗威脅圖譜**：編纂 16 大戰術 (Tactics)、170 項技術 (Techniques) 與 32 種緩解措施。",
        "2. **16 大戰術鏈條**：涵蓋 Reconnaissance, Data Poisoning, Prompt Injection, Model Stealing, Defense Evasion, Exfiltration 至 Impact。",
        "3. **代碼注入與數據污染防禦**：專門反制對抗貼片、訓練集後門與 LLM 指令突圍。",
        "4. **CALDERA Arsenal 外掛**：整合自動化 AI 紅隊演練外掛進行漏洞掃描。"
    ]
    add_formatted_bullets(tb_s6_l.text_frame, s6_l_bullets, font_size=12.5)

    add_icon_card(s6, 8.2, 1.6, 7.0, 5.9, "📊", "ATT&CK vs. ATLAS 比較與紅軍演練", "ATT&CK vs. ATLAS Comparison & Red Teaming", accent_color=DARK_BLUE)
    tb_s6_r = s6.shapes.add_textbox(Inches(8.5), Inches(2.6), Inches(6.4), Inches(4.7))
    s6_r_bullets = [
        "1. **ATT&CK vs. ATLAS 比較**：",
        "   - **ATT&CK**：關注傳統 IT 網路、伺服器與端點漏洞。",
        "   - **ATLAS**：關注 AI 模型、訓練數據管線、向量資料庫與推論 API。",
        "2. **OWASP Top 10 for LLM 對照**：OWASP 為風險分類學，ATLAS 為完整端到端對抗攻擊鏈。",
        "3. **AIEC 紅軍對抗實施**：使用 garak, IBM ART 360, PromptBench 在 Cyber Range 中進行實戰推演。"
    ]
    add_formatted_bullets(tb_s6_r.text_frame, s6_r_bullets, font_size=12.5)
    add_acronym_footer(s6, ["MITRE", "ATLAS", "ATT&CK", "OWASP", "LLM", "API", "garak", "ART", "AIEC"], y_pos=7.65, height=1.0)

    # ----------------------------------------------------
    # SECTION 2: Slides 7-10 (T&E Matrix & Methodologies)
    # ----------------------------------------------------
    # Slide 7
    s7 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s7)
    add_header(s7, "DoD CDAO T&E 四大能力層次", "SECTION 2: T&E MATRIX & METHODOLOGY")
    add_icon_card(s7, 0.8, 1.6, 14.4, 5.9, "📐", "DoD CDAO AI T&E 漸進能力評測階梯", "Progressive T&E Capability Hierarchy", accent_color=BLUE)
    tb_s7 = s7.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(13.8), Inches(4.7))
    s7_bullets = [
        "1. **Level 1: 基礎模型評測 (Base Model T&E)**：驗證演算法精確度、對抗韌性與自然穩健性。",
        "2. **Level 2: 人機系統整合 (HSI T&E)**：評估指揮官與操作員之認知負荷 (NASA-TLX)、眼動追蹤與信任校準 (ECE)。",
        "3. **Level 3: 系統整合評測 (Systems Integration T&E)**：驗證 API 閘道零信任 (SPIFFE/OPA)、工具調用邊界與通訊時延。",
        "4. **Level 4: 作戰性測試與評估 (Operational T&E)**：於 VBS 4 / EADSIM LVC 虛實整合平行戰場中，確效端到端任務完成率 (MSR) 與擊殺鏈閉合速度。"
    ]
    add_formatted_bullets(tb_s7.text_frame, s7_bullets, font_size=13.0)
    add_acronym_footer(s7, ["DoD", "CDAO", "T&E", "HSI", "NASA-TLX", "ECE", "SPIFFE", "OPA", "API", "VBS", "EADSIM", "LVC", "MSR"], y_pos=7.65, height=1.0)

    # Slide 8
    s8 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s8)
    add_header(s8, "JATIC 七大共通構面 (JATIC Common Dimensions)", "SECTION 2: T&E MATRIX & METHODOLOGY")
    add_icon_card(s8, 0.8, 1.6, 7.0, 5.9, "🏛️", "JATIC 七大構面 (1 - 4)", "Robustness, Resiliency, Explainability & Competence", accent_color=BLUE)
    tb_s8_l = s8.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(6.4), Inches(4.7))
    s8_l_bullets = [
        "1. ** Robustness (對抗與自然穩健性)**：對抗貼片與惡劣天候下維持高精度 (Q1, Q2)。",
        "2. ** Resiliency (系統韌性與失效安全)**：極端異常時 100ms 內安全中斷 (Q4)。",
        "3. ** Explainability (可解釋性)**：白箱提供特徵歸因熱力圖與 Point Game Score (Q7)。",
        "4. ** Competence (勝任力與 MSR)**：端到端任務完成率 $\mathrm{MSR} \ge 0.95$ (Q3)。"
    ]
    add_formatted_bullets(tb_s8_l.text_frame, s8_l_bullets, font_size=13.0)

    add_icon_card(s8, 8.2, 1.6, 7.0, 5.9, "⚖️", "JATIC 七大構面 (5 - 7)", "Fairness, Calibration & Drift Detection", accent_color=DARK_BLUE)
    tb_s8_r = s8.shapes.add_textbox(Inches(8.5), Inches(2.6), Inches(6.4), Inches(4.7))
    s8_r_bullets = [
        "5. ** Fairness (公平性與無偏見)**：訓練數據與模型輸出無隱性偏見，合規 AIF360。",
        "6. ** Trust Calibration (信任校準)**：期望校準誤差 $\mathrm{ECE} \le 0.05$，防止過度依賴 (Q5)。",
        "7. ** Drift Detection (漂移監控)**：即時捕捉戰場數據與概念漂移 (Q12, PyOD)。"
    ]
    add_formatted_bullets(tb_s8_r.text_frame, s8_r_bullets, font_size=13.0)
    add_acronym_footer(s8, ["JATIC", "MSR", "ECE", "PyOD", "AIF360"], y_pos=7.65, height=1.0)

    # Slide 9
    s9 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s9)
    add_header(s9, "國防 AI 評測 6 大方法論與 SOP", "SECTION 2: T&E MATRIX & METHODOLOGY")
    add_icon_card(s9, 0.8, 1.6, 14.4, 5.9, "🧪", "六大動態評測方法論全景圖", "Six Dynamic Evaluation Methodologies", accent_color=BLUE)
    tb_s9 = s9.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(13.8), Inches(4.7))
    s9_bullets = [
        "1. **黑箱評測 (Black-box Testing)**：無權重存取條件下，經由 API 注入對抗擾動樣本進行壓力測試。",
        "2. **白箱評測 (White-box Testing)**：調用內部梯度與網路層，使用 SHAP / LIME 產出特徵歸因熱力圖 (Point Game)。",
        "3. **基準測試 (Benchmark Testing)**：使用 ImageNet-C, AgentBench, PromptBench 等國防標準測試集。",
        "4. **紅軍演練 (Red Teaming)**：利用 garak, CALDERA Arsenal 在 Cyber Range 中模擬敵方對抗攻擊與 Prompt 注入。",
        "5. **專家評估 (Human Evaluation)**：指揮官配戴眼動儀與 EEG 評估認知負荷 (NASA-TLX) 與 HMT 信任校準。",
        "6. **營運持續監控 (Continuous Monitoring)**：線上部署 PyOD 與 Alibi 模組，即時警報數據與概念漂移。"
    ]
    add_formatted_bullets(tb_s9.text_frame, s9_bullets, font_size=13.0)
    add_acronym_footer(s9, ["SHAP", "LIME", "garak", "NASA-TLX", "EEG", "HMT", "PyOD", "API"], y_pos=7.65, height=1.0)

    # Slide 10
    s10 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s10)
    add_header(s10, "國防 AI 安全、保密與審計三維矩陣", "SECTION 2: T&E MATRIX & METHODOLOGY")
    add_icon_card(s10, 0.8, 1.6, 14.4, 5.9, "🛡️", "安全 (Security)、保密 (Confidentiality) 與審計 (Auditing) 聯防體系", "3D Defense Matrix", accent_color=BLUE)
    tb_s10 = s10.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(13.8), Inches(4.7))
    s10_bullets = [
        "1. **安全機制 (Security)**：對抗性攻擊防禦 (MITRE ATLAS / CALDERA)、無 GPS Mesh 邊緣硬體防篡改 (Hardware Tamper-Resistance) 與緊急模型自毀 (<100ms Flash/RAM 權重複寫)、JADC2 異質指管對接零信任 API 閘道 (SPIFFE/OPA)。",
        "2. **保密機制 (Confidentiality)**：100% 地端 On-Premise 實體隔離 (Air-Gap) 主權算力、聯邦學習 (Federated Learning)「模型移動，資料不動」與參數融合在地、Gemma 4 本地 LoRA 微調與 GGUF 4-bit 量化。",
        "3. **審計機制 (Auditing)**：資料與模型溯源 (Data & Model Provenance 清洗規則、版本號)、決策可解釋性 (XAI 特徵歸因熱力圖)、信心分數 (Confidence Score) 校準、高保真日誌包含輸入傳感器數據與行動方案 (COA)。"
    ]
    add_formatted_bullets(tb_s10.text_frame, s10_bullets, font_size=13.0)
    add_acronym_footer(s10, ["ATLAS", "SPIFFE", "OPA", "API", "JADC2", "GGUF", "LoRA", "XAI", "COA"], y_pos=7.65, height=1.0)

    # ----------------------------------------------------
    # SECTION 3: Slides 11-16 (System SOPs)
    # ----------------------------------------------------
    # Slide 11: A類
    s11 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s11)
    add_header(s11, "A類 - 電腦視覺與目標偵測評測 SOP", "SECTION 3: SYSTEM-SPECIFIC AI EVALUATION SOP")
    add_icon_card(s11, 0.8, 1.6, 14.4, 5.9, "👁️", "CV / Target Detection Evaluation SOP", accent_color=BLUE)
    tb_s11 = s11.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(13.8), Inches(4.7))
    s11_bullets = [
        "1. **適用範疇與核心威脅**：適用於 YOLO, Mobile SAM, 雷達 ISR 影像。防範對抗貼片 (Adversarial Patch)、FGSM/PGD 漸進擾動與天候干擾。",
        "2. **評測 SOP 與代表性工具鏈**：使用 IBM ART 360 / HEART 對樣本注入 ε 擾動；經由 NRTK 工具包合成 10 等級環境降質數據集壓力測試。",
        "3. **關鍵指標與合格門檻**：",
        "   - **Q1 對抗韌性**：$\mathrm{Acc}_{\mathrm{adv}}/\mathrm{Acc}_{\mathrm{clean}} \ge 0.90 \quad (\epsilon \le 0.05)$",
        "   - **Q2 自然穩健性**：$\Delta \mathrm{mAP} \le 0.10 \quad (10\%\text{ Limit})$",
        "   - **Q7 模型可解釋性**：$\mathrm{Point~Game~Score} \ge 0.85$ (XAITK / SHAP 歸因熱力圖)"
    ]
    add_formatted_bullets(tb_s11.text_frame, s11_bullets, font_size=13.0)
    add_acronym_footer(s11, ["CV", "YOLO", "ISR", "FGSM", "PGD", "ART", "HEART", "NRTK", "mAP", "XAITK", "SHAP"], y_pos=7.65, height=1.0)

    # Slide 12: B類
    s12 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s12)
    add_header(s12, "B類 - 生成式 AI 與大語言模型評測 SOP", "SECTION 3: SYSTEM-SPECIFIC AI EVALUATION SOP")
    add_icon_card(s12, 0.8, 1.6, 14.4, 5.9, "💬", "GenAI / LLM Evaluation SOP", accent_color=DARK_BLUE)
    tb_s12 = s12.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(13.8), Inches(4.7))
    s12_bullets = [
        "1. **適用範疇與核心威脅**：適用於指管對答助手、LLM 情報摘要。防範 Prompt 注入 (Prompt Injection)、角色扮演越獄與虛構幻覺。",
        "2. **評測 SOP 與代表性工具鏈**：使用 garak 漏洞掃描框架執行 10,000 筆測試案例；部署 NeMo Guardrails 與 PromptBench 對照 OWASP LLM Top 10。",
        "3. **關鍵指標與合格門檻**：",
        "   - **Q8 提示越獄與抗注入**：$R_{\mathrm{jailbreak\_def}} \ge 0.99 \quad (99\%)$",
        "   - **Q9 幻覺率與事實忠實度**：$\mathrm{Faithfulness} \ge 0.95 \quad \wedge \quad R_{\mathrm{hallucination}} \le 0.02$"
    ]
    add_formatted_bullets(tb_s12.text_frame, s12_bullets, font_size=13.0)
    add_acronym_footer(s12, ["GenAI", "LLM", "garak", "NeMo", "OWASP"], y_pos=7.65, height=1.0)

    # Slide 13: C類
    s13 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s13)
    add_header(s13, "C類 - 檢索增強生成 RAG 系統評測 SOP", "SECTION 3: SYSTEM-SPECIFIC AI EVALUATION SOP")
    add_icon_card(s13, 0.8, 1.6, 14.4, 5.9, "📚", "RAG Systems Evaluation SOP", accent_color=BLUE)
    tb_s13 = s13.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(13.8), Inches(4.7))
    s13_bullets = [
        "1. **適用範疇與核心威脅**：適用於國防研發技術文件庫、公文探勘。防範段落錯置、虛構引用與資訊降密洩漏 (Declassification Leakage)。",
        "2. **評測 SOP 與代表性工具鏈**：運用 RAGAS 與 TruLens 三元組 (RAG Triad) 對 1,000 組問答對核對；於 Milvus 實施向量 RBAC 標籤動態遮罩。",
        "3. **關鍵指標與合格門檻**：",
        "   - **Q10 檢索精確度與歸屬**：$\mathrm{Context~Precision} \ge 0.90 \quad \wedge \quad \mathrm{Attribution} \ge 0.98$",
        "   - **Q14 防降密洩漏率**：$R_{\mathrm{declass\_leak}} = 0\% \quad (\mathrm{RBAC~Output~Masking})$"
    ]
    add_formatted_bullets(tb_s13.text_frame, s13_bullets, font_size=13.0)
    add_acronym_footer(s13, ["RAG", "RAGAS", "RBAC"], y_pos=7.65, height=1.0)

    # Slide 14: D類
    s14 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s14)
    add_header(s14, "D類 - AI Agent 與多代理協同系統評測 SOP", "SECTION 3: SYSTEM-SPECIFIC AI EVALUATION SOP")
    add_icon_card(s14, 0.8, 1.6, 14.4, 5.9, "🤖", "AI Agent & Multi-Agent Evaluation SOP", accent_color=DARK_BLUE)
    tb_s14 = s14.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(13.8), Inches(4.7))
    s14_bullets = [
        "1. **適用範疇與核心威脅**：適用於自主網路防禦 Agent、戰術排程代理。防範軌跡偏移、越權 API 呼叫 (Tool Misuse) 與惡意指令刪除數據。",
        "2. **評測 SOP 與代表性工具鏈**：使用 AgentBench 記錄 Tool Call 軌跡；經由 Open Policy Agent (OPA) 與 SPIFFE/SPIRE 證書進行零信任策略攔截。",
        "3. **關鍵指標與合格門檻**：",
        "   - **Q11 Agent 工具調用合規**：$R_{\mathrm{unauth\_API}} = 0\% \quad \wedge \quad \mathrm{Task~Success} \ge 0.98$",
        "   - **Q15 系統軌跡可追溯性**：$\mathrm{Log~Coverage} = 100\% \quad \wedge \quad t_{\mathrm{reproduction}} \le 10\mathrm{min}$"
    ]
    add_formatted_bullets(tb_s14.text_frame, s14_bullets, font_size=13.0)
    add_acronym_footer(s14, ["API", "OPA", "SPIFFE", "SPIRE"], y_pos=7.65, height=1.0)

    # Slide 15: E類
    s15 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s15)
    add_header(s15, "E類 - 自主系統與人機協同評測 SOP", "SECTION 3: SYSTEM-SPECIFIC AI EVALUATION SOP")
    add_icon_card(s15, 0.8, 1.6, 14.4, 5.9, "🚁", "Autonomous Systems & HMT Evaluation SOP", accent_color=BLUE)
    tb_s15 = s15.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(13.8), Inches(4.7))
    s15_bullets = [
        "1. **適用範疇與交戰規則 (RoE)**：適用於無人機蜂群、自主武器。劃分 HITL (人在紐中/火力授權)、HOTL (人在紐上/實時監督)、HOOTL (完全自主/航路規劃)。",
        "2. **評測 SOP 與代表性工具鏈**：於 VBS 4 / EADSIM LVC 環境執行 100 次蒙地卡羅模擬；配合 ToAST 與 HITL 物理/邏輯雙重斷路器，遵循 DoDD 3000.09 指令。",
        "3. **關鍵指標與合格門檻**：",
        "   - **Q3 任務完成率**：$\mathrm{MSR} \ge 0.95 \quad (N = 100\text{ Runs})$",
        "   - **Q4 可中止性與失效安全**：$\tau_{\mathrm{abort}} \le 100\text{ms} \quad \wedge \quad \mathrm{FailSafe} = 100\%$",
        "   - **Q5 信任校準與過度依賴**：$\mathrm{ECE} \le 0.05 \quad \wedge \quad R_{\mathrm{overreliance}} \le 0.05$",
        "   - **Q6 認知負荷與適應性**：$\Delta \mathrm{TLX} \ge 0.30 \quad \wedge \quad \Delta t_{\mathrm{decision}} \le 2.0\text{s}$ (NASA-TLX, EEG)"
    ]
    add_formatted_bullets(tb_s15.text_frame, s15_bullets, font_size=12.5)
    add_acronym_footer(s15, ["RoE", "HITL", "HOTL", "HOOTL", "VBS", "EADSIM", "LVC", "ToAST", "DoDD", "MSR", "ECE", "NASA-TLX", "EEG"], y_pos=7.65, height=1.0)

    # Slide 16: F類
    s16 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s16)
    add_header(s16, "F類 - 決策支援與預測分析評測 SOP", "SECTION 3: SYSTEM-SPECIFIC AI EVALUATION SOP")
    add_icon_card(s16, 0.8, 1.6, 14.4, 5.9, "📈", "Predictive Analytics & C2 Evaluation SOP", accent_color=DARK_BLUE)
    tb_s16 = s16.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(13.8), Inches(4.7))
    s16_bullets = [
        "1. **適用範疇與核心威脅**：適用於聲納/雷達特徵識別、擊殺鏈 (Kill Chain) 決策輔助。防範戰場概念漂移 (Concept Drift) 與 OOD 數據高信心誤判。",
        "2. **評測 SOP 與代表性工具鏈**：部署 PyOD 與 Alibi 警報模組；採用 MC-Dropout 生成預測方差；日誌完整留存傳感器數據、Confidence Score 與 COA 建議。",
        "3. **關鍵指標與合格門檻**：",
        "   - **Q12 概念與數據漂移監控率**：$\mathrm{Drift~Recall} \ge 0.95 \quad \wedge \quad t_{\mathrm{alarm}} \le 5\mathrm{min}$",
        "   - **Q13 不確定性量化 (UQ)**：$\mathrm{OOD~Variance~Coverage} \ge 0.95 \quad (95\%)$",
        "   - **Q15 軌跡可追溯性**：高保真日誌包含傳感器數據、Confidence Score 與 COA ($t_{\mathrm{repro}} \le 10\text{min}$)"
    ]
    add_formatted_bullets(tb_s16.text_frame, s16_bullets, font_size=13.0)
    add_acronym_footer(s16, ["PyOD", "OOD", "UQ", "MC-Dropout", "COA"], y_pos=7.65, height=1.0)

    # ----------------------------------------------------
    # SECTION 4: Slides 17-25 (15 Quantitative Metrics)
    # Requirement 5: Move bottom footnote area UP on Slide 17 & 18 for breathing room!
    # Requirement 6: Purge ALL "(LaTeX 渲染公式)" or "(LaTeX 編譯渲染結果)" text strings!
    # ----------------------------------------------------
    s17 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s17)
    add_header(s17, "15 項 AI 量化評測指標與合格門檻總覽", "SECTION 4: 15 QUANTITATIVE EVALUATION METRICS & SOPS")

    tb_s17_m = s17.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(14.4), Inches(4.5))
    tf_17 = tb_s17_m.text_frame
    tf_17.word_wrap = True

    m_summary = [
        "1. **Q1 對抗韌性**：$\mathrm{Acc}_{\mathrm{adv}}/\mathrm{Acc}_{\mathrm{clean}} \ge 0.90$  |  2. **Q2 自然穩健性**：$\Delta \mathrm{mAP} \le 0.10$  |  3. **Q3 任務完成率**：$\mathrm{MSR} \ge 0.95$",
        "4. **Q4 可中止性**：$\tau_{\mathrm{abort}} \le 100\text{ms} \wedge \mathrm{FailSafe} = 100\%$  |  5. **Q5 信任校準**：$\mathrm{ECE} \le 0.05$  |  6. **Q6 認知負荷**：$\Delta \mathrm{TLX} \ge 0.30$",
        "7. **Q7 可解釋性**：$\mathrm{Point~Game} \ge 0.85$  |  8. **Q8 提示越獄**：$R_{\mathrm{jailbreak\_def}} \ge 0.99$  |  9. **Q9 幻覺率**：$\mathrm{Faithfulness} \ge 0.95$",
        "10. **Q10 檢索精確度**：$\mathrm{Precision} \ge 0.90$  |  11. **Q11 Agent合規**：$R_{\mathrm{unauth\_API}} = 0\%$  |  12. **Q12 漂移召回**：$\mathrm{Recall} \ge 0.95$",
        "13. **Q13 不確定性量化**：$\mathrm{OOD~Coverage} \ge 0.95$  |  14. **Q14 防降密**：$R_{\mathrm{declass\_leak}} = 0\%$  |  15. **Q15 軌跡追溯**：$\mathrm{Log~Coverage} = 100\%$"
    ]
    add_formatted_bullets(tf_17, m_summary, font_size=12.5)

    # Requirement 5: Slide 17 Appendix Pointer Card moved UP to y=6.0 with comfortable padding
    card_ptr = s17.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.0), Inches(14.4), Inches(1.25))
    card_ptr.fill.solid()
    card_ptr.fill.fore_color.rgb = RGBColor(239, 246, 255)
    card_ptr.line.color.rgb = BLUE
    card_ptr.line.width = Pt(1.5)

    tf_ptr = card_ptr.text_frame
    tf_ptr.word_wrap = True
    p_ptr_h = tf_ptr.paragraphs[0]
    p_ptr_h.text = "📌 附件參考說明 (Appendix Reference Note & Detailed Walkthrough Pointer):"
    p_ptr_h.font.size = Pt(11)
    p_ptr_h.font.bold = True
    p_ptr_h.font.color.rgb = BLUE
    p_ptr_h.font.name = "微軟正黑體"

    p_ptr_b = tf_ptr.add_paragraph()
    ptr_str = "本頁面為 15 項 AI 量化評測指標與合格門檻總覽。各指標之**詳細測試 SOP 實施步驟**、**代表性測試工具鏈 (ART, garak, RAGAS, PyOD)**、**對抗攻防演練**與 **Obsidian 雙腦筆記架構**，請參閱後附獨立附件檔案 [AIEC_15_Quantitative_Metrics_SOP.pptx](file:///c:/Users/administartor/Downloads/AIEC/AIEC_15_Quantitative_Metrics_SOP.pptx) 與 [[AIEC 15 項量化評測指標與 SOP]]。"
    add_formatted_text(p_ptr_b, ptr_str, font_size=10.0, default_color=NAVY, bold_color=BLUE)

    # Requirement 5: Acronym footer moved to y=7.55 for spacious breathing room!
    add_acronym_footer(s17, ["Acc_adv", "Acc_clean", "mAP", "MSR", "ECE", "garak", "PyOD", "ART", "RAGAS", "SOP"], y_pos=7.55, height=1.1)

    # Metrics list for Slides 18-25
    # Requirement 6: PURGED "(LaTeX 渲染公式)" from all left/right bullets!
    metrics_list = [
        {"id": 1, "name_zh": "對抗韌性", "name_en": "Adversarial Robustness", "def": "模型遭受對抗貼片、FGSM/PGD 漸進式擾動攻擊時維護正確判讀與目標識別的能力。", "latex_calc": r'\mathrm{Robustness~Ratio} = \frac{\mathrm{Acc}_{\mathrm{adv}}(\mathcal{D}_{\mathrm{test}}, \epsilon)}{\mathrm{Acc}_{\mathrm{clean}}(\mathcal{D}_{\mathrm{test}})}', "scope": "A類 電腦視覺 (CV)、目標偵測與影像分類演算法模型", "risk": "敵方附著對抗貼片引發目標誤判、漏報或偽裝欺騙", "sop": "使用 IBM ART 360 / HEART 對輸入樣本注入 ε 漸進式對抗擾動，量測判讀降質與 mAP 變化率。", "tools": "IBM ART 360, HEART Framework, FGSM, PGD Attack Engine", "latex_thresh": r'\mathrm{PASS:}~\frac{\mathrm{Acc}_{\mathrm{adv}}}{\mathrm{Acc}_{\mathrm{clean}}} \geq 0.90 \quad (\epsilon \leq 0.05)', "audit": "產出對抗擾動強度與 Acc 衰減曲線圖，納入 ISO 42001 驗測報告", "std": "MITRE ATLAS (AML.T0015) / NIST AI RMF 1.0 (Measure 2.1)", "acronyms": ["Acc_adv", "Acc_clean", "FGSM", "PGD", "mAP", "ART", "HEART", "MITRE", "ATLAS", "NIST", "RMF", "ISO"]},
        {"id": 2, "name_zh": "自然穩健性", "name_en": "Natural Robustness", "def": "模型面對自然環境干擾（雨雪、煙霧、電戰雜訊、光影突變）時的效能維持能力。", "latex_calc": r'\Delta \mathrm{mAP} = \frac{\mathrm{mAP}_{\mathrm{clean}} - \mathrm{mAP}_{\mathrm{noise}}(\eta)}{\mathrm{mAP}_{\mathrm{clean}}}', "scope": "A類 電腦視覺 (CV)、雷達 ISR 影像與感測器融合系統", "risk": "極端氣候或電戰環境干擾導致目標偵測精度暴降或系統失效", "sop": "透過 NRTK 工具包合成 10 種等級的環境降質數據集進行壓力測試，對比清晰數據集效能。", "tools": "NRTK (Natural Robustness Toolkit), ImageNet-C benchmark", "latex_thresh": r'\mathrm{PASS:}~\Delta \mathrm{mAP} \leq 0.10 \quad (10\%~\mathrm{Limit})', "audit": "記錄不同天候降質條件下之精度衰減曲線，作為部署前確效數據", "std": "JATIC 共通構面 1 / DoD CDAO AI T&E Guidebook", "acronyms": ["mAP", "NRTK", "JATIC", "DoD", "CDAO", "T&E"]},
        {"id": 3, "name_zh": "任務完成率", "name_en": "Mission Success Rate (MSR)", "def": "AI 自主/半自主系統在端到端戰術情境中成功執行完畢並閉合擊殺鏈 (Kill Chain) 的比例。", "latex_calc": r'\mathrm{MSR} = \frac{\sum_{i=1}^{N} S_i}{N}, \quad S_i \in \{0, 1\}', "scope": "E類 自主武器系統、無人載具蜂群與 C2 指管決策系統", "risk": "戰術情境極度複雜導致 AI 決策邏輯死鎖或任務中途失敗", "sop": "於 VBS 4 / EADSIM 虛實整合 (LVC) 平行戰場環境中執行 100 次蒙地卡羅戰術模擬。", "tools": "VBS 4, EADSIM, LVC 平行戰場模擬環境", "latex_thresh": r'\mathrm{PASS:}~\mathrm{MSR} = \frac{\mathrm{Success~Count}}{N} \geq 0.95 \quad (N = 100)', "audit": "保存 100 次 LVC 模擬軌跡與關鍵決策點 Log，符合 OT&E 要求", "std": "DoD T&E Level 4 Operational T&E", "acronyms": ["MSR", "LVC", "VBS", "EADSIM", "DoD", "T&E", "OT&E"]},
        {"id": 4, "name_zh": "可中止性與失效安全", "name_en": "Abortability & Fail-Safe Rate", "def": "當 AI 系統發生異常或接獲人工中斷指令時，即刻安全中斷並進入預設安全保護狀態的能力。", "latex_calc": r'\tau_{\mathrm{abort}} = t_{\mathrm{safe\_state}} - t_{\mathrm{signal\_sent}}, \quad \mathrm{Fail{-}Safe~Rate} = \frac{N_{\mathrm{safe}}}{N_{\mathrm{trigger}}}', "scope": "E類 自主武器系統、無人打擊載具與自動化防禦系統", "risk": "系統異常失控且無法經由手動指令強制中斷，引發非預期災害", "sop": "於模擬任務中隨機注入手動 Stop Signal、通訊斷連與異常偏移，量測安全降級接管回應時間。", "tools": "ToAST (Testing of Autonomous Systems Tool), HITL 硬體斷路器", "latex_thresh": r'\mathrm{PASS:}~\tau_{\mathrm{abort}} \leq 100\mathrm{ms} \quad \wedge \quad \mathrm{Fail{-}Safe~Rate} = 100\%', "audit": "驗證物理/邏輯雙重斷路器回應時間，嚴格遵循 DoDD 3000.09 指令", "std": "DoDD 3000.09 自主武器系統指令", "acronyms": ["HITL", "ToAST", "DoDD", "T&E"]},
        {"id": 5, "name_zh": "信任校準與過度依賴", "name_en": "Trust Calibration & Over-Reliance", "def": "操作員對 AI 輸出信心度的理解符合模型實際能力，防止盲目過度信任或拒絕依賴。", "latex_calc": r'\mathrm{ECE} = \sum_{m=1}^{M} \frac{|B_m|}{N} |\mathrm{acc}(B_m) - \mathrm{conf}(B_m)|, \quad R_{\mathrm{overreliance}} = \frac{N_{\mathrm{blind}}}{N_{\mathrm{false}}}', "scope": "E類 人機協同團隊 (HMT)、指管決策輔助與醫療/後勤 AI", "risk": "操作員對 AI 高信心錯誤提案盲目採納，導致嚴重戰術失誤", "sop": "於 HMT 模擬試驗中故意提供高信心度但錯誤之提案，記錄操作員即時發現與修正之反應率。", "tools": "HMT Evaluation Suite, Trust-in-Automation Scale, ECE Calculator", "latex_thresh": r'\mathrm{PASS:}~\mathrm{ECE} \leq 0.05 \quad \wedge \quad R_{\mathrm{overreliance}} \leq 0.05', "audit": "統計操作員對 AI 提案之修正率與信心度散佈圖，符合 HMT 指南", "std": "JATIC 共通構面 6 / DoD HMT T&E Guidebook", "acronyms": ["ECE", "HMT", "JATIC", "DoD", "T&E"]},
        {"id": 6, "name_zh": "認知負荷與適應性", "name_en": "Cognitive Load & Adaptability", "def": "AI 介面輸出與警報對指揮官或操作員造成的心理負荷程度、決策時延與適應性。", "latex_calc": r'\Delta \mathrm{TLX} = \frac{\mathrm{TLX}_{\mathrm{baseline}} - \mathrm{TLX}_{\mathrm{AI}}}{\mathrm{TLX}_{\mathrm{baseline}}}, \quad \Delta t_{\mathrm{decision}} = t_{\mathrm{response}}', "scope": "E類 人機整合介面 (HSI)、戰術儀表板與指揮管制系統", "risk": "介面資訊過載或警報頻發引發操作員決策恐慌與認知失調", "sop": "操作員配戴眼動儀與 EEG 腦電圖儀完成應變任務後，填寫 NASA-TLX 心理負荷量表。", "tools": "NASA-TLX 心理負荷量表, EEG 腦電儀, Eye-Tracking Suite", "latex_thresh": r'\mathrm{PASS:}~\Delta \mathrm{TLX} \geq 0.30 \quad \wedge \quad \Delta t_{\mathrm{decision}} \leq 2.0\mathrm{s}', "audit": "量測眼動凝視時間與 EEG 腦波數據，確效 HSI 介面優化效益", "std": "DoD T&E Level 2 HSI (Human-Systems Integration)", "acronyms": ["NASA-TLX", "EEG", "HSI", "DoD", "T&E"]},
        {"id": 7, "name_zh": "模型可解釋性與顯著性歸因", "name_en": "Explainability & Point Game", "def": "AI 關鍵決策邏輯機能是否提供可被人類審計與驗證之特徵熱力圖 (Saliency Map)。", "latex_calc": r'\mathrm{Point~Game~Score} = \frac{N_{\mathrm{hit}}(\mathrm{argmax~Saliency} \in \mathrm{ROI})}{N_{\mathrm{total}}}', "scope": "A類 電腦視覺、F類 決策支援與醫療診斷模型", "risk": "黑盒子模型根據背景偽特徵做出判斷，引發不可靠決策", "sop": "白箱調用 XAITK / SHAP / LIME 工具產出特徵歸因熱力圖，比對真實目標重點區域 (Point Game)。", "tools": "XAITK (Explainable AI Toolkit), SHAP, LIME, Grad-CAM", "latex_thresh": r'\mathrm{PASS:}~\mathrm{Point~Game~Score} \geq 0.85 \quad (85\%)', "audit": "匯出目標歸因熱力圖與 Point Game 數據，合規 ISO 42001 透明度", "std": "ISO/IEC 42001 Clause 8.4 透明度要求", "acronyms": ["XAITK", "SHAP", "LIME", "ISO", "AIMS"]},
        {"id": 8, "name_zh": "提示越獄與抗注入能力", "name_en": "Prompt Jailbreak Defense Rate", "def": "LLM 阻絕敵方對抗 Prompt 注入、越獄繞過、角色扮演突圍與護欄穿透的能力。", "latex_calc": r'R_{\mathrm{jailbreak\_def}} = 1 - \frac{N_{\mathrm{successful\_jailbreaks}}}{N_{\mathrm{total\_attacks}}} = \frac{N_{\mathrm{blocked}}}{N_{\mathrm{total\_attacks}}}', "scope": "B類 生成式 AI (GenAI)、LLM 指管對答與情報摘要助手", "risk": "敵方注入對抗指令操控 LLM 輸出危害國防安全之內容或指令", "sop": "使用 garak 自動化漏洞掃描框架執行 10,000 筆測試案例 (Direct/Indirect Injection)。", "tools": "garak (LLM Vulnerability Scanner), NeMo Guardrails, PromptBench", "latex_thresh": r'\mathrm{PASS:}~R_{\mathrm{jailbreak\_def}} \geq 0.99 \quad (99\%)', "audit": "自動生成 garak 漏洞掃描日誌，對照 OWASP LLM Top 10 安全基準", "std": "OWASP LLM Top 10 (LLM01/LLM02) / garak benchmark", "acronyms": ["garak", "NeMo", "LLM", "OWASP", "SOP"]}
    ]

    for idx, m in enumerate(metrics_list):
        s_m = prs.slides.add_slide(blank_layout)
        set_pure_white_bg(s_m)
        add_header(s_m, f"Q{m['id']}. {m['name_zh']} ({m['name_en']})", "SECTION 4: 15 QUANTITATIVE EVALUATION METRICS & SOPS")

        calc_img = render_latex_to_png(m['latex_calc'], f"calc_m_q{m['id']}.png")
        thresh_img = render_latex_to_png(m['latex_thresh'], f"thresh_m_q{m['id']}.png")

        add_icon_card(s_m, 0.8, 1.5, 7.0, 5.9, "📌", "指標定義與計算公式", "Metric Definition & Rendered Math", accent_color=BLUE)
        tb_l = s_m.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(6.4), Inches(1.3))
        # Requirement 6: Pure text without "(LaTeX 渲染公式)"
        left_bullets_top = [f"1. **核心指標定義**：{m['def']}", "2. **量化計算數學公式**："]
        add_formatted_bullets(tb_l.text_frame, left_bullets_top, font_size=12.5)
        s_m.shapes.add_picture(calc_img, Inches(1.3), Inches(3.6), width=Inches(6.0))

        tb_l_bot = s_m.shapes.add_textbox(Inches(1.1), Inches(4.5), Inches(6.4), Inches(2.7))
        left_bullets_bot = [f"3. **適用範疇與風險關切**：{m['scope']}；{m['risk']}", f"4. **對應國際權威標準**：{m['std']}"]
        add_formatted_bullets(tb_l_bot.text_frame, left_bullets_bot, font_size=12.5)

        add_icon_card(s_m, 8.2, 1.5, 7.0, 5.9, "🧪", "驗測 SOP 與合格判定門檻", "Testing SOP & Pass Threshold Math", accent_color=DARK_BLUE)
        tb_r = s_m.shapes.add_textbox(Inches(8.5), Inches(2.5), Inches(6.4), Inches(1.3))
        # Requirement 6: Pure text without "(LaTeX 渲染公式)"
        right_bullets_top = [f"1. **驗測 SOP 實施步驟**：{m['sop']}", f"2. **代表性測試工具鏈**：{m['tools']}", "3. **量化合格門檻公式**："]
        add_formatted_bullets(tb_r.text_frame, right_bullets_top, font_size=12.5)
        s_m.shapes.add_picture(thresh_img, Inches(8.7), Inches(4.3), width=Inches(6.0))

        tb_r_bot = s_m.shapes.add_textbox(Inches(8.5), Inches(5.2), Inches(6.4), Inches(2.0))
        right_bullets_bot = [f"4. **合規稽核與紀錄規範**：{m['audit']}"]
        add_formatted_bullets(tb_r_bot.text_frame, right_bullets_bot, font_size=12.5)

        # Requirement 5: Move acronym footer UP to y=7.55 for Slide 18 & all metric slides!
        add_acronym_footer(s_m, m['acronyms'], y_pos=7.55, height=1.1)

    # ----------------------------------------------------
    # SECTION 5: Slides 26-29 - Platforms & Infrastructure
    # Requirement 7: DELETED Slide 29 (AI CLI與Cron) !
    # Slide 26: 主權 AI 平台與四層 LLM 堆疊
    # Slide 27: 地端 LLM 推論引擎與 Middleware
    # Slide 28: Lattice 戰術 C2 架構與 Menace 邊緣算力節點
    # Slide 29: 聯邦學習 (Federated Learning) 國防保密策略
    # ----------------------------------------------------
    s26 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s26)
    add_header(s26, "主權 AI 平台與四層 LLM 算力堆疊", "SECTION 5: EVALUATION PLATFORMS & AUTOMATED T&E")
    add_icon_card(s26, 0.8, 1.6, 14.4, 5.9, "🏛️", "Sovereign AI Platform & 4-Tier Compute Stack", accent_color=BLUE)
    tb_s26 = s26.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(13.8), Inches(4.7))
    s26_bullets = [
        "1. **民雄院區國家級主權算力**：100% 地端實體隔離 (Air-Gapped) 設施，具備電力韌性、物理防護與演算法安全機制。",
        "2. **Tier 1: 密集型基底模型 (Dense Foundation)**：70B 主權大模型，負責核心情報與高階戰術分析。",
        "3. **Tier 2: 專家混合架構 (MoE Model)**：8x7B 專用模型，針對電戰、圖像識別與資安防禦專精分工。",
        "4. **Tier 3: 邊緣輕量化模型 (Edge Quantized)**：7B / 3B GGUF 量化模型，部署於無人載具與前線節點。",
        "5. **Tier 4: 受控雲端 API 閘道 (Gated Cloud API)**：僅開放非密級資料探勘，設有 SPIFFE/OPA 策略驗證。"
    ]
    add_formatted_bullets(tb_s26.text_frame, s26_bullets, font_size=13.0)
    add_acronym_footer(s26, ["MoE", "GGUF", "API", "SPIFFE", "OPA"], y_pos=7.65, height=1.0)

    s27 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s27)
    add_header(s27, "地端 LLM 推論引擎與 Middleware 工具鏈", "SECTION 5: EVALUATION PLATFORMS & AUTOMATED T&E")
    add_icon_card(s27, 0.8, 1.6, 14.4, 5.9, "⚙️", "On-Prem LLM Inference Engine & Middleware", accent_color=DARK_BLUE)
    tb_s27 = s27.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(13.8), Inches(4.7))
    s27_bullets = [
        "1. **地端模型蒸餾與微調 (On-Prem Distillation & Fine-tuning)**：將 70B 教師模型知識蒸餾 (Distillation) 至 Gemma 4 本地模型，進行 LoRA 輕量化微調。",
        "2. **vLLM 推論引擎**：採用 PagedAttention 記憶體優化，支援多用戶高併發戰術對抗檢索。",
        "3. **llama.cpp / Ollama**：支援 GGUF 4-bit 量化格式，無 Python 依賴，於戰術邊緣硬體實現高 FPS 推論。",
        "4. **Data & Model Provenance 溯源**：紀錄資料清洗規則、模型超參數版本號與 Checkpoint 雜湊值。"
    ]
    add_formatted_bullets(tb_s27.text_frame, s27_bullets, font_size=13.0)
    add_acronym_footer(s27, ["LLM", "vLLM", "GGUF", "LoRA", "SOP"], y_pos=7.65, height=1.0)

    s28 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s28)
    add_header(s28, "Lattice 戰術 C2 架構與 Menace 邊緣算力節點", "SECTION 5: EVALUATION PLATFORMS & AUTOMATED T&E")
    add_icon_card(s28, 0.8, 1.6, 14.4, 5.9, "📡", "Tactical C2 Architecture & Edge Compute Nodes", accent_color=BLUE)
    tb_s28 = s28.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(13.8), Inches(4.7))
    s28_bullets = [
        "1. **JADC2 全領域指管對接**：整合異質指管系統（如將本地飛彈陣地數據與商用 C2 平台對接）。",
        "2. **Menace 邊緣算力節點**：於前線戰術車輛與防空陣地部署微型計算模組，提供即時辨識與決策支援。",
        "3. **零信任 API 閘道邊界驗證**：採用 SPIFFE/SPIRE 跨網段服務數位身分認證與 OPA 微秒級策略過濾。",
        "4. **通訊中斷保護**：具備邊緣斷網自主運轉與復網後數據自動差分同步能力。"
    ]
    add_formatted_bullets(tb_s28.text_frame, s28_bullets, font_size=13.0)
    add_acronym_footer(s28, ["C2", "JADC2", "API", "SPIFFE", "SPIRE", "OPA"], y_pos=7.65, height=1.0)

    s29 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s29)
    add_header(s29, "聯邦學習 (Federated Learning) 國防保密策略", "SECTION 5: EVALUATION PLATFORMS & AUTOMATED T&E")
    add_icon_card(s29, 0.8, 1.6, 14.4, 5.9, "🤝", "Federated Learning Strategy & Parameter Fusion", accent_color=BLUE)
    tb_s29 = s29.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(13.8), Inches(4.7))
    s29_bullets = [
        "1. **「模型移動，資料不動」原則**：跨防空陣地、雷達站與艦隊聯合訓練時，原始作戰數據嚴禁傳出本地戰術網段。",
        "2. **參數融合在地 (Parameter Fusion)**：各節點僅回傳經訓練後之梯度 (Gradient) 與權重更新，由中央節點進行 FedAvg 彙整。",
        "3. **差分隱私 (Differential Privacy, DP)**：於上傳梯度中注入校準雜訊，防止敵方利用逆向推導還原原始特徵。",
        "4. **同態加密 (Homomorphic Encryption)**：對上傳參數進行同態加密，中央節點直接在密文空間執行參數融合。"
    ]
    add_formatted_bullets(tb_s29.text_frame, s29_bullets, font_size=13.0)
    add_acronym_footer(s29, ["DP", "FedAvg"], y_pos=7.65, height=1.0)

    # ----------------------------------------------------
    # SECTION 6: Slides 30-33 - Advanced Defense AI Security & Infrastructure
    # Requirement 10: Compile ALL LaTeX math expressions on Slides 30-33 into rendered math images!
    # Requirement 8: Slide 30 Footnote -> Add YOLO, SAM!
    # Requirement 9: Slide 31 Footnote -> Add ROE / RoE!
    # ----------------------------------------------------

    # Render LaTeX images for Section 6
    img_tau_abort = render_latex_to_png(r'\tau_{\mathrm{abort}} \leq 100\mathrm{ms}', 'sec6_tau_abort.png')
    img_r_declass = render_latex_to_png(r'R_{\mathrm{declass\_leak}} = 0\%', 'sec6_r_declass.png')
    img_ece_val = render_latex_to_png(r'\mathrm{ECE} \leq 0.05', 'sec6_ece.png')

    # Slide 30 (NCSIST 藍圖 - Requirement 8: Add YOLO & SAM to footnote)
    s30 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s30)
    add_header(s30, "NCSIST AIEC 國防 AI 評測總體藍圖與 5 大戰術柱石", "SECTION 6: ADVANCED DEFENSE AI SECURITY & INFRASTRUCTURE")

    add_icon_card(s30, 0.8, 1.5, 7.0, 5.9, "🏛️", "NCSIST AIEC 架構圖與主動防禦", "NCSIST Defense Architecture Blueprint", accent_color=BLUE)
    tb_s30_l = s30.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(6.4), Inches(4.7))
    s30_l_bullets = [
        "1. **代碼注入與數據污染主動防禦**：針對視覺辨識 (YOLO, Mobile SAM) 與 LLM 模型，建立代碼注入 (Code Injection) 與數據污染 (Poisoning) 過濾閘門。",
        "2. **Cyber Range 對抗演練**：結合 [[MITRE ATLAS 人工智慧對抗威脅矩陣]] 與 CALDERA Arsenal 外掛，建立動態資安脆弱性模擬。",
        "3. **國家級機構中樞**：由國家中山科學研究院 (NCSIST) 主導主權算力與國防 AIEC 評測中心規範放行。"
    ]
    add_formatted_bullets(tb_s30_l.text_frame, s30_l_bullets, font_size=12.5)

    add_icon_card(s30, 8.2, 1.5, 7.0, 5.9, "🛡️", "底層 5 大戰術 SOP 柱石", "Five Tactical Operational Pillar SOPs", accent_color=DARK_BLUE)
    tb_s30_r = s30.shapes.add_textbox(Inches(8.5), Inches(2.5), Inches(6.4), Inches(4.7))
    s30_r_bullets = [
        "1. **擬定 RoE 交戰規則**：明確定義人機協同與可控性 (HITL/HOTL/HOOTL)。",
        "2. **驗證分層式 AI 架構**：確效系統能在極端作戰條件下完成任務。",
        "3. **執行紅隊對抗/演訓**：測試系統承受攻擊與受損後之彈性恢復力。",
        "4. **研發/軍事資料集分級**：落實清洗規則、版本控管與 Data/Model Provenance。",
        "5. **供應鏈安全檢查**：審查開源模型與第三方 SDK 後門及邏輯合規。"
    ]
    add_formatted_bullets(tb_s30_r.text_frame, s30_r_bullets, font_size=12.5)
    # Requirement 8: Add YOLO and SAM to Slide 30 Footnote!
    add_acronym_footer(s30, ["NCSIST", "AIEC", "ATLAS", "YOLO", "SAM", "RoE", "HITL", "HOTL", "HOOTL", "CMMC", "TRL", "SOP"], y_pos=7.55, height=1.1)

    # Slide 31 (RoE 邊界 - Requirement 9: Add ROE / RoE to footnote)
    s31 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s31)
    add_header(s31, "國防 AIEC 核心任務與交戰規則 (RoE) 授權邊界", "SECTION 6: ADVANCED DEFENSE AI SECURITY & INFRASTRUCTURE")

    add_icon_card(s31, 0.8, 1.5, 7.0, 5.9, "🎯", "國防 AIEC 五大核心管制任務", "Five Core Control Tasks of Defense AIEC", accent_color=BLUE)
    tb_s31_l = s31.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(6.4), Inches(4.7))
    s31_l_bullets = [
        "1. **制定交戰規則 (RoE)**：劃分全系統人機自動化權能授權邊界。",
        "2. **審核主權基礎設施**：審核民雄院區三層算力與電力/物理韌性。",
        "3. **策劃 LVC 紅軍測試**：於 VBS 4 / EADSIM 平行戰場測試對抗防禦。",
        "4. **RAG 降密稽核**：審查研發資料集探勘，防範降密洩漏：",
        "5. **供應鏈後門審查**：嚴格審查開源模型與 SDK 演算法後門。"
    ]
    add_formatted_bullets(tb_s31_l.text_frame, s31_l_bullets, font_size=12.5)
    # Embed rendered LaTeX image for RAG anti-declassification rate on Slide 31
    s31.shapes.add_picture(img_r_declass, Inches(1.3), Inches(4.5), width=Inches(5.0))

    add_icon_card(s31, 8.2, 1.5, 7.0, 5.9, "⚖️", "RoE 人機授權三階權能邊界", "Three-Tier Autonomy Boundary (RoE)", accent_color=DARK_BLUE)
    tb_s31_r = s31.shapes.add_textbox(Inches(8.5), Inches(2.5), Inches(6.4), Inches(4.7))
    s31_r_bullets = [
        "1. **Human-in-the-loop (HITL / 人在紐中)**：最終火力打擊與目標授權必須由人類指揮官手動執行。",
        "2. **Human-on-the-loop (HOTL / 人在紐上)**：AI 執行戰術建議，人類具備實時監控與強制中斷權 (Abort Button)。",
        "3. **Human-out-of-the-loop (HOOTL / 完全自主)**：僅限於蜂群無人機航路規劃與偵察等非致傷性任務。"
    ]
    add_formatted_bullets(tb_s31_r.text_frame, s31_r_bullets, font_size=12.5)
    # Requirement 9: Add ROE / RoE to Slide 31 Footnote!
    add_acronym_footer(s31, ["ROE", "RoE", "HITL", "HOTL", "HOOTL", "LVC", "VBS", "EADSIM", "RAG", "RBAC", "TRL"], y_pos=7.55, height=1.1)

    # Slide 32 (邊緣自毀 - Requirement 10: Render LaTeX math equations)
    s32 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s32)
    add_header(s32, "戰術邊緣硬體安全與模型緊急自毀機制", "SECTION 6: ADVANCED DEFENSE AI SECURITY & INFRASTRUCTURE")

    add_icon_card(s32, 0.8, 1.5, 7.0, 5.9, "🛡️", "無 GPS 網狀通訊與硬體級防篡改", "Hardware Tamper-Resistance in Mesh Networks", accent_color=BLUE)
    tb_s32_l = s32.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(6.4), Inches(4.7))
    s32_l_bullets = [
        "1. **網狀通訊被俘獲風險 (Mesh Network)**：在前線電戰干擾與無 GPS 網狀環境下，無人機與邊緣運算設備極易遭敵方俘獲。",
        "2. **硬體級防篡改 (Tamper-Resistance)**：採用軍規 Secure Enclave / TPM 晶片與物理開箱感測電路 (Tamper Mesh)。",
        "3. **物理防護障壁**：設備遭受異常解體或外殼穿透時，立即切斷主板電源並引發邏輯保護。"
    ]
    add_formatted_bullets(tb_s32_l.text_frame, s32_l_bullets, font_size=12.5)

    add_icon_card(s32, 8.2, 1.5, 7.0, 5.9, "💣", "模型緊急自毀與物理零化", "Emergency Self-Destruct & Zeroization", accent_color=DARK_BLUE)
    tb_s32_r = s32.shapes.add_textbox(Inches(8.5), Inches(2.5), Inches(6.4), Inches(2.0))
    s32_r_bullets = [
        "1. **緊急權重複寫指令**：當系統偵測到物理入侵、通訊斷連超過安全時限或接獲中斷信號時，極速觸發自毀回應時間：",
        "2. **Flash/RAM 雜訊衝刷**：以隨機雜訊迅速覆寫模型記憶體與 Flash 快閃記憶體，徹底撕毀演算法與權重。",
        "3. **物理零化 (Zeroization)**：隨機化金鑰熔絲 (Key Fusing)，防止敵方逆向工程與模型竊取。"
    ]
    add_formatted_bullets(tb_s32_r.text_frame, s32_r_bullets, font_size=12.5)
    # Requirement 10: Render LaTeX math equation img_tau_abort on Slide 32!
    s32.shapes.add_picture(img_tau_abort, Inches(8.7), Inches(3.2), width=Inches(5.0))
    add_acronym_footer(s32, ["GPS", "Mesh", "TPM", "RAM", "Flash"], y_pos=7.55, height=1.1)

    # Slide 33 (地端蒸餾與 Provenance - Requirement 10: Render LaTeX math equations)
    s33 = prs.slides.add_slide(blank_layout)
    set_pure_white_bg(s33)
    add_header(s33, "地端模型蒸餾、資料與模型溯源 SOP", "SECTION 6: ADVANCED DEFENSE AI SECURITY & INFRASTRUCTURE")

    add_icon_card(s33, 0.8, 1.5, 7.0, 5.9, "⚙️", "地端模型蒸餾與輕量化微調", "On-Premise Distillation & LoRA Fine-Tuning", accent_color=BLUE)
    tb_s33_l = s33.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(6.4), Inches(4.7))
    s33_l_bullets = [
        "1. **知識蒸餾 (Knowledge Distillation)**：將地端 70B 主權教師模型之推理能力蒸餾至 7B / 3B 輕量化邊緣學生模型。",
        "2. **Gemma 4 LoRA 本地微調**：針對戰術任務進行輕量化適應微調 (LoRA)，僅更新微小參數比例，兼顧時延與精度。",
        "3. **GGUF 4-bit 量化**：透過 llama.cpp 實施 GGUF 量化部署，適應邊緣低功耗晶片。"
    ]
    add_formatted_bullets(tb_s33_l.text_frame, s33_l_bullets, font_size=12.5)

    add_icon_card(s33, 8.2, 1.5, 7.0, 5.9, "📊", "Data & Model Provenance 溯源與 Confidence", "Data & Model Provenance & Confidence Calibration", accent_color=DARK_BLUE)
    tb_s33_r = s33.shapes.add_textbox(Inches(8.5), Inches(2.5), Inches(6.4), Inches(2.2))
    s33_r_bullets = [
        "1. **Data Provenance (資料溯源)**：詳細記錄訓練資料集來源、邏輯檢查規則、清洗演算法版本號 (`v1.4.2`) 與標籤品質。",
        "2. **Model Provenance (模型溯源)**：追溯模型訓練超參數與 Checkpoint 雜湊值；若發現戰術誤判，可於 10 分鐘內回溯隔離受污染數據。",
        "3. **信心分數 (Confidence Score) 校準**：所有模型輸出伴隨 Confidence Score，確保期望校準誤差指標滿定："
    ]
    add_formatted_bullets(tb_s33_r.text_frame, s33_r_bullets, font_size=12.5)
    # Requirement 10: Render LaTeX math equation img_ece_val on Slide 33!
    s33.shapes.add_picture(img_ece_val, Inches(8.7), Inches(4.7), width=Inches(4.5))
    add_acronym_footer(s33, ["LLM", "LoRA", "GGUF", "ECE", "SOP"], y_pos=7.55, height=1.1)

    out_dir = r'c:\Users\administartor\Downloads\AIEC'
    out_path = os.path.join(out_dir, 'AIEC_AI_Evaluation_30_Slides_NanoBanana.pptx')
    prs.save(out_path)
    print(f'Successfully generated master 33-slide presentation at: {out_path}')

if __name__ == '__main__':
    generate_33_slides_master_deck()
